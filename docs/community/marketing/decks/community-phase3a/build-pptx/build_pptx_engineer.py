#!/usr/bin/env python3
"""
Build OCM-Community-Engineer.pptx — 13-slide community deck (Option A v0.4:
"The Day a CVE Drops").

Story arc (locked, NARRATIVE-OPTION-A.md v0.4):

   1  HERO          It's 2am. A CVE drops.                  (single-line, whole gradient)
   2  PAIN          Your release isn't a thing. It's a scavenger hunt.
   3  NO NAME       The release has no name.                (delayed beat)
   4  PACK          Pack.                                   [constructor YAML]
   5  SIGN          What gets signed and travels.           [descriptor YAML]
   6  TRAVELS       Travels.                                [3 transport patterns]
   7  DEPLOY        Apply once. The controllers take over.  [Controllers box]
   8  COMPOSE       A product is a tree.                    [product → notes + postgres]
   9  BUMP          Bump the product. Everything follows.   [Day-2 diff + cascade]
  10  ODG           The scanner speaks libraries. OCM speaks components.
                                                            [side-by-side, shared digest]
  11  2AM REDO      It's 2am. You already know.             [4 cards w/ arrows]
  12  CLOSE         A release is a thing, not a scavenger hunt.   [CTA layout, signature]
  13  ADOPT         Two paths. Pick the one that fits Monday.     [Q&A backdrop]

ADR-0016 ownership annotations (slide 10): the image manifest is unchanged;
OCM publishes a side-car referrer manifest with artifactType
'application/vnd.ocm.software.ownership.v1+json' whose annotations carry
software.ocm.component.{name,version} and software.ocm.artifact.
Discoverable via OCI Referrers API (oras discover --artifact-type).

This script imports helpers from the architect deck's build script
(palette, banner, brand row, YAML block, layout helpers, sanity check)
so brand updates propagate automatically.

Usage:
    python3 build_pptx_engineer.py
"""
from __future__ import annotations

import sys
from pathlib import Path

SCRIPT_DIR = Path(__file__).resolve().parent
ARCHITECT_BUILD_DIR = (
    SCRIPT_DIR.parent.parent / "architect-phase2a" / "build-pptx"
)
sys.path.insert(0, str(ARCHITECT_BUILD_DIR))

import build_pptx_architect_external as arch  # noqa: E402

from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.util import Emu, Pt  # noqa: E402


# Convenience aliases — pulled from the architect build module
C = arch.C
px = arch.px
add_textbox = arch.add_textbox
add_brand_row = arch.add_brand_row
add_banner_full_bleed = arch.add_banner_full_bleed
add_yaml_block = arch.add_yaml_block
add_diagram = arch.add_diagram
add_footer_caption = arch.add_footer_caption
add_left_bullets = arch.add_left_bullets
set_text = arch.set_text
set_split_gradient_title = arch.set_split_gradient_title
set_blue_box_bullets = arch.set_blue_box_bullets
set_action_path_lines = arch.set_action_path_lines
delete_placeholder = arch.delete_placeholder
layouts_by_name = arch.layouts_by_name
open_template_as_pptx = arch.open_template_as_pptx
sanity_check = arch.sanity_check
THEME_DIR = arch.THEME_DIR
SLIDE_W_PX = arch.SLIDE_W_PX
SLIDE_H_PX = arch.SLIDE_H_PX

OUTPUT_PPTX = SCRIPT_DIR.parent / "OCM-Community-Engineer.pptx"


# =============================================================================
# Slide builders
# =============================================================================

def build_slide_1_hero(prs, layouts):
    """1 HERO — single line, whole-line gradient. 'Where is it running?'"""
    s = prs.slides.add_slide(layouts["Hero"])
    add_banner_full_bleed(s, THEME_DIR / "OCM-Banner.png")
    # Whole-line gradient: empty prefix, full title as 'noun'.
    set_split_gradient_title(s, 1, prefix="", noun="It's 2am. A CVE drops.")
    # Title-line-2 placeholder kept empty so layout doesn't error out.
    set_text(s, 2, "", color=C.WHITE, align_left=True)
    set_text(s, 3, "Where is it running?", color=C.CYAN)
    set_text(s, 4,
             "Open Component Model — open source, NeoNephos Foundation.",
             color=C.WHITE)
    add_brand_row(s)


def build_slide_2_pain(prs, layouts):
    """2 PAIN — Your release isn't a thing. It's a scavenger hunt.
    'Thirty minutes to find the pieces' promoted to a standalone line."""
    s = prs.slides.add_slide(layouts["Plain"])
    set_text(s, 1, "THE PAIN")
    set_text(s, 2, "Your release isn't a thing. It's a scavenger hunt.")
    set_blue_box_bullets(s, 10, [
        "The image is in GHCR. The chart is in another registry.",
        "The SBOM is on an S3 bucket someone set up two years ago.",
        "The Terraform that wired it together is in a fourth repo.",
        "The CVE doesn't care.",
    ])
    # Standalone "thirty minutes" line below the bullets — sets up the
    # slide-10 "thirty seconds" callback.
    tb, tf = add_textbox(s, 120, 900, SLIDE_W_PX - 240, 80)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Thirty minutes to find the pieces."
    r.font.name = "Aptos"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.italic = True
    r.font.color.rgb = C.BLUE


def build_slide_3_no_name(prs, layouts):
    """3 NO NAME — The release has no name. (Delayed beat — no answer yet.)"""
    s = prs.slides.add_slide(layouts["Plain / Compact"])
    set_text(s, 1, "THE FRAGMENTATION")
    set_text(s, 2, "The release has no name.")
    delete_placeholder(s, 10)

    # Three short body beats, centred, large.
    tb, tf = add_textbox(s, 120, 540, SLIDE_W_PX - 240, 280)
    for i, line in enumerate([
        "Every artifact has an identity.",
        "None of them name the release.",
        "",
        "Mirror the image — the reference changes.",
        "Mirror the chart — the reference changes.",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0) if i == 0 else Pt(12)
        r = p.add_run()
        r.text = line
        r.font.name = "Aptos"
        r.font.size = Pt(30)
        if line and not line.endswith("."):
            r.font.color.rgb = C.GREY_MID
        else:
            r.font.color.rgb = C.BLACK
        if i in (0, 1):
            r.font.bold = True
            r.font.color.rgb = C.BLUE_MID

    # Grey aside at the bottom
    add_footer_caption(
        s, 920,
        "Cosign signs each piece. None of them sign the release as one "
        "named, location-independent unit.",
        italic=True,
    )


def _build_yaml_two_pane(s, *, title_eyebrow, title_text,
                          yaml_lines, cli_text,
                          right_lines, right_lead=None):
    """Common layout for slides 4 and 5: title + eyebrow on top,
    YAML on the left, CLI + caption on the right."""
    set_text(s, 1, title_eyebrow)
    set_text(s, 2, title_text)
    delete_placeholder(s, 10)

    # YAML block on left
    add_yaml_block(s, x_px=120, y_px=540, w_px=1080, h_px=440,
                    yaml_lines=yaml_lines, font_size=17)

    # CLI on right (large, blue, monospace)
    cli_tb, cli_tf = add_textbox(s, 1240, 560, 560, 80)
    p = cli_tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = cli_text
    r.font.name = "Consolas"
    r.font.size = Pt(32)
    r.font.bold = True
    r.font.color.rgb = C.BLUE

    # Right caption lines (≥22pt for legibility)
    cap_y = 660
    cap_tb, cap_tf = add_textbox(s, 1240, cap_y, 560, 320)
    if right_lead is not None:
        p = cap_tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = right_lead
        r.font.name = "Aptos"
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.color.rgb = C.BLUE
        first_added = False
    else:
        first_added = True

    for i, line in enumerate(right_lines):
        if not first_added and i == 0:
            p = cap_tf.paragraphs[0]
            first_added = True
        else:
            p = cap_tf.add_paragraph()
            p.space_before = Pt(14)
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.name = "Aptos"
        r.font.size = Pt(22)
        r.font.color.rgb = C.BLACK


def build_slide_4_pack(prs, layouts):
    """4 PACK — single-word title. Constructor YAML left, ocm add cv right.
    The component name in the YAML is the reveal slide 3 deferred."""
    s = prs.slides.add_slide(layouts["Plain"])
    yaml_lines = [
        ("components:",                                       C.BLUE),
        ("  - name: github.com/acme.org/helloworld",          C.BLACK),
        ("    version: 1.0.0",                                C.BLACK),
        ("    provider:",                                     C.BLACK),
        ("      name: acme.org",                              C.BLACK),
        ("    resources:",                                    C.BLUE),
        ("      - name: image",                               C.BLACK),
        ("        type: ociImage",                            C.BLACK),
        ("        access:",                                   C.BLACK),
        ("          type: OCIImage/v1",                       C.BLACK),
        ("          imageReference: "
         "ghcr.io/stefanprodan/podinfo:6.9.1",                C.BLACK),
    ]
    _build_yaml_two_pane(
        s,
        title_eyebrow="PACK",
        title_text="Pack.",
        yaml_lines=yaml_lines,
        cli_text="$ ocm add cv",
        right_lines=[
            "Produces a CTF archive.",
            "Portable. OCI-compatible.",
            "Move it with ocm transfer.",
        ],
    )


def build_slide_5_sign(prs, layouts):
    """5 SIGN — descriptor YAML left, ocm sign cv right.
    Same shape as slide 4 (input → output mental model)."""
    s = prs.slides.add_slide(layouts["Plain"])
    # Trimmed descriptor — mirrors architect deck slide 6 pattern.
    yaml_lines = [
        ("component:",                                        C.BLUE),
        ("  name: github.com/acme.org/helloworld",            C.BLACK),
        ("  version: 1.0.0",                                  C.BLACK),
        ("  provider:",                                       C.BLACK),
        ("    name: acme.org",                                C.BLACK),
        ("  resources:",                                      C.BLUE),
        ("    - name: image",                                 C.BLACK),
        ("      type: ociImage",                              C.BLACK),
        ("      digest:",                                     C.BLACK),
        ("        hashAlgorithm: SHA-256",                    C.BLACK),
        ("        value: 70a2577d7b…",                        C.BLACK),
        ("      access:",                                     C.BLACK),
        ("        type: OCIImage/v1",                         C.BLACK),
        ("signatures:",                                       C.BLUE),
        ("  - name: default",                                 C.BLACK),
        ("    digest: { value: 70a2577d7b… }",                C.BLACK),
        ("    signature: { value: <PEM block> }",             C.BLACK),
    ]
    _build_yaml_two_pane(
        s,
        title_eyebrow="SIGN",
        title_text="What gets signed and travels.",
        yaml_lines=yaml_lines,
        cli_text="$ ocm sign cv",
        right_lines=[
            "OCM signs the descriptor,",
            "not the individual files.",
            "",
            "The descriptor carries digests",
            "for every artifact — one signature,",
            "tamper-evident over the release.",
        ],
    )


def build_slide_6_travels(prs, layouts):
    """6 TRAVELS — single-word title, three transport patterns."""
    s = prs.slides.add_slide(layouts["Plain / Compact"])
    set_text(s, 1, "TRANSPORT")
    set_text(s, 2, "Travels.")
    set_blue_box_bullets(s, 10, [
        "Registry → Registry. Public to private, cloud to cloud.",
        "Registry → CTF archive. The whole release in one filesystem.",
        "CTF → Registry. Behind the air gap. No callback upstream.",
    ])
    add_footer_caption(
        s, 920,
        "Same identity in OCI, S3, or on USB. Verify with the same key.",
        italic=True,
    )


def _draw_controllers_box(slide, *, x_px: int, y_px: int,
                            w_px: int, h_px: int):
    """Single 'OCM Controllers' box with four verbs and CR labels.

    Layout inside the box:
       OCM Controllers (centred header)
       ──────────────────────
       pull descriptor      Repository
       verify signature     Component
       resolve resources    Resource
       apply & reconcile    Deployer
    """
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  px(x_px), px(y_px),
                                  px(w_px), px(h_px))
    box.fill.solid()
    box.fill.fore_color.rgb = C.GREY_SOFT
    box.line.color.rgb = C.BLUE
    box.line.width = Pt(2.0)
    # Remove default text frame, we add our own textboxes inside.
    box.text_frame.text = ""

    # Header
    head_tb, head_tf = add_textbox(slide, x_px + 30, y_px + 24,
                                    w_px - 60, 60)
    p = head_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "OCM Controllers"
    r.font.name = "Aptos"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = C.BLUE

    # Horizontal rule under header
    rule = slide.shapes.add_connector(1,
        px(x_px + 60), px(y_px + 90),
        px(x_px + w_px - 60), px(y_px + 90))
    rule.line.color.rgb = C.BLUE
    rule.line.width = Pt(1.0)

    # Four rows: verb (left, large blue) + CR (right, grey)
    rows = [
        ("pull descriptor",   "Repository"),
        ("verify signature",  "Component"),
        ("resolve resources", "Resource"),
        ("apply & reconcile", "Deployer"),
    ]
    row_h = (h_px - 130) // len(rows)
    row_top = y_px + 110
    for i, (verb, cr) in enumerate(rows):
        row_y = row_top + i * row_h
        # Verb on the left
        v_tb, v_tf = add_textbox(slide, x_px + 60, row_y,
                                  (w_px // 2) + 100, row_h)
        v_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = v_tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = verb
        r.font.name = "Aptos"
        r.font.size = Pt(24)
        r.font.color.rgb = C.BLACK
        # CR on the right, in light grey
        c_tb, c_tf = add_textbox(slide, x_px + (w_px // 2) + 200, row_y,
                                  (w_px // 2) - 260, row_h)
        c_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c_tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = cr
        r.font.name = "Aptos"
        r.font.size = Pt(20)
        r.font.italic = True
        r.font.color.rgb = C.GREY_MID


def build_slide_7_deploy(prs, layouts):
    """7 DEPLOY — Apply once. The controllers take over.
    Single box, four verbs, CR labels in light grey."""
    s = prs.slides.add_slide(layouts["Plain"])
    set_text(s, 1, "DEPLOY")
    set_text(s, 2, "Apply once. The controllers take over.")
    delete_placeholder(s, 10)

    # Component CR label above the box
    in_tb, in_tf = add_textbox(s, 0, 540, SLIDE_W_PX, 50)
    p = in_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Component CR"
    r.font.name = "Consolas"
    r.font.size = Pt(20)
    r.font.italic = True
    r.font.color.rgb = C.GREY_MID

    # Centred controllers box
    box_w = 800
    box_h = 360
    box_x = (SLIDE_W_PX - box_w) // 2
    _draw_controllers_box(s, x_px=box_x, y_px=605,
                            w_px=box_w, h_px=box_h)

    # Cluster label below the box
    out_tb, out_tf = add_textbox(s, 0, 985, SLIDE_W_PX, 50)
    p = out_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Cluster"
    r.font.name = "Aptos"
    r.font.size = Pt(20)
    r.font.italic = True
    r.font.color.rgb = C.GREY_MID


def build_slide_8_compose(prs, layouts):
    """8 COMPOSE — A product is a tree.
    Composition tree: product → notes + postgres. Defines the noun
    'product component' before slide 9 says 'bump the product'."""
    s = prs.slides.add_slide(layouts["Plain"])
    set_text(s, 1, "COMPOSE")
    set_text(s, 2, "A product is a tree.")
    delete_placeholder(s, 10)

    # Tree, centred horizontally and vertically in the body region.
    # Parent in brand blue + bold; children in brand grey, indented.
    tree_lines = [
        ("github.com/acme/sovereign/product : 1.0.0",     True),   # parent
        ("",                                              False),
        ("   ├── github.com/acme/sovereign/notes    : 1.0.0", False),
        ("   └── github.com/acme/sovereign/postgres : 1.0.0", False),
    ]
    tb, tf = add_textbox(s, 200, 580, SLIDE_W_PX - 400, 300)
    for i, (line, is_parent) in enumerate(tree_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0) if i == 0 else Pt(12)
        r = p.add_run()
        r.text = line
        r.font.name = "Consolas"
        r.font.size = Pt(28)
        if is_parent:
            r.font.bold = True
            r.font.color.rgb = C.BLUE
        elif line.strip().startswith("├") or line.strip().startswith("└"):
            r.font.color.rgb = C.GREY_MID
        else:
            r.font.color.rgb = C.BLACK

    # Two-line caption below the tree
    cap_tb, cap_tf = add_textbox(s, 120, 920, SLIDE_W_PX - 240, 100)
    caption_lines = [
        "Leaf components carry resources. The product references them.",
        "One name, one signature, covers the whole tree.",
    ]
    for i, line in enumerate(caption_lines):
        p = cap_tf.paragraphs[0] if i == 0 else cap_tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0) if i == 0 else Pt(8)
        r = p.add_run()
        r.text = line
        r.font.name = "Aptos"
        r.font.size = Pt(22)
        r.font.italic = (i == 1)
        r.font.color.rgb = C.GREY_MID if i == 1 else C.BLACK


def build_slide_9_bump(prs, layouts):
    """9 BUMP — Day 2 only. YAML diff + kro/OCM/Flux cascade.
    Composition tree moved to slide 8 — this slide now does ONE job."""
    s = prs.slides.add_slide(layouts["Plain"])
    set_text(s, 1, "DAY 2")
    set_text(s, 2, "Bump the product. Everything follows.")
    delete_placeholder(s, 10)

    # Left half — YAML diff, big and centred-vertical
    yaml_lines = [
        ("spec:",                                  C.BLUE),
        ("  version: 1.1.0   # was: 1.0.0",        C.BLACK),
    ]
    add_yaml_block(s, x_px=120, y_px=560, w_px=900, h_px=200,
                    yaml_lines=yaml_lines, font_size=28)

    # Right half — numbered cascade
    cas_tb, cas_tf = add_textbox(s, 1080, 540, 720, 420)
    cascade = [
        ("1.", "kro re-renders the RGD with the new spec.version."),
        ("2.", "Component CR's semver updates; controllers resolve "
                "and verify the new version."),
        ("3.", "Resource CRs resolve new digests for the child "
                "charts and images."),
        ("4.", "Flux HelmReleases roll the new artifacts."),
    ]
    for i, (num, text) in enumerate(cascade):
        p = cas_tf.paragraphs[0] if i == 0 else cas_tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0) if i == 0 else Pt(14)
        r = p.add_run()
        r.text = num + "  "
        r.font.name = "Aptos"
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.color.rgb = C.BLUE
        r = p.add_run()
        r.text = text
        r.font.name = "Aptos"
        r.font.size = Pt(20)
        r.font.color.rgb = C.BLACK

    # Punch line below — full width, blue, prominent
    tb, tf = add_textbox(s, 120, 920, SLIDE_W_PX - 240, 80)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "You changed one line. The cluster did the rest."
    r.font.name = "Aptos"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.italic = True
    r.font.color.rgb = C.BLUE


# Shared digest string for slides 10 and 11 — visual continuity.
# When the audience sees the same digest twice, the join lands without
# needing words.
_SHARED_IMAGE_REF = "eu.gcr.io/acme/notes"
_SHARED_DIGEST = "sha256:70a2577d…"
_SHARED_FULL_REF = f"{_SHARED_IMAGE_REF}@{_SHARED_DIGEST}"


def _emit_console_lines(tf, lines, *, font_size=18):
    """Helper: emit a list of (text, color, bold) tuples as Consolas
    paragraphs into a text frame. Lines with text=='' produce a blank
    paragraph for spacing."""
    for i, item in enumerate(lines):
        text, color, bold = item if len(item) == 3 else (*item, False)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0) if i == 0 else Pt(4)
        if not text:
            continue
        r = p.add_run()
        r.text = text
        r.font.name = "Consolas"
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.color.rgb = color


def build_slide_10_odg(prs, layouts):
    """10 ODG — side-by-side: scanner output (left) / oras discover (right).
    Same digest string highlighted in brand blue on BOTH sides — the
    audience's eye traces the join.

    Mechanism per ADR-0016: image manifest is unchanged; OCM publishes a
    separate ownership-referrer manifest with artifactType
    application/vnd.ocm.software.ownership.v1+json; that referrer's
    annotations carry software.ocm.component.{name,version} and
    software.ocm.artifact. Discoverable via the OCI Referrers API,
    concretely via `oras discover --artifact-type ...`."""
    s = prs.slides.add_slide(layouts["Plain"])
    set_text(s, 1, "ODG")
    set_text(s, 2, "The scanner speaks libraries. OCM speaks components.")
    delete_placeholder(s, 10)

    # Two columns side-by-side
    col_y = 540
    col_h = 420
    col_w = 900
    gap = 60
    total_w = 2 * col_w + gap
    start_x = (SLIDE_W_PX - total_w) // 2  # ≈ 30

    # ---------- LEFT COLUMN — WHAT THE SCANNER REPORTS ----------
    # Header
    lh_tb, lh_tf = add_textbox(s, start_x, col_y, col_w, 50)
    p = lh_tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "WHAT THE SCANNER REPORTS"
    r.font.name = "Aptos"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = C.BLUE
    # Rule under header
    rule = s.shapes.add_connector(1, px(start_x), px(col_y + 50),
                                   px(start_x + col_w), px(col_y + 50))
    rule.line.color.rgb = C.BLUE
    rule.line.width = Pt(1.25)
    # Body — scanner output style
    lb_tb, lb_tf = add_textbox(s, start_x, col_y + 70, col_w, col_h - 70)
    _emit_console_lines(lb_tf, [
        ("CVE-2026-XXXX in libfoo",        C.BLACK,    True),
        ("",                                C.BLACK,    False),
        (_SHARED_IMAGE_REF,                 C.BLUE,     False),
        ("@" + _SHARED_DIGEST,              C.BLUE,     True),
        ("",                                C.BLACK,    False),
        ("package · version · digest",      C.GREY_MID, False),
    ], font_size=20)

    # ---------- RIGHT COLUMN — WHAT ODG DISCOVERS ----------
    rx = start_x + col_w + gap
    # Header
    rh_tb, rh_tf = add_textbox(s, rx, col_y, col_w, 50)
    p = rh_tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "WHAT ODG DISCOVERS"
    r.font.name = "Aptos"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = C.BLUE
    rule = s.shapes.add_connector(1, px(rx), px(col_y + 50),
                                   px(rx + col_w), px(col_y + 50))
    rule.line.color.rgb = C.BLUE
    rule.line.width = Pt(1.25)
    # Body — oras command + result
    rb_tb, rb_tf = add_textbox(s, rx, col_y + 70, col_w, col_h - 70)
    _emit_console_lines(rb_tf, [
        ("$ oras discover \\",                                  C.BLACK,    True),
        ("    " + _SHARED_FULL_REF + " \\",                     C.BLUE,     True),
        ("    --artifact-type \\",                              C.BLACK,    False),
        ("      application/vnd.ocm.software.ownership.v1+json",
                                                                 C.BLACK,    False),
        ("",                                                     C.BLACK,    False),
        ("→ annotations:",                                       C.BLUE,     True),
        ("    software.ocm.component.name",                      C.BLUE,     False),
        ("      = github.com/acme/notes",                        C.BLACK,    False),
        ("    software.ocm.component.version",                   C.BLUE,     False),
        ("      = 1.0.0",                                        C.BLACK,    False),
        ("    software.ocm.artifact",                            C.BLUE,     False),
        ("      = { name: notes-image, kind: resource }",        C.BLACK,    False),
    ], font_size=15)

    # Footer — the architectural insight
    foot_tb, foot_tf = add_textbox(s, 120, 990, SLIDE_W_PX - 240, 50)
    p = foot_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "The image is unchanged. Ownership rides in a side-car."
    r.font.name = "Aptos"
    r.font.size = Pt(20)
    r.font.italic = True
    r.font.color.rgb = C.GREY_MID


def build_slide_11_2am_redo(prs, layouts):
    """11 2AM REDO — four stacked cards with arrows. Mirrors architect
    slide 11 (Deploy) card-chain visual. The digest from slide 10
    propagates through the cards — visual continuity."""
    s = prs.slides.add_slide(layouts["Plain"])
    set_text(s, 1, "AT 2AM")
    set_text(s, 2, "It's 2am. You already know.")
    delete_placeholder(s, 10)

    # Four stacked cards. Slide body region: y ≈ 540 → 960 (420px height).
    # 4 cards × 90px + 3 gaps × 30px arrow region = 450px. Close.
    card_x = 200
    card_w = SLIDE_W_PX - 400
    card_h = 95
    gap = 25
    first_y = 540

    cards = [
        ("SCANNER",
         [("CVE-2026-XXXX in libfoo", False),
          ("found in image " + _SHARED_IMAGE_REF + "@" + _SHARED_DIGEST, True)]),
        ("ODG",
         [("oras discover the image with the OCM ownership artifactType", False),
          ("→ component github.com/acme/notes : 1.0.0 (resource notes-image)", False)]),
        ("DASHBOARD",
         [("owner = team-notes · env = eu-prod-12 · triaged = no", False)]),
        ("ACTION",
         [("Page team-notes. Patch on the shelf. Bump product. Done.", False)]),
    ]

    for i, (label, body_lines) in enumerate(cards):
        y = first_y + i * (card_h + gap)
        # Card outline
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   px(card_x), px(y),
                                   px(card_w), px(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = C.GREY_SOFT
        card.line.color.rgb = C.BLUE
        card.line.width = Pt(1.5)
        # Clear card's default text frame (we use our own)
        card.text_frame.text = ""
        # Label (small, brand blue, top-left)
        lbl_tb, lbl_tf = add_textbox(s, card_x + 30, y + 12,
                                      200, 28)
        p = lbl_tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = label
        r.font.name = "Aptos"
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = C.BLUE
        # Body
        body_tb, body_tf = add_textbox(s, card_x + 30, y + 38,
                                        card_w - 60, card_h - 40)
        for j, (line, is_digest_line) in enumerate(body_lines):
            p = body_tf.paragraphs[0] if j == 0 else body_tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(0) if j == 0 else Pt(2)
            # Highlight the digest line in brand blue for visual continuity
            # with slide 10.
            if is_digest_line:
                # Split into prefix + image_ref + digest so we can colour
                # the digest in brand blue while leaving prefix in black.
                prefix = "found in image "
                r = p.add_run()
                r.text = prefix
                r.font.name = "Aptos"
                r.font.size = Pt(18)
                r.font.color.rgb = C.BLACK
                r = p.add_run()
                r.text = _SHARED_IMAGE_REF + "@" + _SHARED_DIGEST
                r.font.name = "Consolas"
                r.font.size = Pt(18)
                r.font.bold = True
                r.font.color.rgb = C.BLUE
            else:
                r = p.add_run()
                r.text = line
                r.font.name = "Aptos"
                r.font.size = Pt(18)
                r.font.color.rgb = C.BLACK

        # Arrow between cards (not after the last one)
        if i < len(cards) - 1:
            arrow_x = SLIDE_W_PX // 2
            arrow_y_top = y + card_h + 4
            arrow_y_bot = y + card_h + gap - 4
            line_shape = s.shapes.add_connector(1,
                px(arrow_x), px(arrow_y_top),
                px(arrow_x), px(arrow_y_bot - 6))
            line_shape.line.color.rgb = C.BLUE_MID
            line_shape.line.width = Pt(2.5)
            # Arrowhead (downward triangle)
            tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                px(arrow_x - 8), px(arrow_y_bot - 8),
                px(16), px(10))
            tri.fill.solid()
            tri.fill.fore_color.rgb = C.BLUE_MID
            tri.line.fill.background()
            tri.rotation = 180

    # Footer — Thirty seconds callback (bookend with slide 2)
    foot_tb, foot_tf = add_textbox(s, 120, 1000, SLIDE_W_PX - 240, 50)
    p = foot_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Thirty seconds. Coffee's still warm."
    r.font.name = "Aptos"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.italic = True
    r.font.color.rgb = C.BLUE


def build_slide_12_close(prs, layouts):
    """12 CLOSE — CTA brand layout, signature line dominant, three doors.
    Speaker closes the talk here."""
    s = prs.slides.add_slide(layouts["CTA"])
    set_text(s, 1, "A release is a thing, not a scavenger hunt.",
              color=C.WHITE)
    set_action_path_lines(s, 2, [
        ("Star and watch",  "github.com/open-component-model"),
        ("Try the tutorial", "ocm.software/docs/getting-started"),
        ("Talk to us",       "community channels on the website"),
    ])
    add_brand_row(s)


def build_slide_13_adopt(prs, layouts):
    """13 ADOPT — two paths, Q&A backdrop. Mirrors architect slide 14 layout."""
    s = prs.slides.add_slide(layouts["Plain / Compact"])
    set_text(s, 1, "ADOPT")
    set_text(s, 2, "Two paths. Pick the one that fits Monday.")
    delete_placeholder(s, 10)

    col_y = 560
    col_h = 380
    col_w = 820
    gap = 80
    total_w = 2 * col_w + gap
    start_x = (SLIDE_W_PX - total_w) // 2

    columns = [
        ("FROM ZERO — 10 MINUTES",
         ["Install the ocm CLI.",
          "Write a 10-line component-constructor.yaml.",
          "ocm add cv → sign → verify → transfer → verify.",
          "Signed at source, verified at destination, "
          "transferable anywhere."]),
        ("ON YOUR PLATFORM — WIRE IN",
         ["Install OCM controllers; point at your registry.",
          "Apply a Component CR — verified, reconciling.",
          "Compose with Argo, Flux as you already do.",
          "Your next release ships as an OCM component."]),
    ]
    for i, (header, lines) in enumerate(columns):
        x = start_x + i * (col_w + gap)
        head_tb, head_tf = add_textbox(s, x, col_y, col_w, 56)
        hp = head_tf.paragraphs[0]
        hp.alignment = PP_ALIGN.LEFT
        hr = hp.add_run()
        hr.text = header
        hr.font.name = "Aptos"
        hr.font.size = Pt(22)
        hr.font.bold = True
        hr.font.color.rgb = C.BLUE
        rule = s.shapes.add_connector(1, px(x), px(col_y + 56),
                                       px(x + col_w), px(col_y + 56))
        rule.line.color.rgb = C.BLUE
        rule.line.width = Pt(1.25)
        body_tb, body_tf = add_textbox(s, x, col_y + 72, col_w, col_h - 72)
        for j, line in enumerate(lines):
            p = body_tf.paragraphs[0] if j == 0 else body_tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(0) if j == 0 else Pt(10)
            r = p.add_run()
            r.text = line
            r.font.name = "Aptos"
            r.font.size = Pt(20)
            r.font.color.rgb = C.BLACK


# =============================================================================
# Main
# =============================================================================

def main():
    prs = open_template_as_pptx()
    layouts = layouts_by_name(prs)

    expected = {"Hero", "CTA", "Content / 3-Column",
                "Content / Diagram", "Content / Diagram Compact",
                "Content / Tiles", "Content / 2-Column",
                "Section Divider", "Plain", "Plain / Compact"}
    missing = expected - set(layouts)
    if missing:
        sys.exit(f"template missing expected layouts: {missing}")

    build_slide_1_hero(prs, layouts)
    build_slide_2_pain(prs, layouts)
    build_slide_3_no_name(prs, layouts)
    build_slide_4_pack(prs, layouts)
    build_slide_5_sign(prs, layouts)
    build_slide_6_travels(prs, layouts)
    build_slide_7_deploy(prs, layouts)
    build_slide_8_compose(prs, layouts)
    build_slide_9_bump(prs, layouts)
    build_slide_10_odg(prs, layouts)
    build_slide_11_2am_redo(prs, layouts)
    build_slide_12_close(prs, layouts)
    build_slide_13_adopt(prs, layouts)

    sanity_check(prs)
    prs.save(str(OUTPUT_PPTX))
    print(f"Wrote {OUTPUT_PPTX} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()

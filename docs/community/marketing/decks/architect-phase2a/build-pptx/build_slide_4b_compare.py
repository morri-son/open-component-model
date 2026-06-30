"""Generate slide 4b ("How OCM compares") as a standalone single-slide .pptx.

This is a Phase 2B addition to the architect deck. It produces ONE slide using
the same OCM-Master.potx template, palette, and Plain layout as the rest of
the deck — so when you drag-insert this slide into the SharePoint copy between
current slides 4 and 5, it inherits the deck's visual idiom.

Rendering approach: PowerPoint native table object (shapes.add_table) so the
matrix gets row dividers, column structure, and cell padding for free — no
manual textbox positioning, no orphan placeholder issues.

Run from the build-pptx/ directory:

    python build_slide_4b_compare.py

Output: ../OCM-Story-Architect-External-Slide-4b.pptx (one slide only).
"""

from __future__ import annotations

import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt
from lxml import etree


# --- Paths & geometry mirroring build_pptx_architect_external.py -----------

SCRIPT_DIR = Path(__file__).resolve().parent
DECK_DIR = SCRIPT_DIR.parent
POTX_PATH = DECK_DIR / "OCM-Master.potx"
RASTER_DIR = SCRIPT_DIR / "_raster"
RASTER_DIR.mkdir(exist_ok=True)
OUTPUT_PPTX = DECK_DIR / "OCM-Story-Architect-External-Slide-4b.pptx"

SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
PX = 9525


def px(n: float) -> Emu:
    return Emu(int(n * PX))


# --- Brand palette ---------------------------------------------------------

class C:
    BLUE       = RGBColor(0x0F, 0x6B, 0xFF)
    BLUE_MID   = RGBColor(0x0A, 0x3A, 0x99)
    BLUE_SOFT  = RGBColor(0xE8, 0xEF, 0xFF)
    GREY_MID   = RGBColor(0x6B, 0x72, 0x80)
    GREY_SOFT  = RGBColor(0xF3, 0xF4, 0xF6)
    GREY_RULE  = RGBColor(0xE5, 0xE7, 0xEB)
    BLACK      = RGBColor(0x00, 0x00, 0x00)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)


# --- Template loader (same trick as the main build script) -----------------

def open_template_as_pptx() -> Presentation:
    if not POTX_PATH.exists():
        sys.exit(f"{POTX_PATH} not found — expected the OCM-Master template.")
    tmp_pptx = RASTER_DIR / "_potx_loaded_4b.pptx"
    with zipfile.ZipFile(POTX_PATH, "r") as src:
        data = {n: src.read(n) for n in src.namelist()}
    ct = data["[Content_Types].xml"].decode("utf-8")
    ct = ct.replace(
        "presentationml.template.main+xml",
        "presentationml.presentation.main+xml",
    )
    data["[Content_Types].xml"] = ct.encode("utf-8")
    with zipfile.ZipFile(tmp_pptx, "w", zipfile.ZIP_DEFLATED) as out:
        for name, blob in data.items():
            out.writestr(name, blob)
    return Presentation(str(tmp_pptx))


# --- Helpers ---------------------------------------------------------------

def find_placeholder(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def delete_placeholder(slide, idx: int):
    ph = find_placeholder(slide, idx)
    if ph is None:
        return
    el = ph._element
    el.getparent().remove(el)


def set_text(slide, idx: int, text: str):
    ph = find_placeholder(slide, idx)
    if ph is None:
        return
    ph.text_frame.text = text


def set_speaker_notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


def style_cell_text(cell, text: str, *, size: int, color: RGBColor,
                    bold: bool = False, italic: bool = False,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
                    font: str = "Aptos"):
    """Write text into a table cell and apply OCM-deck typography."""
    cell.vertical_anchor = anchor
    tf = cell.text_frame
    tf.margin_left = Emu(int(0.15 * 914400))   # ~0.15"
    tf.margin_right = Emu(int(0.15 * 914400))
    tf.margin_top = Emu(int(0.05 * 914400))
    tf.margin_bottom = Emu(int(0.05 * 914400))
    tf.word_wrap = True
    # Clear existing content
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def set_cell_fill(cell, color: RGBColor):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def clear_cell_borders(cell):
    """Remove all four borders on a cell (PowerPoint table cells inherit
    style borders from the table style; we want our own row-rule discipline)."""
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("lnL", "lnR", "lnT", "lnB"):
        # Remove any existing border element
        for el in tcPr.findall(qn(f"a:{tag}")):
            tcPr.remove(el)
        # Add a "no fill" border element to override the style
        ln = etree.SubElement(tcPr, qn(f"a:{tag}"))
        ln.set("w", "0")
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")
        noFill = etree.SubElement(ln, qn("a:noFill"))


def set_cell_bottom_border(cell, color: RGBColor, weight_emu: int = 9525):
    """Apply a single bottom border at the given colour + weight (default 1pt)."""
    clear_cell_borders(cell)
    tcPr = cell._tc.get_or_add_tcPr()
    # Remove the existing noFill bottom border we just added
    for el in tcPr.findall(qn("a:lnB")):
        tcPr.remove(el)
    ln = etree.SubElement(tcPr, qn("a:lnB"))
    ln.set("w", str(weight_emu))
    ln.set("cap", "flat")
    ln.set("cmpd", "sng")
    ln.set("algn", "ctr")
    solidFill = etree.SubElement(ln, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2]))


def set_cell_top_border(cell, color: RGBColor, weight_emu: int = 9525):
    """Apply a single top border at the given colour + weight. Used to mark
    the OCM row's visual separation from the per-artifact group above."""
    tcPr = cell._tc.get_or_add_tcPr()
    # Remove any existing top border (including any noFill from clear_cell_borders)
    for el in tcPr.findall(qn("a:lnT")):
        tcPr.remove(el)
    ln = etree.SubElement(tcPr, qn("a:lnT"))
    ln.set("w", str(weight_emu))
    ln.set("cap", "flat")
    ln.set("cmpd", "sng")
    ln.set("algn", "ctr")
    solidFill = etree.SubElement(ln, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2]))


def disable_table_style(table):
    """Strip the built-in table style so our manually-set fills and borders
    aren't overridden. Sets firstRow/bandRow/etc. flags to false and removes
    the styleId reference."""
    tbl = table._tbl
    tblPr = tbl.find(qn("a:tblPr"))
    if tblPr is None:
        return
    # Turn off all banding / first-row / first-col special styling
    for attr in ("firstRow", "firstCol", "lastRow", "lastCol",
                 "bandRow", "bandCol"):
        tblPr.set(attr, "0")
    # Remove the tableStyleId child (which references the master's style)
    for el in tblPr.findall(qn("a:tableStyleId")):
        tblPr.remove(el)


# --- Slide builder ---------------------------------------------------------

SPEAKER_NOTES_4B = (
    "Set the comparative anchor an architect-track audience expects. Each tool in the room operates on a different unit; OCM operates one level up.\n"
    "• cosign / sigstore - signs one OCI artifact. Strong per-image trust. Doesn't bundle. Doesn't travel across registries without re-sign or `cosign copy`. OCM uses Sigstore as one of its signing schemes for the component descriptor.\n"
    "• SLSA / in-toto - attests the build that produced one artifact. Provenance, not bundling. Not natively air-gap; needs a separate transport story. OCM carries SLSA/in-toto attestations as resources inside the component.\n"
    "• SBOM / OCI 1.1 referrers - inventories one artifact's contents and attaches it to that artifact's digest. Discovery, not bundling. Doesn't span a multi-artifact release. OCM carries SBOMs as resources; the descriptor names which SBOM belongs to which artifact.\n"
    "• OCM - signs THE COMPONENT, a named versioned bundle of artifacts plus access paths. One signature covers every digest. Location-independent: access fields rewritten on transfer; signature still verifies. Air-gap native: CTF round-trip with no callback to source.\n"
    "The 'partial' cells are calibrated honesty - SLSA attestations CAN travel with their subject if you choose to, OCI 1.1 referrers ARE digest-addressable so partially location-independent. We don't overclaim.\n"
    "Close the slide with the band line: 'OCM rides on top. It doesn't replace the per-artifact tools - it adds the release-level envelope they don't.' The next slide (Constructor) shows what that envelope looks like."
)


# Matrix data — single source of truth.
HEADERS = ["", "WHAT IT SIGNS", "LOCATION-INDEPENDENT", "AIR-GAP NATIVE"]
ROWS = [
    ("cosign / sigstore",   "one OCI artifact",         "no",       "no"),
    ("SLSA / in-toto",      "one build's provenance",   "no",       "partial"),
    ("SBOM / OCI 1.1 refs", "one artifact's contents",  "partial",  "no"),
    ("OCM",                 "a component (the bundle)", "yes",      "yes"),
]


def build_slide_4b(prs):
    layouts = {layout.name: layout for layout in prs.slide_layouts}
    layout = layouts.get("Plain")
    if layout is None:
        sys.exit("template missing 'Plain' layout")

    s = prs.slides.add_slide(layout)

    # Use the template's eyebrow + title placeholders (inherit master fonts).
    set_text(s, 1, "HOW OCM COMPARES")
    set_text(s, 2, "Composes with what's there.")

    # Delete the Body placeholder so "Text hinzufügen" doesn't render.
    delete_placeholder(s, 10)

    # --- Table -------------------------------------------------------------
    # 5 rows (1 header + 4 data) x 4 columns (label + 3 metrics).
    # Position: y=440..960, leaving caption space below.
    rows = 1 + len(ROWS)
    cols = 4
    table_left = px(120)
    table_top = px(440)
    table_w = px(SLIDE_W_PX - 240)
    table_h = px(420)

    shape = s.shapes.add_table(rows, cols, table_left, table_top, table_w, table_h)
    table = shape.table
    disable_table_style(table)

    # Column widths: label column wider, value columns narrower.
    # Total ~1680px = label(420) + signs(540) + loc(380) + air(340).
    table.columns[0].width = px(420)
    table.columns[1].width = px(540)
    table.columns[2].width = px(380)
    table.columns[3].width = px(340)

    # Row heights: header thinner, data rows taller for readability.
    table.rows[0].height = px(70)
    for i in range(1, rows):
        table.rows[i].height = px(80)

    # --- Header row -------------------------------------------------------
    for c_idx, label in enumerate(HEADERS):
        cell = table.cell(0, c_idx)
        clear_cell_borders(cell)
        set_cell_fill(cell, C.WHITE)
        # Header cells: mid-blue all-caps, bold, ~Pt18
        style_cell_text(cell, label, size=18, color=C.BLUE_MID, bold=True,
                        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM)
        # Brand-blue bottom border on the whole header row
        set_cell_bottom_border(cell, C.BLUE, weight_emu=19050)  # 2pt

    # --- Data rows --------------------------------------------------------
    for r_idx, row in enumerate(ROWS, start=1):
        is_ocm = (row[0] == "OCM")
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            clear_cell_borders(cell)
            # All rows on white. OCM row gets brand-blue bold type — colour
            # signals "this is us" without filling the row. The visual
            # divider for the OCM row is a 1pt brand-blue TOP border (added
            # below the loop) — that separates the per-artifact group above
            # from OCM below without shouting.
            set_cell_fill(cell, C.WHITE)
            if is_ocm:
                text_color = C.BLUE
                bold = True
            elif val == "partial":
                # Calibrated honesty — grey-mid, same weight as yes/no.
                text_color = C.GREY_MID
                bold = False
            else:
                text_color = C.BLACK
                bold = False
            size = 22 if c_idx == 0 else 20
            style_cell_text(cell, val, size=size, color=text_color,
                            bold=bold, italic=False, align=PP_ALIGN.LEFT,
                            anchor=MSO_ANCHOR.MIDDLE)
            # Bottom rule: light grey between data rows, none on the last
            # row (OCM) — the OCM row's separation comes from its TOP rule,
            # set immediately after this loop.
            if not is_ocm:
                set_cell_bottom_border(cell, C.GREY_RULE, weight_emu=6350)  # ~0.67pt

    # Top border on the OCM row — 1pt brand-blue rule that visually
    # separates the per-artifact group above from OCM below. Replaces the
    # full-row blue fill we used before, which read as a vendor-matrix
    # highlight bar.
    for c_idx in range(cols):
        cell = table.cell(rows - 1, c_idx)
        set_cell_top_border(cell, C.BLUE, weight_emu=12700)  # 1pt

    # --- Caption below the table -----------------------------------------
    # Caption sits ~80px below the table, mid-blue, single line.
    cap_top = table_top + table_h + px(60)
    cap = s.shapes.add_textbox(table_left, cap_top, table_w, px(60))
    cap.text_frame.word_wrap = True
    cap.text_frame.margin_left = cap.text_frame.margin_right = 0
    cap.text_frame.margin_top = cap.text_frame.margin_bottom = 0
    p = cap.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = ("OCM rides on top. It doesn't replace the per-artifact tools — "
              "it adds the release-level envelope they don't.")
    r.font.name = "Aptos"
    r.font.size = Pt(20)
    r.font.bold = False
    r.font.color.rgb = C.BLUE_MID

    set_speaker_notes(s, SPEAKER_NOTES_4B)


def main():
    prs = open_template_as_pptx()

    # Remove any sample slides the template might carry.
    while len(prs.slides) > 0:
        slide_id = prs.slides._sldIdLst[0]
        prs.slides._sldIdLst.remove(slide_id)

    build_slide_4b(prs)
    prs.save(str(OUTPUT_PPTX))
    print(f"Wrote {OUTPUT_PPTX} ({len(prs.slides)} slide)")


if __name__ == "__main__":
    main()


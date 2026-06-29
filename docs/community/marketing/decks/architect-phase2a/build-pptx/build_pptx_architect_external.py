#!/usr/bin/env python3
"""
Build OCM-Story-Architect-External.pptx — 15-slide trunk + appendix.

Trunk story arc (locked, autonomous-loop handoff §4):

  1  Pain          You ship pieces. / Nothing carries the release.
  2  Cause         In every existing tool, identity is bound to location.
  3  Insight       Identity that travels with the artifact.
  4  Positioning   One wrapper. All artifacts. Signed once.
  5  Constructor   What you write. (YAML)
  6  Descriptor    What gets signed and travels. (YAML)
  7  Overview      THE FOUR MOVES — Pack · Sign · Transport · Deploy.
  8  Compose       Service carries resources. Product carries references.
  9  Sign          Same signed object. Three trust models.
 10  Transport     Three patterns, one command.
 11  Deploy        Repository → Component → Resource → Deployer.
 12  Composition   One product. Three components. One line to upgrade.
 13  Adoption      Two paths to a first OCM component in production.
 14  What's sharp  Three honest edges.
 15  CTA           Evaluate · Pilot · Engage.

Appendix (pull-on-demand, not in main render order narration):

 16  Replication   Alongside the chain. Not within it.

No warm-ups, no appendix mentions in the main 15-slide narration,
no hidden trademarks (the user copies trademark slides in from the
exec deck after the loop completes — handoff §2).

Usage:
    .venv/bin/python build_pptx_architect_external.py
"""
from __future__ import annotations

import shutil
import subprocess
import sys
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

from speaker_notes import SPEAKER_NOTES


# -----------------------------------------------------------------------------
# Paths
# -----------------------------------------------------------------------------

SCRIPT_DIR = Path(__file__).resolve().parent
DECK_DIR = SCRIPT_DIR.parent
DIAGRAMS_DIR = DECK_DIR / "diagrams"
EXEC_DIAGRAMS_DIR = DECK_DIR.parent / "exec-phase1" / "diagrams"
EXEC_BUILD_DIR = DECK_DIR.parent / "exec-phase1" / "build-pptx"
ASSETS_DIR = DECK_DIR.parent.parent / "assets"
THEME_DIR = DECK_DIR / "theme"
RASTER_DIR = SCRIPT_DIR / "_raster"

# Reuse the exec deck's native "Pack · Sign · Transport · Deploy" rendering.
# The slide-7 PPT (which the user crafted in PowerPoint and we extracted)
# matches this function's output 1:1, with the same icons + sovereign-cloud
# target glyph. By calling it directly we keep one source of truth.
if str(EXEC_BUILD_DIR) not in sys.path:
    sys.path.insert(0, str(EXEC_BUILD_DIR))
from slide_6_native import add_pack_sign_transport_deploy_native

POTX_PATH = DECK_DIR / "OCM-Master.potx"
OUTPUT_PPTX = DECK_DIR / "OCM-Story-Architect-External.pptx"

RASTER_DIR.mkdir(exist_ok=True)


# -----------------------------------------------------------------------------
# Slide geometry — 16:9 @ 1920x1080
# -----------------------------------------------------------------------------

SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
PX = 9525  # 1 px in EMU at 96 dpi

def px(n: float) -> Emu:
    return Emu(int(n * PX))


# -----------------------------------------------------------------------------
# OCM brand palette
# -----------------------------------------------------------------------------

class C:
    BLUE       = RGBColor(0x0F, 0x6B, 0xFF)
    BLUE_MID   = RGBColor(0x0A, 0x3A, 0x99)
    CYAN       = RGBColor(0x5C, 0xD6, 0xFF)
    GREY_MID   = RGBColor(0x6B, 0x72, 0x80)
    BLUE_NIGHT = RGBColor(0x0A, 0x15, 0x30)
    GREY_SOFT  = RGBColor(0xF3, 0xF4, 0xF6)
    BLACK      = RGBColor(0x00, 0x00, 0x00)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# -----------------------------------------------------------------------------
# SVG -> PNG rasterization (for diagrams)
# -----------------------------------------------------------------------------

def rasterize_svg(svg_path: Path, target_w_px: int) -> Path:
    if not svg_path.exists():
        raise FileNotFoundError(svg_path)
    out = RASTER_DIR / (svg_path.stem + f"_{target_w_px}.png")
    if out.exists() and out.stat().st_mtime >= svg_path.stat().st_mtime:
        return out
    subprocess.run(
        ["rsvg-convert", "--width", str(target_w_px), "--keep-aspect-ratio",
         str(svg_path), "-o", str(out)],
        check=True, capture_output=True,
    )
    return out


def rasterize_svg_recolored(svg_path: Path, target_w_px: int,
                            color_hex: str,
                            stroke_width: float | None = None) -> Path:
    """Rasterize a `currentColor`-based SVG (Tabler icon family) with an
    explicit stroke/fill colour. Mirrors the exec build script's helper of
    the same name — duplicated here to keep the architect script self-
    contained (no import-time dependency on exec build internals beyond
    slide_6_native, which is intentional reuse). Falls back to on-the-fly
    rasterisation under `_raster/`; does NOT honour exec's prebuilt-icon
    shortcut, because the architect deck doesn't ship its own prebuilt
    library yet.
    """
    if not svg_path.exists():
        raise FileNotFoundError(svg_path)
    colour = color_hex.lstrip("#").upper()
    sw_tag = (f"_sw{stroke_width:g}".replace(".", "p")
              if stroke_width is not None else "")
    out = RASTER_DIR / f"{svg_path.stem}_{target_w_px}_{colour}{sw_tag}.png"
    if out.exists() and out.stat().st_mtime >= svg_path.stat().st_mtime:
        return out
    src = svg_path.read_text(encoding="utf-8")
    patched = src.replace("currentColor", f"#{colour}")
    if stroke_width is not None:
        import re as _re
        patched = _re.sub(
            r'stroke-width="[^"]*"',
            f'stroke-width="{stroke_width:g}"',
            patched,
        )
    tmp = RASTER_DIR / f"{svg_path.stem}_{colour}{sw_tag}.svg"
    tmp.write_text(patched, encoding="utf-8")
    subprocess.run(
        ["rsvg-convert", "--width", str(target_w_px), "--keep-aspect-ratio",
         str(tmp), "-o", str(out)],
        check=True, capture_output=True,
    )
    return out


def find_diagram(name: str) -> Path | None:
    for candidate in (DIAGRAMS_DIR / name,
                      DIAGRAMS_DIR / "architect" / name,
                      EXEC_DIAGRAMS_DIR / name):
        if candidate.exists():
            return candidate
    return None


# -----------------------------------------------------------------------------
# Open .potx as .pptx
# -----------------------------------------------------------------------------

def open_template_as_pptx() -> Presentation:
    if not POTX_PATH.exists():
        raise FileNotFoundError(
            f"{POTX_PATH} not found — run `python build_potx.py` first."
        )
    tmp_pptx = RASTER_DIR / "_potx_loaded.pptx"
    with zipfile.ZipFile(POTX_PATH, "r") as src:
        data = {n: src.read(n) for n in src.namelist()}
    ct = data["[Content_Types].xml"].decode("utf-8")
    ct = ct.replace(
        "presentationml.template.main+xml",
        "presentationml.presentation.main+xml",
    )
    data["[Content_Types].xml"] = ct.encode("utf-8")
    with zipfile.ZipFile(tmp_pptx, "w", zipfile.ZIP_DEFLATED) as out:
        for name, blob in data.items():
            out.writestr(name, blob)
    return Presentation(str(tmp_pptx))


# -----------------------------------------------------------------------------
# Placeholder helpers
# -----------------------------------------------------------------------------

def find_placeholder(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(f"placeholder idx={idx} not found on slide using "
                   f"layout {slide.slide_layout.name!r}")


def delete_placeholder(slide, idx: int):
    ph = find_placeholder(slide, idx)
    sp = ph._element
    sp.getparent().remove(sp)


def set_text(slide, idx: int, text: str, *, color: RGBColor | None = None,
             align_left: bool = False):
    ph = find_placeholder(slide, idx)
    tf = ph.text_frame
    tf.clear()
    paragraphs = text.split("\n")
    for i, line in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if align_left:
            p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        if color is not None:
            run.font.color.rgb = color


def set_split_gradient_title(slide, idx: int, prefix: str, noun: str,
                              align_left: bool = True):
    ph = find_placeholder(slide, idx)
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align_left:
        p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = prefix
    r1.font.color.rgb = C.WHITE
    r2 = p.add_run()
    r2.text = noun
    rPr = r2._r.get_or_add_rPr()
    for tag in ("solidFill", "gradFill", "noFill"):
        existing = rPr.find(f"{{{A_NS}}}{tag}")
        if existing is not None:
            rPr.remove(existing)
    grad_xml = (
        f'<a:gradFill xmlns:a="{A_NS}" flip="none" rotWithShape="1">'
        '<a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="FFFFFF"/></a:gs>'
        '<a:gs pos="35000"><a:srgbClr val="5CD6FF"/></a:gs>'
        '<a:gs pos="75000"><a:srgbClr val="0F6BFF"/></a:gs>'
        '</a:gsLst>'
        '<a:lin ang="0" scaled="1"/>'
        '</a:gradFill>'
    )
    rPr.insert(0, etree.fromstring(grad_xml))


def set_blue_box_bullets(slide, idx: int, items: list[str],
                          *, font_size: int | None = None):
    """Render a body placeholder as a list with blue square bullets.

    Each item may optionally start with **anchor** ... — the anchor (the part
    before the first ' — ') is rendered bold-blue; the rest in black.

    font_size (pt) overrides the placeholder's default size. Pass None to
    keep the master's size.
    """
    ph = find_placeholder(slide, idx)
    tf = ph.text_frame
    tf.clear()
    for i, body in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        bullet = p.add_run()
        bullet.text = "▪  "
        bullet.font.color.rgb = C.BLUE
        bullet.font.bold = True
        if font_size is not None:
            bullet.font.size = Pt(font_size)
        # Split on first em-dash if present, to bold the anchor.
        sep = " — "
        if sep in body:
            anchor, rest = body.split(sep, 1)
            r_anchor = p.add_run()
            r_anchor.text = anchor
            r_anchor.font.color.rgb = C.BLUE
            r_anchor.font.bold = True
            if font_size is not None:
                r_anchor.font.size = Pt(font_size)
            r_rest = p.add_run()
            r_rest.text = sep + rest
            r_rest.font.color.rgb = C.BLACK
            if font_size is not None:
                r_rest.font.size = Pt(font_size)
        else:
            text_run = p.add_run()
            text_run.text = body
            text_run.font.color.rgb = C.BLACK
            if font_size is not None:
                text_run.font.size = Pt(font_size)


def set_action_path_lines(slide, idx: int, lines: list[tuple[str, str]],
                            sep: str = " — "):
    """CTA-style: ACTION (cyan, bold) + sep + path (white)."""
    ph = find_placeholder(slide, idx)
    tf = ph.text_frame
    tf.clear()
    for i, (action, path) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run()
        r1.text = action
        r1.font.color.rgb = C.CYAN
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = sep + path
        r2.font.color.rgb = C.WHITE


# -----------------------------------------------------------------------------
# Static decoration helpers
# -----------------------------------------------------------------------------

def add_banner_full_bleed(slide, image_path: Path):
    if not image_path.exists():
        return
    pic = slide.shapes.add_picture(str(image_path), 0, 0,
                                    width=px(SLIDE_W_PX),
                                    height=px(SLIDE_H_PX))
    spTree = pic._element.getparent()
    spTree.remove(pic._element)
    insert_at = 0
    for i, el in enumerate(spTree):
        if el.tag.endswith("}grpSpPr"):
            insert_at = i + 1
            break
    spTree.insert(insert_at, pic._element)


def add_brand_row(slide):
    ocm_svg = ASSETS_DIR / "ocm" / "ocm-horizontal-white.svg"
    nn_svg = ASSETS_DIR / "neonephos" / "neonephos-foundation-horizontal-white.svg"
    ocm_png = rasterize_svg(ocm_svg, target_w_px=400)
    nn_png = rasterize_svg(nn_svg, target_w_px=400)
    ocm_pic = slide.shapes.add_picture(
        str(ocm_png), px(96), px(SLIDE_H_PX - 56 - 76), height=px(76)
    )
    nn_pic = slide.shapes.add_picture(
        str(nn_png), px(96), px(SLIDE_H_PX - 56 - 52), height=px(52)
    )
    nn_pic.left = px(SLIDE_W_PX - 96) - nn_pic.width


def add_diagram(slide, svg_path: Path | None,
                 x_px: int, y_px: int,
                 max_w_px: int, max_h_px: int):
    if svg_path is None or not svg_path.exists():
        return
    try:
        delete_placeholder(slide, 10)
    except KeyError:
        pass
    png = rasterize_svg(svg_path, target_w_px=max_w_px)
    pic = slide.shapes.add_picture(str(png), px(x_px), px(y_px),
                                    width=px(max_w_px))
    if pic.height > px(max_h_px):
        ratio = px(max_h_px) / pic.height
        pic.height = px(max_h_px)
        pic.width = int(pic.width * ratio)
    pic.left = px(x_px) + (px(max_w_px) - pic.width) // 2
    pic.top  = px(y_px) + (px(max_h_px) - pic.height) // 2


def set_speaker_notes(slide, text: str):
    """Embed presenter prompts into the slide's notes pane.

    python-pptx exposes notes via slide.notes_slide.notes_text_frame. We
    clear whatever the master left there and write one paragraph per
    newline-separated line. The result lives in the speaker view in
    PowerPoint and in printed handouts. Canonical full notes (with timers,
    Q&A prep, stage directions) stay in ../notes/SPEAKER-NOTES-ARCHITECT-
    EXTERNAL.md; what's embedded here is the trimmed beat-list."""
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        # Notes pane uses the master's default font; no explicit styling.


def add_footer_caption(slide, y_px: int, text: str, *, italic: bool = False):
    tb = slide.shapes.add_textbox(px(120), px(y_px),
                                   px(SLIDE_W_PX - 240), px(60))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(18)
    r.font.italic = italic
    r.font.color.rgb = C.GREY_MID


def add_textbox(slide, x_px: int, y_px: int, w_px: int, h_px: int):
    """Helper: returns a textbox with margins zeroed and word-wrap on."""
    tb = slide.shapes.add_textbox(px(x_px), px(y_px), px(w_px), px(h_px))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    return tb, tf


def add_yaml_block(slide, x_px: int, y_px: int, w_px: int, h_px: int,
                    yaml_lines: list[tuple[str, RGBColor]],
                    *, font_size: int = 16,
                    bg_color: RGBColor = C.GREY_SOFT,
                    border: bool = True):
    """Render a YAML-style block. Each line is a (text, color) pair.
    Background is a rounded rectangle in GREY_SOFT by default."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  px(x_px), px(y_px), px(w_px), px(h_px))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    if border:
        box.line.color.rgb = C.GREY_MID
        box.line.width = Pt(0.75)
    else:
        box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(140000)
    tf.margin_top = tf.margin_bottom = Emu(120000)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (line, color) in enumerate(yaml_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(font_size)
        r.font.color.rgb = color


def add_callout(slide, x_px: int, y_px: int, w_px: int, h_px: int,
                 anchor: str, body: str,
                 *, anchor_color: RGBColor = C.BLUE):
    """Small explanatory callout box used next to YAML blocks.
    Anchor word bold-blue, body in black. White background, blue left bar."""
    # left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  px(x_px), px(y_px), px(6), px(h_px))
    bar.fill.solid()
    bar.fill.fore_color.rgb = anchor_color
    bar.line.fill.background()
    # text
    tb, tf = add_textbox(slide, x_px + 18, y_px, w_px - 18, h_px)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r_a = p.add_run()
    r_a.text = anchor
    r_a.font.name = "Aptos"
    r_a.font.size = Pt(16)
    r_a.font.bold = True
    r_a.font.color.rgb = anchor_color
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(4)
    r_b = p2.add_run()
    r_b.text = body
    r_b.font.name = "Aptos"
    r_b.font.size = Pt(15)
    r_b.font.color.rgb = C.GREY_MID


def add_left_bullets(slide, x_px: int, y_px: int, w_px: int, h_px: int,
                      items: list[str], *,
                      font_size: int = 20, line_spacing: float = 1.35):
    """Render a custom textbox with blue-bullet + bold-anchor + body bullets,
    used when the layout body placeholder is too narrow / wide for the slide.
    Each item: '**Anchor** — body.' uses ' — ' as separator."""
    tb, tf = add_textbox(slide, x_px, y_px, w_px, h_px)
    sep = " — "
    for i, body in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_before = Pt(0) if i == 0 else Pt(10)
        p.space_after = Pt(0)
        bullet = p.add_run()
        bullet.text = "▪  "
        bullet.font.name = "Aptos"
        bullet.font.size = Pt(font_size)
        bullet.font.color.rgb = C.BLUE
        bullet.font.bold = True
        if sep in body:
            anchor, rest = body.split(sep, 1)
            r_a = p.add_run()
            r_a.text = anchor
            r_a.font.name = "Aptos"
            r_a.font.size = Pt(font_size)
            r_a.font.color.rgb = C.BLUE
            r_a.font.bold = True
            r_b = p.add_run()
            r_b.text = sep + rest
            r_b.font.name = "Aptos"
            r_b.font.size = Pt(font_size)
            r_b.font.color.rgb = C.BLACK
        else:
            r = p.add_run()
            r.text = body
            r.font.name = "Aptos"
            r.font.size = Pt(font_size)
            r.font.color.rgb = C.BLACK


# -----------------------------------------------------------------------------
# Layout lookup
# -----------------------------------------------------------------------------

def layouts_by_name(prs: Presentation) -> dict[str, object]:
    return {l.name: l for l in prs.slide_masters[0].slide_layouts}


def customize_hero_for_architect(prs):
    layouts = layouts_by_name(prs)
    hero = layouts.get("Hero")
    if hero is None:
        return
    size_changed = False
    for ph in hero.placeholders:
        try:
            idx = ph.placeholder_format.idx
        except Exception:
            continue
        if idx not in (1, 2):
            continue
        try:
            ph.width = px(1824)
        except Exception:
            pass
        sp = ph._element
        for def_rpr in sp.iter(f"{{{A_NS}}}defRPr"):
            sz = def_rpr.get("sz")
            if sz == "13200":
                def_rpr.set("sz", "11500")
                size_changed = True
    if not size_changed:
        print("NOTE: Hero font size unchanged (sz=13200 not found on "
              "layout placeholders); width widen applied.")


def shrink_plain_titles(prs, target_pt: int = 56):
    """Drop the Plain & Plain/Compact title size so longer titles don't wrap
    into the body slot. target_pt expressed in standard pt (e.g. 56 -> 5600).
    Hazard mitigation per handoff §5 — content-first deck, geometry second."""
    layouts = layouts_by_name(prs)
    target_sz = str(target_pt * 100)
    for layout_name in ("Plain", "Plain / Compact"):
        layout = layouts.get(layout_name)
        if layout is None:
            continue
        for ph in layout.placeholders:
            try:
                idx = ph.placeholder_format.idx
            except Exception:
                continue
            if idx != 2:  # title
                continue
            sp = ph._element
            for def_rpr in sp.iter(f"{{{A_NS}}}defRPr"):
                # The Plain title is 74pt = 7400 in lstStyle.
                if def_rpr.get("sz") == "7400":
                    def_rpr.set("sz", target_sz)


def shrink_diagram_title(prs, target_pt: int = 44):
    """Content/Diagram title slot is only h=80px but defaults to 74pt — any
    multi-line title clips off the bottom. Slides 5/6/7/12 share this layout
    in the architect deck, so we shrink the title to fit a single line of
    most reasonable titles, or two short lines."""
    layouts = layouts_by_name(prs)
    layout = layouts.get("Content / Diagram")
    if layout is None:
        return
    target_sz = str(target_pt * 100)
    for ph in layout.placeholders:
        try:
            idx = ph.placeholder_format.idx
        except Exception:
            continue
        if idx != 2:
            continue
        sp = ph._element
        for def_rpr in sp.iter(f"{{{A_NS}}}defRPr"):
            if def_rpr.get("sz") == "7400":
                def_rpr.set("sz", target_sz)


# =============================================================================
# Build the deck
# =============================================================================

def build():
    prs = open_template_as_pptx()
    customize_hero_for_architect(prs)
    layouts = layouts_by_name(prs)

    expected = {"Hero", "CTA", "Content / 3-Column",
                "Content / 3-Column Tall Title",
                "Content / Diagram", "Content / Diagram Compact",
                "Content / Tiles", "Content / 2-Column",
                "Section Divider", "Plain", "Plain / Compact"}
    missing = expected - set(layouts)
    if missing:
        sys.exit(f"template missing expected layouts: {missing}")

    build_slide_1_pain(prs, layouts)
    build_slide_2_cause(prs, layouts)
    build_slide_3_insight(prs, layouts)
    build_slide_4_positioning(prs, layouts)
    build_slide_5_constructor(prs, layouts)
    build_slide_6_descriptor(prs, layouts)
    build_slide_7_overview(prs, layouts)
    build_slide_8_pack(prs, layouts)
    build_slide_9_sign(prs, layouts)
    build_slide_10_transport(prs, layouts)
    build_slide_11_deploy(prs, layouts)
    build_slide_12_composition(prs, layouts)
    build_slide_14_adoption(prs, layouts)        # was 14, now slide 13
    build_slide_13_whats_sharp(prs, layouts)     # was 13, now slide 14
    build_slide_15_cta(prs, layouts)
    build_slide_16_appendix_replication(prs, layouts)  # appendix, out of main arc

    # Embed condensed presenter prompts into the notes pane of each slide.
    for idx, slide in enumerate(prs.slides, start=1):
        note = SPEAKER_NOTES.get(idx)
        if note:
            set_speaker_notes(slide, note)

    sanity_check(prs)
    prs.save(str(OUTPUT_PPTX))
    print(f"Wrote {OUTPUT_PPTX} ({len(prs.slides)} slides)")


# -----------------------------------------------------------------------------
# Slide builders
# -----------------------------------------------------------------------------

def build_slide_1_pain(prs, layouts):
    """1 PAIN — Hero. Two-line gradient title. Subtitle. Footer. Brand row."""
    s = prs.slides.add_slide(layouts["Hero"])
    add_banner_full_bleed(s, THEME_DIR / "OCM-Banner.png")
    set_text(s, 1, "You ship pieces.", color=C.WHITE, align_left=True)
    set_split_gradient_title(s, 2, prefix="",
                              noun="Nothing carries the release.")
    set_text(s, 3,
             "You sign the pieces. Nothing signs the release.",
             color=C.CYAN)
    set_text(s, 4,
             "Open Component Model — open source, NeoNephos Foundation.",
             color=C.WHITE)
    add_brand_row(s)


def build_slide_2_cause(prs, layouts):
    """2 CAUSE — Diagnosis. Three bullets, footer caption."""
    s = prs.slides.add_slide(layouts["Plain / Compact"])
    set_text(s, 1, "DIAGNOSIS")
    set_text(s, 2, "In every existing tool, identity is bound to location.")
    set_blue_box_bullets(s, 10, [
        "OCI image — identified by registry/repo:tag. Mirror it; every downstream reference is invalidated.",
        "Helm chart — identified by repo URL + name + version. Mirror the repo; pulls fail.",
        "SBOM — linked to its subject by file path or naming convention. Move the artifact; the link dangles.",
    ])
    # Punchline moved to speaker notes: "Cosign attestations sign each
    # piece. None of them sign the release as one named, location-
    # independent unit."


def _draw_coordinate_travel_diagram(slide, *, x_px: int, y_px: int,
                                     w_px: int, h_px: int):
    """Native-PPT visual for slide 3: one coordinate at top, three registry
    cylinders below, per-cylinder access label.

    Typography (after row-14 readability review):
      - Chip:    20pt (was 17)
      - Silos:   20pt (was 18)
      - Access:  18pt (was 13)
    Bottom 'same digest' annotation removed — redundant with the left-side
    Digest bullet on the slide.
    """
    # ---- Coordinate chip at top -----------------------------------------
    chip_w = 600
    chip_h = 64
    chip_x = x_px + (w_px - chip_w) // 2
    chip_y = y_px
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   px(chip_x), px(chip_y),
                                   px(chip_w), px(chip_h))
    chip.fill.solid()
    chip.fill.fore_color.rgb = C.GREY_SOFT
    chip.line.color.rgb = C.BLUE
    chip.line.width = Pt(1.5)
    tf = chip.text_frame
    tf.margin_left = tf.margin_right = Emu(80000)
    tf.margin_top = tf.margin_bottom = Emu(40000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "github.com/acme.org/helloworld:1.0.0"
    r.font.name = "Consolas"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = C.BLUE

    # ---- Three registry cylinders ---------------------------------------
    cyl_w = 160
    cyl_h = 140
    gap = (w_px - 3 * cyl_w) // 4
    cyl_y = y_px + chip_h + 80
    cyls = [
        ("EU reg",    "ghcr.io/eu/..."),
        ("US reg",    "ghcr.io/us/..."),
        ("Air-gap",   "registry.local/..."),
    ]
    chip_bottom_x = chip_x + chip_w // 2
    chip_bottom_y = chip_y + chip_h
    cyl_centers_x = []
    for i, (name, access_ref) in enumerate(cyls):
        cx = x_px + gap + i * (cyl_w + gap)
        cyl_centers_x.append(cx + cyl_w // 2)
        cyl = slide.shapes.add_shape(MSO_SHAPE.CAN,
                                      px(cx), px(cyl_y),
                                      px(cyl_w), px(cyl_h))
        cyl.fill.solid()
        cyl.fill.fore_color.rgb = C.WHITE
        cyl.line.color.rgb = C.BLUE
        cyl.line.width = Pt(1.5)
        ctf = cyl.text_frame
        ctf.margin_left = ctf.margin_right = Emu(40000)
        ctf.margin_top = ctf.margin_bottom = Emu(40000)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = name
        cr.font.name = "Aptos"
        cr.font.size = Pt(20)
        cr.font.bold = True
        cr.font.color.rgb = C.BLUE
        # Access reference under each cylinder
        atb, atf = add_textbox(slide, cx - 40, cyl_y + cyl_h + 16,
                                cyl_w + 80, 36)
        ap = atf.paragraphs[0]
        ap.alignment = PP_ALIGN.CENTER
        ar = ap.add_run()
        ar.text = f"access: {access_ref}"
        ar.font.name = "Consolas"
        ar.font.size = Pt(18)
        ar.font.color.rgb = C.GREY_MID

    # ---- Connector lines from chip bottom to each cylinder top ----------
    for cyl_cx in cyl_centers_x:
        line_shape = slide.shapes.add_connector(1,
            px(chip_bottom_x), px(chip_bottom_y),
            px(cyl_cx), px(cyl_y))
        line_shape.line.color.rgb = C.GREY_MID
        line_shape.line.width = Pt(1.25)


def build_slide_3_insight(prs, layouts):
    """3 INSIGHT — The Hinge. Left bullets, right ASCII diagram, footer."""
    s = prs.slides.add_slide(layouts["Plain"])
    set_text(s, 1, "THE HINGE")
    set_text(s, 2, "Identity that travels with the artifact.")
    delete_placeholder(s, 10)

    add_left_bullets(
        s, x_px=120, y_px=540, w_px=940, h_px=420,
        items=[
            "Component identity — name and version of the component. "
            "Globally unique. Location-agnostic.",
            "Digest — every resource inside the component carries a "
            "content hash. Computed once.",
            "Access — where the resource currently lives. "
            "Rewritten on transfer. Digest stays.",
        ],
        font_size=28,
    )

    # RIGHT half: native PowerPoint visual showing coordinate-travel.
    _draw_coordinate_travel_diagram(s,
                                     x_px=1060, y_px=540,
                                     w_px=820, h_px=440)

    # Footer caption — the slide's load-bearing one-liner. Brand blue so
    # it reads as the conclusion to the diagram above. Same treatment as
    # slide 12 footer (consistent "this is the punchline" signal across
    # the deck).
    footer_tb, footer_tf = add_textbox(s, 120, 990, SLIDE_W_PX - 240, 60)
    fp = footer_tf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = "Move the artifact. The digest stays. Only the access changes."
    fr.font.name = "Aptos"
    fr.font.size = Pt(24)
    fr.font.color.rgb = C.BLUE


def build_slide_4_positioning(prs, layouts):
    """4 POSITIONING — Where OCM Sits. Three columns.

    Columns named to define the 'component' noun before slides 5+ rely on it:
    a component is the (artifact-format-agnostic, location-agnostic, signed)
    wrapper around resources."""
    s = prs.slides.add_slide(layouts["Content / 3-Column Tall Title"])
    set_text(s, 1, "WHERE OCM SITS")
    set_text(s, 2, "Wraps every artifact. Signs the whole release.")
    set_text(s, 10, "ANY FORMAT")
    set_text(s, 11, "OCI, Helm, configs, SBOMs, npm, maven, binaries.\n"
                     "Artifact type is free-form; access types are pluggable.")
    set_text(s, 12, "ANY LOCATION")
    set_text(s, 13, "Component identity travels.\n"
                     "The component carries its name across registries.")
    set_text(s, 14, "ONE SIGNATURE")
    set_text(s, 15, "Covers every digest in the component.\n"
                     "Survives transport.")
    # Punchline moved to speaker notes: "A component is the unit you
    # sign, transport, and deploy. The next slides show how it's built."


def build_slide_5_constructor(prs, layouts):
    """5 CONSTRUCTOR — What you write.

    Pt16 YAML, ~17 lines showing the two pack-time mechanisms:
    `input:` (by value) and `access:` (by reference). componentReferences
    composition was deliberately dropped from this slide on 2026-06-24 —
    it lives on slide 8 (COMPOSE), where it gets its own visual treatment.
    Earlier iteration carried a componentReferences stanza here, but the
    slide tries to teach two mechanisms and showing a third concept
    diluted the input-vs-access distinction.
    """
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "CONSTRUCTOR")
    set_text(s, 2, "What you write.")
    delete_placeholder(s, 10)

    K = C.BLUE_MID
    V = C.BLACK
    COM = C.GREY_MID
    yaml_lines = [
        ("components:",                                                       K),
        ("- name: github.com/acme.org/helloworld",                            V),
        ("  version: 1.0.0",                                                  V),
        ("  provider:",                                                       K),
        ("    name: acme.org",                                                V),
        ("  resources:",                                                      K),
        ("    - name: mylocalfile",                                           V),
        ("      type: blob",                                                  V),
        ("      input:                # Embed by value",                      COM),
        ("        type: File/v1",                                             V),
        ("        path: ./my-local-resource.txt",                             V),
        ("    - name: image",                                                 V),
        ("      type: ociImage",                                              V),
        ("      version: 1.0.0",                                              V),
        ("      access:               # Reference external artifact",         COM),
        ("        type: OCIImage/v1",                                         V),
        ("        imageReference: ghcr.io/stefanprodan/podinfo:6.9.1",        V),
    ]
    # 17 lines @ Pt16 @ ls 1.20 ~ 408pt ~ 544px + ~60px padding ~ 604px.
    # Box y=300, h=620 fits with comfortable bottom padding now that
    # componentReferences is gone.
    add_yaml_block(s, x_px=120, y_px=300, w_px=1180, h_px=620,
                    yaml_lines=yaml_lines, font_size=16)

    # Callouts aligned to their YAML sections.
    # Line height @ Pt16 ls 1.20 ~ 25.6px. Box top y=300 + ~30px top padding.
    # Line 9 (input:) starts at ~y=508; line 15 (access:) at ~y=662.
    # Right-side callouts removed — the inline YAML comments
    # (`# Embed by value`, `# Reference external artifact`) carry the
    # input-vs-access distinction without a horizontal saccade across
    # the slide. Architects read comments inline as they parse YAML
    # structure left-to-right; callouts were redundant friction.


def build_slide_6_descriptor(prs, layouts):
    """6 DESCRIPTOR — What gets signed and travels.

    Technically-correct signing semantics (per signing-and-verification-concept.md):
    - Each resource carries a content digest. Those digests are inputs to the
      canonical descriptor normalization.
    - OCM signs ONE hash: the SHA-256 of the canonicalized descriptor.
    - That single hash is what the signature value covers.
    - access fields are excluded from normalization, so signature survives transfer.

    Callouts ordered to match the descriptor flow: access first (it appears first
    in the YAML, and it is the location pointer), then digest (the content
    identity), then signature (one hash over the whole canonicalized descriptor).
    """
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "DESCRIPTOR")
    set_text(s, 2, "What gets signed and travels.")
    delete_placeholder(s, 10)

    K = C.BLUE_MID
    V = C.BLACK
    COM = C.GREY_MID
    SIG = C.BLUE
    # Color discipline (audited, semantics-based):
    #   K (dark blue)   — every YAML key, INCLUDING `signatures:` itself.
    #                     The key is structural — the architect's eye-
    #                     anchor for "this is YAML scaffolding".
    #   V (black)       — values that are metadata about the component
    #                     (name, version, type, etc.). NOT the cryptographic
    #                     payload.
    #   COM (grey)      — inline comments. Kept on the same line as the key
    #                     they annotate; padded with spaces for readability.
    #   SIG (brand blue)— RESERVED for "content the signature covers /
    #                     content the signature IS": digest values + the
    #                     signature value + the contents of the
    #                     signatures: block (name, algorithm, etc.). The
    #                     rule is "wherever brand blue appears, this is
    #                     the cryptographic payload". Architects can scan
    #                     for blue and see exactly what's signed.
    yaml_lines = [
        ("component:                              # (fields trimmed)",                       COM),
        ("  name: github.com/acme.org/helloworld",                                           V),
        ("  version: 1.0.0",                                                                 V),
        ("  resources:",                                                                     K),
        ("    - name: image",                                                                V),
        ("      type: ociImage",                                                             V),
        ("      access:                          # excluded — rewritten on transfer",        COM),
        ("        type: OCIImage/v1",                                                        V),
        ("        imageReference: ghcr.io/.../podinfo@sha256:8fa569...",                     V),
        ("",                                                                                 V),
        ("      digest:                          # content identity — input to descriptor hash", COM),
        ("        hashAlgorithm: SHA-256",                                                   SIG),
        ("        value: 262578cde928d5c9eba3bce0...",                                       SIG),
        ("",                                                                                 V),
        ("signatures:                            # one hash over the canonicalized descriptor", COM),
        ("  - name: acme-release-key",                                                       SIG),
        ("    digest:                            # of the descriptor",                       COM),
        ("      hashAlgorithm: SHA-256",                                                     SIG),
        ("      value: a4b1c2d3e4f5...",                                                     SIG),
        ("    signature:",                                                                   K),
        ("      algorithm: RSASSA-PSS",                                                      SIG),
        ("      value: <hex-encoded signature>",                                             SIG),
    ]
    # Note on the inline-comment-on-same-line trick: a single tuple gets
    # one color from add_yaml_block. So when a line has both a key and a
    # trailing # comment, we color the WHOLE line with the tone that best
    # serves the line's job. For `component:` and the digest/signatures
    # keys with inline comments, COM (grey) reads as the line's tone —
    # the architect sees the comment first, the key second. That's the
    # correct reading order for an annotated key: "this thing, with this
    # caveat". For uncommented keys (`resources:`, `signature:`), K (dark
    # blue) renders as pure structure.
    #
    # 22 lines @ Pt14 @ ls 1.20 ~ 403pt ~ 540px-ish at Pt14. Box height
    # 680 with padding fits comfortably.
    add_yaml_block(s, x_px=120, y_px=280, w_px=1500, h_px=680,
                    yaml_lines=yaml_lines, font_size=14)

    # Callouts aligned with YAML sections.
    # Line height @ Pt16 ls 1.20 ~ 25.6px. Box top y=280 + ~30px padding.
    # Line 7 (access:) at ~y=464; line 10 (digest:) at ~y=541;
    # line 14 (signatures:) at ~y=643.
    # Right-side callouts removed for slide 6 — inline YAML comments now
    # carry the access-excluded / digest-as-input / signature-as-hash
    # meaning, and brand-blue highlighting (signatures block + digest
    # values) marks what's new vs the constructor on slide 5.

    # No footer caption — at 18pt it's unreadable at presentation distance.
    # Speaker delivers "Sign the descriptor hash, not the access." verbally.


def build_slide_7_overview(prs, layouts):
    """7 OVERVIEW — Pack · Sign · Transport · Deploy.

    Native PowerPoint shapes via the exec deck's slide_6_native module.
    The user crafted this slide by hand in PowerPoint after copying the
    exec slide-8 layout; the PPT XML extraction confirmed the geometry
    matches `add_pack_sign_transport_deploy_native(s, x=60, y=240,
    w=1800, h=780, ...)` 1:1 (cards at 337.5×375, top stripe h≈4,
    icons at +30/+34 from card top-left, etc.).

    Eyebrow is THE FOUR MOVES (the architect-deck framing — docs frame
    OCM as a lifecycle, not a CLI verb list). The CLI verbs you'll
    actually type are `ocm add cv` / `ocm sign cv` / `ocm transfer cv`
    / kubectl apply against the Deployer — speaker notes bridge to that.

    Card body text matches the exec deck's selling copy (the user kept
    the exec wording when porting the slide; the architect-deck wording
    earlier proposed was a deck-author refinement that didn't land).
    """
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "THE FOUR MOVES")
    set_text(s, 2, "Pack · Sign · Transport · Deploy.")
    icons_dir = EXEC_DIAGRAMS_DIR / "icons"
    add_pack_sign_transport_deploy_native(
        s,
        x=60, y=240, w=1800, h=780,
        icons_dir=icons_dir,
        rasterize_recolored=rasterize_svg_recolored,
        # No cards_override — defaults to the exec CARDS list, which
        # matches the extracted slide-7 text.
    )


def build_slide_8_pack(prs, layouts):
    """8 COMPOSE — Components compose. Service vs product/release.

    Repurposed from Pack to introduce composition before slide 12 needs it.
    Real architecture: service components carry resources; product/release
    components have no resources of their own — they reference services
    by name and version. Two YAML snippets side by side show the shape.

    Terminology aligned to docs (`how-to/model-products.md`): SERVICE
    component (was "LEAF") carries resources; PRODUCT component composes
    services by `componentReferences:`. Closing phrase "One release unit,
    transferable, signable end-to-end" anchors the release noun the deck
    title rests on.
    """
    s = prs.slides.add_slide(layouts["Content / Diagram Compact"])
    set_text(s, 1, "COMPOSE")
    set_text(s, 2, "Service carries resources. Product carries references.")
    delete_placeholder(s, 10)

    # Intro text aligned with title (x=120).
    # Trailing "…" rhymes with the `...` continuation lines in the service
    # YAML below — same idiom across two registers: "more types possible"
    # in the prose, "more fields per resource" in the YAML. "by name and
    # version" dropped — the right-box YAML already shows the mechanism.
    intro_tb, intro_tf = add_textbox(s, 120, 280, 1680, 90)
    intro_tf.word_wrap = True
    ip = intro_tf.paragraphs[0]
    ip.alignment = PP_ALIGN.LEFT
    ip.line_spacing = 1.30
    for text, kind in [
        ("Service components ", "accent"),
        ("carry resources — images, charts, configs, SBOMs, …\n", "normal"),
        ("A product component ", "accent"),
        ("composes other components.", "normal"),
        (" One release unit, transferable, signable end-to-end.", "normal"),
    ]:
        r = ip.add_run()
        r.text = text
        r.font.name = "Aptos"
        r.font.size = Pt(22)
        if kind == "accent":
            r.font.bold = True
            r.font.color.rgb = C.BLUE
        else:
            r.font.color.rgb = C.BLACK

    K = C.BLUE_MID
    V = C.BLACK
    COM = C.GREY_MID

    # YAML notes:
    #   - Top header comments dropped on both boxes — the title and the
    #     SERVICE/PRODUCT labels already name what each box is; the headers
    #     were saying the same thing a third time.
    #   - First trimmed resource carries the full comment ("# type,
    #     input/access, digest trimmed"); the other three trimmed resources
    #     just use a bare `...` on a separate indented line. Teaches the
    #     convention once, then uses shorthand. Cuts visual weight without
    #     losing the "more fields per resource" signal.
    leaf_yaml = [
        ("components:",                                      K),
        ("  - name: acme.org/sovereign/notes",               V),
        ("    version: 1.0.0",                               V),
        ("    resources:",                                   K),
        ("      - name: image       # OCI image",            V),
        ("        # type, input/access, digest trimmed",    COM),
        ("      - name: chart       # Helm chart",           V),
        ("        ...",                                      COM),
        ("  - name: acme.org/sovereign/postgres",            V),
        ("    version: 1.0.0",                               V),
        ("    resources:",                                   K),
        ("      - name: image       # OCI image",            V),
        ("        ...",                                      COM),
        ("      - name: chart       # Helm chart",           V),
        ("        ...",                                      COM),
        ("# product references both, by name and version",   COM),
    ]
    product_yaml = [
        ("components:",                                        K),
        ("  - name: acme.org/sovereign/product",                V),
        ("    version: 1.0.0",                                  V),
        ("    componentReferences:",                            K),
        ("      - name: notes",                                 V),
        ("        componentName: acme.org/sovereign/notes",     V),
        ("        version: 1.0.0",                              V),
        ("      - name: postgres",                              V),
        ("        componentName: acme.org/sovereign/postgres",  V),
        ("        version: 1.0.0",                              V),
        ("# no resources of its own — pure composition",      COM),
    ]
    # Two side-by-side boxes, leaner widths, arrow between.
    box_w = 760
    box_h = 560
    box_y = 460
    arrow_w = 100
    arrow_h = 80
    gap = 30
    total_w = 2 * box_w + 2 * gap + arrow_w
    left_x = (SLIDE_W_PX - total_w) // 2
    right_x = left_x + box_w + gap + arrow_w + gap
    arrow_x = left_x + box_w + gap
    arrow_y = box_y + (box_h - arrow_h) // 2

    add_yaml_block(s, x_px=left_x, y_px=box_y, w_px=box_w, h_px=box_h,
                    yaml_lines=leaf_yaml, font_size=17)
    add_yaml_block(s, x_px=right_x, y_px=box_y, w_px=box_w, h_px=box_h,
                    yaml_lines=product_yaml, font_size=17)

    # Labels above each block.
    ll_tb, ll_tf = add_textbox(s, left_x, box_y - 50, box_w, 36)
    lp = ll_tf.paragraphs[0]; lp.alignment = PP_ALIGN.LEFT
    lr = lp.add_run(); lr.text = "SERVICE"
    lr.font.name = "Aptos"; lr.font.size = Pt(22); lr.font.bold = True
    lr.font.color.rgb = C.GREY_MID

    rl_tb, rl_tf = add_textbox(s, right_x, box_y - 50, box_w, 36)
    rp = rl_tf.paragraphs[0]; rp.alignment = PP_ALIGN.LEFT
    rr = rp.add_run(); rr.text = "PRODUCT"
    rr.font.name = "Aptos"; rr.font.size = Pt(22); rr.font.bold = True
    rr.font.color.rgb = C.BLUE

    # Composition arrow service → product.
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                px(arrow_x), px(arrow_y),
                                px(arrow_w), px(arrow_h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = C.BLUE
    arrow.line.fill.background()


def build_slide_9_sign(prs, layouts):
    """9 SIGN — One signature shape. Three signing options.

    All three algorithms shown are stable on the same v1alpha1 surface:
    Plain RSA (bare public-key pinning, no PKI), GPG (OpenPGP keyrings),
    Sigstore (keyless via OIDC). PEM-encoded RSA exists but is genuinely
    experimental (CLI prints `experimental` warnings on sign/verify per
    bindings/go/rsa/signing/handler/handler.go:114,181) — flagged on
    slide 14, not shown here.
    """
    s = prs.slides.add_slide(layouts["Content / 3-Column"])
    set_text(s, 1, "SIGN")
    set_text(s, 2, "Same signed object. Three signing options.")
    set_text(s, 10, "RSA")
    set_text(s, 11, "Bare public-key pinning.\n"
                     "If you already rotate a signing key.")
    set_text(s, 12, "GPG")
    set_text(s, 13, "OpenPGP keys, ASCII-armored.\n"
                     "If your team runs a keyring.")
    set_text(s, 14, "SIGSTORE")
    set_text(s, 15, "Keyless via OIDC + Rekor.\n"
                     "If you already trust your identity provider.")
    # Punchline moved to speaker notes: "Same descriptor hash. Three
    # ways to vouch for it. Pick what your org already runs."


def build_slide_10_transport(prs, layouts):
    """10 TRANSPORT — Three patterns. One command."""
    s = prs.slides.add_slide(layouts["Content / 3-Column"])
    set_text(s, 1, "TRANSPORT")
    set_text(s, 2, "Three patterns. One command.")
    set_text(s, 10, "REGISTRY → REGISTRY")
    set_text(s, 11, "Promote across stages.\n"
                     "Source registry to target registry.")
    set_text(s, 12, "REGISTRY → CTF")
    set_text(s, 13, "Export to a local archive.\n"
                     "Hand-carry across the boundary.")
    set_text(s, 14, "CTF → REGISTRY")
    set_text(s, 15, "Air-gap import.\n"
                     "Verify on arrival. No callback to source.")
    # CTF acronym is used in the third column without definition — add a
    # footer below the columns so an architect glancing at the slide
    # doesn't have to wait for the speaker to define it. Wording matches
    # the website (transfer-concept.md): "filesystem-based OCM repository".
    # Left-aligned to anchor with the title; Pt20 regular (not italic) so
    # it reads as a definition the audience can use, not a disclaimer.
    ctf_tb, ctf_tf = add_textbox(s, 120, 980, SLIDE_W_PX - 240, 60)
    cp = ctf_tf.paragraphs[0]
    cp.alignment = PP_ALIGN.LEFT
    cr = cp.add_run()
    cr.text = ("CTF = Common Transport Format — a filesystem-based "
               "OCM repository, portable via any transfer mechanism.")
    cr.font.name = "Aptos"
    cr.font.size = Pt(20)
    cr.font.color.rgb = C.GREY_MID

    # AIR-GAP tag above the third column — pulls the dramatic move out
    # of the trio. Air-gap is the slide's headline use case but visually
    # the three columns weigh the same; this tag fixes that. Positioned
    # just below the title bottom (~508) and above the third column header.
    # Sized Pt22 (one Pt below the Pt23 column-header weight) so it sits
    # in the same typographic family as the column titles, not below them.
    tag_tb, tag_tf = add_textbox(s, 1320, 520, 320, 40)
    tp = tag_tf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    tr = tp.add_run()
    tr.text = "AIR-GAP"
    tr.font.name = "Aptos"
    tr.font.size = Pt(22)
    tr.font.bold = True
    tr.font.color.rgb = C.BLUE
    # Punchline moved to speaker notes: "Access fields rewrite at
    # transfer. Digests don't. Signature still verifies — anywhere."


def _render_chain_cards(slide, *,
                        cards,
                        cards_y, card_w, card_h, gap,
                        stripe_h,
                        label_size, label_color,
                        body_size, body_color,
                        stripe_color,
                        arrow_color, arrow_stroke_pt,
                        arrow_y_offset,
                        shadow=True,
                        label_pad_x=30, label_pad_y=42,
                        body_pad_y=110,
                        body_align=PP_ALIGN.LEFT):
    """Render a horizontal chain of cards in the slide-7/11 family.

    Each card is a borderless rounded rectangle with a flush-top brand
    rectangle (the "stripe"), an optional soft drop shadow, an ALL-CAPS
    label, and a multi-line body. Adjacent cards are connected by a thin
    arrow (stem rectangle + triangular head).

    Args:
        slide: the python-pptx slide to draw on.
        cards: list of (label, body_runs) where body_runs is a list of
            (text, bold) tuples — one run per body line.
        cards_y: vertical position of the top of the card row (slide-px).
        card_w / card_h: per-card width and height (slide-px).
        gap: horizontal gap between adjacent cards.
        stripe_h: height of the top stripe (slide-px); 4-5 typical.
        label_size / label_color: label typography (Aptos bold).
        body_size / body_color: body typography (Aptos regular, bold runs
            opt-in via the per-line `bold` flag in `body_runs`).
        stripe_color: top-stripe fill color (RGBColor).
        arrow_color / arrow_stroke_pt: arrow fill + line weight.
        arrow_y_offset: arrow vertical center relative to card top.
        shadow: whether to attach the soft outerShdw drop shadow.
        label_pad_x / label_pad_y: label position inside the card
            (offset from card top-left).
        body_pad_y: body vertical offset from card top.
        body_align: paragraph alignment for both label and body
            (label inherits this — keep them aligned together).

    Used by slide 11 (DEPLOY chain — all brand colors) and slide 16
    (Replication appendix — chain in grey + Replication card pulled out
    in brand). Same shape, different color palettes via the parameters.
    """
    total_w = len(cards) * card_w + (len(cards) - 1) * gap
    start_x = (SLIDE_W_PX - total_w) // 2

    for i, (label, body_runs) in enumerate(cards):
        x = start_x + i * (card_w + gap)

        # Card body — borderless rounded rect with optional drop shadow.
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            px(x), px(cards_y), px(card_w), px(card_h),
        )
        # Corner radius ~14 px regardless of card size (matches slide_6_native).
        card.adjustments[0] = 14.0 / min(card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = C.GREY_SOFT
        card.line.fill.background()  # no outline

        if shadow:
            # Soft outerShdw — same params as slide_6_native (blur 3 px,
            # offset 0,3 px straight down, 30% black).
            spPr = card._element.spPr
            for old in spPr.findall(f"{{{A_NS}}}effectLst"):
                spPr.remove(old)
            effectLst = etree.SubElement(spPr, f"{{{A_NS}}}effectLst")
            outerShdw = etree.SubElement(effectLst, f"{{{A_NS}}}outerShdw")
            outerShdw.set("blurRad", "28575")
            outerShdw.set("dist",    "28575")
            outerShdw.set("dir",     "5400000")
            outerShdw.set("rotWithShape", "0")
            clr = etree.SubElement(outerShdw, f"{{{A_NS}}}srgbClr")
            clr.set("val", "000000")
            alpha = etree.SubElement(clr, f"{{{A_NS}}}alpha")
            alpha.set("val", "30000")

        # Top stripe — flush rectangle on the card's top edge.
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            px(x), px(cards_y), px(card_w), px(stripe_h),
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = stripe_color
        stripe.line.fill.background()

        # Label — ALL-CAPS Aptos bold, left-aligned by default.
        lbl_tb = slide.shapes.add_textbox(
            px(x + label_pad_x), px(cards_y + label_pad_y),
            px(card_w - 2 * label_pad_x), px(50),
        )
        lbl_tf = lbl_tb.text_frame
        lbl_tf.margin_left = lbl_tf.margin_right = 0
        lbl_tf.margin_top = lbl_tf.margin_bottom = 0
        lbl_tf.word_wrap = True
        lp = lbl_tf.paragraphs[0]
        lp.alignment = body_align
        lr = lp.add_run()
        lr.text = label
        lr.font.name = "Aptos"
        lr.font.size = Pt(label_size)
        lr.font.bold = True
        lr.font.color.rgb = label_color

        # Body — one paragraph per body run (each is a separate line).
        # Bold can be set per run for emphasis (e.g., slide 11 Component
        # card bolds "Verifies its signature." to mark the verification
        # cliff edge).
        bdy_tb = slide.shapes.add_textbox(
            px(x + label_pad_x), px(cards_y + body_pad_y),
            px(card_w - 2 * label_pad_x), px(card_h - body_pad_y - 20),
        )
        bdy_tf = bdy_tb.text_frame
        bdy_tf.margin_left = bdy_tf.margin_right = 0
        bdy_tf.margin_top = bdy_tf.margin_bottom = 0
        bdy_tf.word_wrap = True
        for j, (run_text, run_bold) in enumerate(body_runs):
            p = bdy_tf.paragraphs[0] if j == 0 else bdy_tf.add_paragraph()
            p.alignment = body_align
            p.line_spacing = 1.25
            r = p.add_run()
            r.text = run_text
            r.font.name = "Aptos"
            r.font.size = Pt(body_size)
            r.font.color.rgb = body_color
            if run_bold:
                r.font.bold = True

        # Arrow into the next card — stem rectangle + triangle head
        # (MSO_SHAPE.RIGHT_TRIANGLE rotated 90° = right-pointing).
        if i < len(cards) - 1:
            arrow_y = cards_y + arrow_y_offset
            arrow_x1 = x + card_w + 10
            arrow_x2 = x + card_w + gap - 10
            stem_h = arrow_stroke_pt * 0.9
            head_w = 16
            stem_w = max(0, (arrow_x2 - arrow_x1) - head_w + 1)
            if stem_w > 0:
                stem = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    px(arrow_x1), px(arrow_y - stem_h / 2.0),
                    px(stem_w), px(stem_h),
                )
                stem.fill.solid()
                stem.fill.fore_color.rgb = arrow_color
                stem.line.fill.background()
            # Triangle head: build via freeform for a symmetric isoceles.
            head_x = arrow_x2 - head_w
            head_y_top = arrow_y - 7
            head_y_bot = arrow_y + 7
            builder = slide.shapes.build_freeform(
                px(head_x), px(arrow_y), scale=1.0,
            )
            builder.add_line_segments([
                (px(arrow_x2), px(head_y_top)),
                (px(arrow_x2), px(head_y_bot)),
            ], close=True)
            head = builder.convert_to_shape()
            head.fill.solid()
            head.fill.fore_color.rgb = arrow_color
            head.line.fill.background()


def build_slide_11_deploy(prs, layouts):
    """11 DEPLOY — Four CRs, one chain. Slide-7 card family.

    Cards are borderless rounded rectangles with a brand-blue top stripe,
    soft drop shadow, ALL-CAPS mid-blue labels left-aligned, dark-grey
    bodies left-aligned. Connecting arrows are thin (stem + triangle head)
    in brand blue.

    The Component card's second body line ("Verifies its signature.") is
    bold — the security-architect cue that tells the audience where the
    chain stops on bad signatures. Other three cards stay regular weight.

    Geometry per the user's PPT extraction (slide-7 card family at a
    smaller scale appropriate for a four-CR chain without icons).

    Replication is the fifth controller; it sits alongside the chain
    rather than within it. Moved to appendix slide 16 — Q&A backup, not
    part of the four-card architectural message.
    """
    s = prs.slides.add_slide(layouts["Plain"])
    set_text(s, 1, "DEPLOY")
    set_text(s, 2, "OCM controllers verify and apply.")
    delete_placeholder(s, 10)

    _render_chain_cards(
        s,
        cards=[
            ("REPOSITORY", [("Where component versions live.", False)]),
            ("COMPONENT",  [("Pulls one version.", False),
                            ("Verifies its signature.", True)]),
            ("RESOURCE",   [("One artifact, by digest.", False)]),
            ("DEPLOYER",   [("Applies it to the cluster.", False)]),
        ],
        cards_y=540, card_w=361, card_h=265, gap=76,
        stripe_h=5,
        label_size=30, label_color=C.BLUE_MID,
        body_size=22, body_color=RGBColor(0x33, 0x33, 0x33),
        stripe_color=C.BLUE,
        arrow_color=C.BLUE, arrow_stroke_pt=2.0,
        arrow_y_offset=132,  # arrow vertical center relative to card top
    )

    # Replication moved to appendix slide 16 — Q&A backup, not part of
    # the four-card chain. Speaker keeps it in their back pocket for
    # cluster-side mirroring questions.

    # Punchline moved to speaker notes: "Pluggable at the Deployer tier:
    # built-in for raw manifests; Flux for HelmReleases; Argo for
    # Applications."


def build_slide_12_composition(prs, layouts):
    """12 DAY 2 — Bump the product. Everything follows.

    Composition is already introduced on slide 8. This slide does ONE job:
    the day-2 upgrade mechanic. Same product, bumped version, references
    follow, signature is recomputed end-to-end.

    Layout changes 2026-06-24:
    - External DAY 1 / DAY 2 labels above each YAML block dropped, AND
      the `# day N — product X` comment lines inside each YAML are also
      dropped. The slide-level `DAY 2` eyebrow + the "bump version" arrow
      label carry the same information without YAML chrome.
    - Each YAML now ends with a 3-line `signatures:` block so the footer
      ("Every digest pinned by the signature. The cluster cannot drift.")
      is grounded in what the slide visually shows. The day-2 signature
      value is highlighted in brand blue alongside the version changes,
      so the visual diff reads three changes that tell one story: bump
      one line, the whole chain re-signs.
    - Arrow between the blocks carries the label "bump version" — the
      arrow is the operator action, the label names it.
    - Footer caption beneath both blocks delivers the differentiator
      that previously lived only in speaker notes.
    """
    s = prs.slides.add_slide(layouts["Content / Diagram Compact"])
    set_text(s, 1, "DAY 2")
    set_text(s, 2, "Bump the product version.\nEverything follows.")
    delete_placeholder(s, 10)

    K = C.BLUE_MID
    V = C.BLACK
    COM = C.GREY_MID
    HIGH = C.BLUE

    left_yaml = [
        ("component:",                                                K),
        ("  name: acme.org/sovereign/product",                        V),
        ("  version: 1.0.0",                                          V),
        ("  componentReferences:",                                    K),
        ("    - name: notes",                                         V),
        ("      componentName: acme.org/sovereign/notes",             V),
        ("      version: 1.0.0",                                      V),
        ("      digest:                                               # of the child descriptor", COM),
        ("        hashAlgorithm: SHA-256",                            V),
        ("        value: 7a1b2c3d4e...",                              V),
        ("    - name: postgres",                                      V),
        ("      componentName: acme.org/sovereign/postgres",          V),
        ("      version: 1.0.0",                                      V),
        ("      digest:",                                             COM),
        ("        hashAlgorithm: SHA-256",                            V),
        ("        value: f5e4d3c2b1...",                              V),
        ("signatures:",                                               K),
        ("  - name: acme-release-key",                                V),
        ("    signature:",                                            K),
        ("      algorithm: RSASSA-PSS",                               V),
        ("      value: a4b1c2d3e5f6789abc012345def04691...",          V),
    ]
    right_yaml = [
        ("component:",                                                K),
        ("  name: acme.org/sovereign/product",                        V),
        ("  version: 1.1.0",                                          HIGH),
        ("  componentReferences:",                                    K),
        ("    - name: notes",                                         V),
        ("      componentName: acme.org/sovereign/notes",             V),
        ("      version: 1.1.0",                                      HIGH),
        ("      digest:                                               # of the child descriptor", COM),
        ("        hashAlgorithm: SHA-256",                            V),
        ("        value: 9b8a7c6d5e...",                              HIGH),
        ("    - name: postgres",                                      V),
        ("      componentName: acme.org/sovereign/postgres",          V),
        ("      version: 1.0.0",                                      V),
        ("      digest:",                                             COM),
        ("        hashAlgorithm: SHA-256",                            V),
        ("        value: f5e4d3c2b1...",                              V),
        ("signatures:",                                               K),
        ("  - name: acme-release-key",                                V),
        ("    signature:",                                            K),
        ("      algorithm: RSASSA-PSS",                               V),
        ("      value: 9c2af18b3e7d52914a8c6b0f1d2e8f37...",          HIGH),
    ]
    # 21 lines @ Pt17 — reduced from 15 lines @ Pt20 to fit reference
    # digests. Box grown to 600px, top moved up to 370 to keep footer
    # at y=985.
    # YAML carries componentName: on each reference + algorithm: in the
    # signature: sub-block per spec — without the mediaType: line, which
    # adds verification chrome without signal at this altitude.
    box_w = 740
    box_h = 600
    box_y = 370
    arrow_w = 120
    arrow_h = 100
    gap = 20
    total_w = 2 * box_w + 2 * gap + arrow_w
    left_x = (SLIDE_W_PX - total_w) // 2
    right_x = left_x + box_w + gap + arrow_w + gap
    arrow_x = left_x + box_w + gap
    arrow_y = box_y + (box_h - arrow_h) // 2

    add_yaml_block(s, x_px=left_x, y_px=box_y, w_px=box_w, h_px=box_h,
                    yaml_lines=left_yaml, font_size=17)
    add_yaml_block(s, x_px=right_x, y_px=box_y, w_px=box_w, h_px=box_h,
                    yaml_lines=right_yaml, font_size=17)

    # Arrow with "bump version" label above it. Label sits just above the
    # arrow's vertical center so the eye reads label-then-arrow as one unit.
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                px(arrow_x), px(arrow_y),
                                px(arrow_w), px(arrow_h))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = C.BLUE
    arrow.line.fill.background()

    label_tb, label_tf = add_textbox(s, arrow_x - 30, arrow_y - 60,
                                       arrow_w + 60, 50)
    lp = label_tf.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run()
    lr.text = "bump version"
    lr.font.name = "Aptos"
    lr.font.size = Pt(20)
    lr.font.bold = True
    lr.font.color.rgb = C.BLUE

    # Footer caption — the differentiator. Single line, centred under both
    # YAML blocks, in brand blue so it reads as the conclusion of the diff
    # rather than chrome. Brand blue matches the highlighted day-2 changes
    # above (version bumps + new signature value) — same color register
    # says "these changes produce this consequence". This sentence answers
    # the most common architect challenge to OCM ("how is this different
    # from helm upgrade?") with a property the audience can verify from
    # the diagram above.
    footer_tb, footer_tf = add_textbox(s, 120, 985, SLIDE_W_PX - 240, 60)
    fp = footer_tf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = "Every digest pinned by the signature. The cluster cannot drift."
    fr.font.name = "Aptos"
    fr.font.size = Pt(24)
    fr.font.color.rgb = C.BLUE


def build_slide_13_whats_sharp(prs, layouts):
    """13 WHAT'S SHARP — Three real edges from the docs/repo.

    All three are doc- and code-confirmed limitations a first-time architect
    will hit:
      1. Transfer defaults to descriptor-only.
         (transfer-concept.md:42 — air-gap requires --copy-resources)
      2. Controllers ship as v1alpha1 — pin to specific release tags.
         (kubernetes/controller/api/v1alpha1)
      3. Helm-deploy adds kro + Flux dependencies.
         (ocm-controllers.md:116 — Helm path requires kro + Flux alongside)
    """
    s = prs.slides.add_slide(layouts["Plain / Compact"])
    set_text(s, 1, "WHAT'S SHARP")
    set_text(s, 2, "Three honest edges.")
    set_blue_box_bullets(s, 10, [
        "Transfer defaults — copies only the descriptor. "
        "For air-gap, pass --copy-resources so the bytes travel too.",
        "Controllers are v1alpha1 — the CRD surface can move. "
        "Pin to specific release tags in your platform installs.",
        "Helm-deploy adds kro + Flux — the OCM controllers don't ship "
        "them. Bring your existing GitOps engine.",
    ], font_size=26)
    # Punchline moved to speaker notes: "Honest now beats apologetic
    # later. Plan for the trim edge."


def build_slide_14_adoption(prs, layouts):
    """14 ADOPTION — Two paths to a first OCM component.
    Uses Plain/Compact layout with custom 2-column textboxes (the
    Content/2-Column layout has only body placeholders, no header slots)."""
    s = prs.slides.add_slide(layouts["Plain / Compact"])
    set_text(s, 1, "ADOPTION")
    set_text(s, 2, "Two paths to a first OCM component.")
    delete_placeholder(s, 10)

    # Column geometry: two equal columns under the title.
    col_y = 500
    col_h = 380
    col_w = 820
    gap = 80
    total_w = 2 * col_w + gap
    start_x = (SLIDE_W_PX - total_w) // 2  # 100 if SLIDE_W_PX=1920

    columns = [
        ("FROM ZERO — CLI",
         ["Pack one component. Sign it.",
          "Air-gap CTF round-trip.",
          "Verify on the other side.",
          "Thirty minutes on a laptop."]),
        ("ON YOUR CLUSTER — CONTROLLERS",
         ["Helm-install the OCM controllers.",
          "Point them at your registry.",
          "Deploy a component.",
          "Thirty minutes on any cluster."]),
    ]
    for i, (header, lines) in enumerate(columns):
        x = start_x + i * (col_w + gap)
        # Header
        head_tb, head_tf = add_textbox(s, x, col_y, col_w, 56)
        hp = head_tf.paragraphs[0]
        hp.alignment = PP_ALIGN.LEFT
        hr = hp.add_run()
        hr.text = header
        hr.font.name = "Aptos"
        hr.font.size = Pt(22)
        hr.font.bold = True
        hr.font.color.rgb = C.BLUE
        # Thin horizontal rule under header
        rule = s.shapes.add_connector(1, px(x), px(col_y + 56),
                                       px(x + col_w), px(col_y + 56))
        rule.line.color.rgb = C.BLUE
        rule.line.width = Pt(1.25)
        # Body lines
        body_tb, body_tf = add_textbox(s, x, col_y + 72, col_w, col_h - 72)
        for j, line in enumerate(lines):
            p = body_tf.paragraphs[0] if j == 0 else body_tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(0) if j == 0 else Pt(10)
            r = p.add_run()
            r.text = line
            r.font.name = "Aptos"
            r.font.size = Pt(22)
            r.font.color.rgb = C.BLACK

    # Punchline moved to speaker notes: "Pick the path. Coexists with
    # cosign, Argo, Flux, Kyverno — OCM signs the descriptor; your
    # existing controls stay in place."


def build_slide_15_cta(prs, layouts):
    """15 CTA — Architect-shaped: Evaluate / Pilot / Engage.

    Reshaped from the developer-arc (Try it / Build with us / Talk to us)
    to mirror the architect's evaluation arc (read the spec → scoped pilot
    → engage with standard stewardship). Mirrors the exec deck's
    Pilot/Evaluate/Engage pattern. Layout placeholders only, no custom
    textboxes."""
    s = prs.slides.add_slide(layouts["CTA"])
    set_text(s, 1, "Ship the release as one unit.", color=C.WHITE)
    set_action_path_lines(s, 2, [
        ("Evaluate", "ocm.software · run conformance/scenarios/sovereign"),
        ("Pilot",    "github.com/open-component-model · one product, one team"),
        ("Engage",   "community channels on the website · NeoNephos Foundation"),
    ])
    add_brand_row(s)


def build_slide_16_appendix_replication(prs, layouts):
    """16 APPENDIX — Replication (pull-on-demand).

    Out of the main 15-slide arc. Pulled only if the audience asks about
    cluster-side mirroring or repo-to-repo transfer without the CLI.

    Visually consistent with slide 11: same card family (rounded rect, no
    border, top stripe, soft shadow, ALL-CAPS left-aligned label, dark-
    grey body). Two distinct treatments separate the "chain echo" from
    the "Replication highlight":

      Top row — four chain cards (Repository / Component / Resource /
        Deployer) in GREY. Grey top stripe, grey labels, grey body. Same
        card shape as slide 11 but visually dimmed — "this is the chain
        from slide 11, not the focus of this slide". Smaller cards too
        (~330×200) reinforcing the "echo" reading.

      Bottom — single Replication card in BRAND COLORS. Brand-blue top
        stripe, mid-blue label, dark-grey body. Wider (~700×260) and
        offset below the chain. This is the slide's actual subject.

    Footer caption beneath the Replication card delivers the one-line
    "controller equivalent of `ocm transfer cv`" framing.

    Replication facts verified against website/content/docs/reference/
    kubernetes-api/replication.md and the user's confirmation that
    `status.lastTransferredDigest` is correct.
    """
    s = prs.slides.add_slide(layouts["Plain"])
    set_text(s, 1, "APPENDIX · REPLICATION")
    set_text(s, 2, "Alongside the chain. Not within it.")
    delete_placeholder(s, 10)

    # --- Top row: dimmed chain in slide-7/11 family with grey stripe ---
    _render_chain_cards(
        s,
        cards=[
            ("REPOSITORY", [("Where component versions live.", False)]),
            ("COMPONENT",  [("Pulls + verifies.", False)]),
            ("RESOURCE",   [("One artifact, by digest.", False)]),
            ("DEPLOYER",   [("Applies it to the cluster.", False)]),
        ],
        cards_y=520, card_w=330, card_h=200, gap=60,
        stripe_h=4,
        label_size=22, label_color=C.GREY_MID,
        body_size=16, body_color=C.GREY_MID,
        stripe_color=C.GREY_MID,
        arrow_color=C.GREY_MID, arrow_stroke_pt=1.8,
        arrow_y_offset=100,
        shadow=True,
        label_pad_x=24, label_pad_y=32,
        body_pad_y=90,
    )

    # --- Bottom: highlighted Replication card in slide-7/11 family ---
    rep_w = 700
    rep_h = 260
    rep_x = (SLIDE_W_PX - rep_w) // 2
    rep_y = 520 + 200 + 60  # chain row bottom + 60px breathing space

    rep_card = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        px(rep_x), px(rep_y), px(rep_w), px(rep_h),
    )
    rep_card.adjustments[0] = 14.0 / min(rep_w, rep_h)
    rep_card.fill.solid()
    rep_card.fill.fore_color.rgb = C.GREY_SOFT
    rep_card.line.fill.background()

    # Drop shadow (same as chain cards above).
    spPr = rep_card._element.spPr
    for old in spPr.findall(f"{{{A_NS}}}effectLst"):
        spPr.remove(old)
    effectLst = etree.SubElement(spPr, f"{{{A_NS}}}effectLst")
    outerShdw = etree.SubElement(effectLst, f"{{{A_NS}}}outerShdw")
    outerShdw.set("blurRad", "28575")
    outerShdw.set("dist",    "28575")
    outerShdw.set("dir",     "5400000")
    outerShdw.set("rotWithShape", "0")
    clr = etree.SubElement(outerShdw, f"{{{A_NS}}}srgbClr")
    clr.set("val", "000000")
    alpha = etree.SubElement(clr, f"{{{A_NS}}}alpha")
    alpha.set("val", "30000")

    # Brand-blue top stripe.
    rep_stripe = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        px(rep_x), px(rep_y), px(rep_w), px(5),
    )
    rep_stripe.fill.solid()
    rep_stripe.fill.fore_color.rgb = C.BLUE
    rep_stripe.line.fill.background()

    # Label "REPLICATION" left-aligned mid-blue 30pt bold.
    rlbl_tb = s.shapes.add_textbox(
        px(rep_x + 30), px(rep_y + 42), px(rep_w - 60), px(50),
    )
    rlbl_tf = rlbl_tb.text_frame
    rlbl_tf.margin_left = rlbl_tf.margin_right = 0
    rlbl_tf.margin_top = rlbl_tf.margin_bottom = 0
    rlbl_tf.word_wrap = True
    lp = rlbl_tf.paragraphs[0]
    lp.alignment = PP_ALIGN.LEFT
    lr = lp.add_run()
    lr.text = "REPLICATION"
    lr.font.name = "Aptos"
    lr.font.size = Pt(30)
    lr.font.bold = True
    lr.font.color.rgb = C.BLUE_MID

    # Body — two lines, left-aligned dark-grey 22pt.
    rbdy_tb = s.shapes.add_textbox(
        px(rep_x + 30), px(rep_y + 110), px(rep_w - 60), px(rep_h - 130),
    )
    rbdy_tf = rbdy_tb.text_frame
    rbdy_tf.margin_left = rbdy_tf.margin_right = 0
    rbdy_tf.margin_top = rbdy_tf.margin_bottom = 0
    rbdy_tf.word_wrap = True
    body_lines = [
        "Transfers a resolved component version from one OCM repository to another.",
        "Records status.lastTransferredDigest. Same digest → no-op.",
    ]
    DARK_GREY = RGBColor(0x33, 0x33, 0x33)
    for j, line in enumerate(body_lines):
        p = rbdy_tf.paragraphs[0] if j == 0 else rbdy_tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        r = p.add_run()
        r.text = line
        r.font.name = "Aptos"
        r.font.size = Pt(22)
        r.font.color.rgb = DARK_GREY

    # Footer caption — at y=1010 (NOT y=1090; the PPT user had it off-
    # slide at y=1090 which is below the 1080-px slide bottom). Mid-blue
    # Pt20 centered, not bold. Footer states the mechanism (what the CR
    # binds and what invariant it maintains) — the use cases (promotion,
    # mirroring, in-cluster air-gap) are implicit because slide 10 has
    # already named them as `ocm transfer cv` use cases.
    footer_tb, footer_tf = add_textbox(
        s, 120, 1010, SLIDE_W_PX - 240, 50,
    )
    fp = footer_tf.paragraphs[0]
    fp.alignment = PP_ALIGN.CENTER
    fr = fp.add_run()
    fr.text = ("Controller-shaped equivalent of the OCM CLI's `ocm transfer cv` — "
               "point it at a source `Component` and a target `Repository`, "
               "and it keeps them in sync.")
    fr.font.name = "Aptos"
    fr.font.size = Pt(20)
    fr.font.color.rgb = C.BLUE_MID



# -----------------------------------------------------------------------------
# Sanity check
# -----------------------------------------------------------------------------

# -----------------------------------------------------------------------------

def sanity_check(prs):
    issues = []
    for i, s in enumerate(prs.slides, 1):
        if s.element.get("show") == "0":
            continue
        title_bottom_px = None
        for ph in s.placeholders:
            if ph.placeholder_format.idx == 2 and ph.top is not None:
                title_bottom_px = (ph.top + ph.height) / PX
        if title_bottom_px is None:
            continue
        for shape in s.shapes:
            if not shape.has_text_frame:
                continue
            if shape.is_placeholder:
                continue
            if shape.top is None:
                continue
            shape_top_px = shape.top / PX
            if shape_top_px < title_bottom_px - 4:
                issues.append(
                    f"slide {i}: shape top {shape_top_px:.0f}px overlaps "
                    f"title bottom {title_bottom_px:.0f}px"
                )
    if issues:
        print("WARNING: layout overlaps detected:")
        for issue in issues:
            print(f"  - {issue}")
    else:
        print("Layout sanity check passed.")


# =============================================================================
if __name__ == "__main__":
    if not shutil.which("rsvg-convert"):
        sys.exit("rsvg-convert not found; install via `brew install librsvg`")
    build()

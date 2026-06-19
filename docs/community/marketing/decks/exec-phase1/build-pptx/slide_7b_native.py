"""
slide_7b_native.py — Slide 7b "Sovereign-ready: trust, but verify." rendered
as native PowerPoint shapes (rounded rectangles, lines, freeform paths,
textboxes, recoloured icons), not a rasterized SVG.

Why native: SVG rasterization throws away typography, links, and
edit-in-PowerPoint affordances. For diagrams that are mostly boxes + text +
a single connecting arrow (like this one), drawing them with python-pptx
primitives keeps the slide hand-editable and crisp at any zoom.

Composition (1:1 with diagrams/06-sovereign-airgap.svg, viewBox 1600×760):

  Source side (left, no container)
    ─ "SOURCE" eyebrow (grey, letter-spaced)
    ─ "Pack · Sign" headline (mid-blue 32pt bold)
    ─ Public-registry pillar (white rounded rect, soft shadow)
        · ico-registry (brand-blue)
        · "Public registry" + "Build artifacts live here."
        · Signed-component glyph (single box: white body + grey stub +
          hairline divider + brand-blue lock badge + name + version)

  Trust boundary (vertical dashed brand-blue line, "TRUST BOUNDARY" label
  above)

  Target side (right, dashed brand-blue rounded container = "AIR-GAPPED")
    ─ "SOVEREIGN TARGET" eyebrow (brand-blue, letter-spaced)
    ─ "Verify · Deploy" headline (mid-blue 32pt bold)
    ─ Air-gapped container with "AIR-GAPPED" label inside top-left
    ─ Receiving component glyph (identical to source-side glyph)
    ─ Three local-consumer tiles in a row (Local registry / K8s cluster /
      Auditor) each with icon + label + green check-badge
    ─ "Trust, but verify." + "No callback to source. Day-2 ops included."
    ─ Bottom accent: light-blue rounded rect with
      "SAME IDENTITY · SAME SIGNATURE · ANY LOCATION"

  Crossing arrow (rendered last so it draws on top of the air-gapped
  container's dashed border): cubic-bezier freeform path from source-
  component right edge to receiving-component left edge, brand-blue 3.5pt,
  triangular arrowhead at the target end.

  "TRANSPORT" label above the arrow's mid-flight.

Usage:
    from slide_7b_native import add_sovereign_airgap_native
    add_sovereign_airgap_native(slide, x=60, y=240, w=1800, h=780,
                                 icons_dir=ICONS_DIR,
                                 rasterize_recolored=rasterize_svg_recolored)
"""
from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


PX = 9525


def px(n: float) -> Emu:
    return Emu(int(n * PX))


# OCM brand palette (mirror of build_pptx C class — kept local so this module
# is self-contained and importable by both deck builders).
BLUE       = RGBColor(0x0F, 0x6B, 0xFF)   # brand-blue
BLUE_MID   = RGBColor(0x0A, 0x3A, 0x99)   # mid blue used for headlines
BLUE_LIGHT = RGBColor(0xE8, 0xF0, 0xFF)   # accent / soft pillar stroke
GREY_MID   = RGBColor(0x6B, 0x72, 0x80)   # secondary text
GREY_HAIR  = RGBColor(0xE5, 0xE7, 0xEB)   # hairline divider
GREY_STUB  = RGBColor(0xF3, 0xF4, 0xF6)   # version stub fill
GREEN_CHK  = RGBColor(0x16, 0xA3, 0x4A)   # check-badge fill (deck exception)
BLACK      = RGBColor(0x00, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)


def _styled_run(p, text, *, size_pt, bold=False, color=BLACK,
                font="Inter", italic=False, all_caps=False,
                letter_spacing=None):
    """Add a run with the OCM type system applied."""
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size_pt)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if all_caps or letter_spacing is not None:
        rPr = r._r.get_or_add_rPr()
        if all_caps:
            rPr.set("cap", "all")
        if letter_spacing is not None:
            rPr.set("spc", str(letter_spacing))
    return r


def _add_textbox(slide, x, y, w, h, *, anchor="t", align=None, margins=0):
    """Borderless textbox at the given coords, with consistent inset margins."""
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = px(margins)
    tf.margin_top = tf.margin_bottom = px(margins)
    tf.word_wrap = True
    if anchor == "ctr":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "b":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    if align is not None:
        tf.paragraphs[0].alignment = align
    return tf


def _add_rounded_rect(slide, x, y, w, h, *,
                      fill=WHITE, stroke=BLUE, stroke_pt=1.5,
                      corner_pct=0.04, dash=False):
    """Rounded rectangle with explicit fill, stroke, and corner radius."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h),
    )
    shape.adjustments[0] = corner_pct
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = stroke
        shape.line.width = Pt(stroke_pt)
        if dash:
            # python-pptx doesn't expose dash style directly on Line; patch XML.
            from pptx.oxml.ns import qn
            ln = shape.line._get_or_add_ln()
            # Remove any existing prstDash to avoid duplicates.
            for old in ln.findall(qn("a:prstDash")):
                ln.remove(old)
            from lxml import etree
            prstDash = etree.SubElement(ln, qn("a:prstDash"))
            prstDash.set("val", "dash")
    return shape


def _add_rect(slide, x, y, w, h, *, fill=WHITE, stroke=None, stroke_pt=1.0):
    """Plain rectangle (no rounded corners). Used for the version stub fill."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h),
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = stroke
        shape.line.width = Pt(stroke_pt)
    return shape


def _add_line(slide, x1, y1, x2, y2, *, color=BLUE, width_pt=1.5, dash=False):
    """Straight line connector between two points."""
    from pptx.enum.shapes import MSO_CONNECTOR
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, px(x1), px(y1), px(x2), px(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    if dash:
        from pptx.oxml.ns import qn
        from lxml import etree
        ln = line.line._get_or_add_ln()
        for old in ln.findall(qn("a:prstDash")):
            ln.remove(old)
        prstDash = etree.SubElement(ln, qn("a:prstDash"))
        prstDash.set("val", "dash")
    return line


def _add_oval(slide, cx, cy, r, *, fill=BLUE, stroke=None, stroke_pt=1.0):
    """Filled circle centred at (cx, cy) with radius r."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, px(cx - r), px(cy - r), px(2 * r), px(2 * r),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = stroke
        shape.line.width = Pt(stroke_pt)
    return shape


from icon_strokes import STROKE_THIN


def add_sovereign_airgap_native(slide, *, x, y, w, h,
                                  icons_dir: Path,
                                  rasterize_recolored,
                                  icon_stroke=STROKE_THIN):
    """Draw the sovereign-ready air-gap diagram natively into the given slot.

    icon_stroke selects Tabler-icon weight (icon_strokes.STROKE_THIN /
    STROKE_REGULAR / STROKE_BOLD). Default STROKE_THIN.

    Parameters:
        slide: python-pptx Slide object.
        x, y, w, h: diagram slot in slide coordinates (px @ 96 dpi).
        icons_dir: folder containing the Tabler-style SVGs used as glyphs.
        rasterize_recolored: callable (svg_path, target_w_px, color_hex) -> Path
            that returns a rasterized PNG of the SVG in the requested colour.
            Passed in (rather than imported) so this module stays free of
            build-script-specific dependencies.

    Layout:
        Source SVG viewBox is 1600×760. We map (svg_x, svg_y) into the slot
        via uniform scale fx = w/1600, fy = h/760. Coordinates inside this
        function are kept in SVG space and transformed at draw-time via the
        local _sx / _sy helpers, so the math stays legible against the SVG.
    """
    fx = w / 1600.0
    fy = h / 760.0

    def _sx(sx):
        return x + sx * fx

    def _sy(sy):
        return y + sy * fy

    def _sw(sw):
        return sw * fx

    def _sh(sh):
        return sh * fy

    # ------------------------------------------------------------------
    # SOURCE SIDE (left)
    # ------------------------------------------------------------------
    # SVG group at translate(80,100). Inside that group:
    #   "SOURCE" eyebrow text at local (0, 40) → absolute (80, 140).
    #   "Pack · Sign" headline at local (0, 78) → absolute (80, 178).
    # SVG <text y=...> sets the baseline; we render in textboxes which
    # anchor at the top, so we pull each baseline up by ~font-size to land
    # the visual position in the same place. 20pt eyebrow → ~24px climb;
    # 32pt headline → ~38px climb. We chose textbox tops accordingly:
    #   eyebrow top = svg_y(40) - 24 = 116
    #   headline top = svg_y(78) - 38 = 140
    tf = _add_textbox(slide, _sx(80), _sy(116), _sw(420), _sh(32))
    _styled_run(tf.paragraphs[0], "SOURCE",
                size_pt=14, bold=True, color=GREY_MID,
                all_caps=True, letter_spacing="300")

    tf = _add_textbox(slide, _sx(80), _sy(140), _sw(420), _sh(48))
    _styled_run(tf.paragraphs[0], "Pack · Sign",
                size_pt=22, bold=True, color=BLUE_MID)

    # Public-registry pillar — rounded rect 320×200 at SVG (80, 240).
    pillar_x, pillar_y, pillar_w, pillar_h = 80, 240, 320, 200
    _add_rounded_rect(slide,
                      _sx(pillar_x), _sy(pillar_y),
                      _sw(pillar_w), _sh(pillar_h),
                      fill=WHITE, stroke=BLUE_LIGHT, stroke_pt=2.0,
                      corner_pct=0.06)

    # Pillar header: ico-registry icon + "Public registry" + subtext.
    reg_icon = icons_dir / "registry.svg"
    if reg_icon.exists():
        png = rasterize_recolored(reg_icon, 96, "0F6BFF", stroke_width=icon_stroke)
        slide.shapes.add_picture(
            str(png),
            px(_sx(pillar_x + 20)), px(_sy(pillar_y + 20)),
            width=px(_sw(56)), height=px(_sh(56)),
        )
    tf = _add_textbox(slide,
                      _sx(pillar_x + 100), _sy(pillar_y + 28),
                      _sw(220), _sh(28))
    _styled_run(tf.paragraphs[0], "Public registry",
                size_pt=14, bold=True, color=BLUE_MID)
    tf = _add_textbox(slide,
                      _sx(pillar_x + 100), _sy(pillar_y + 56),
                      _sw(220), _sh(24))
    _styled_run(tf.paragraphs[0], "Build artifacts live here.",
                size_pt=10, color=GREY_MID)

    # Source-side component glyph at SVG (95, 350) — 290×70 footprint.
    # We need uniform scale for the glyph itself so the visual proportions
    # of the badge / divider / version text stay correct. The glyph is built
    # in SVG coordinates; the helper takes SVG-space (gx, gy) and we scale
    # via _sx/_sy/_sw/_sh inside it. Because the glyph composition uses many
    # small offsets, we instead pass scaled coordinates and accept that the
    # glyph helper draws in the *target* coordinate system once we wrap it
    # into a small adapter below.
    src_glyph_x, src_glyph_y = 95, 350
    _draw_glyph_in_slot(
        slide, src_glyph_x, src_glyph_y,
        sx=_sx, sy=_sy, sw=_sw, sh=_sh,
        icons_dir=icons_dir, rasterize_recolored=rasterize_recolored,
        icon_stroke=icon_stroke,
    )

    # ------------------------------------------------------------------
    # TRUST BOUNDARY
    # ------------------------------------------------------------------
    # SVG group translate(570,80); line from local (80, 40)..(80, 640).
    # Absolute: (650, 120) → (650, 720). Label at local (80, 30) → (650, 110).
    boundary_x = 650
    _add_line(slide,
              _sx(boundary_x), _sy(120),
              _sx(boundary_x), _sy(720),
              color=BLUE, width_pt=1.8, dash=True)
    tf = _add_textbox(slide,
                      _sx(boundary_x - 120), _sy(86),
                      _sw(240), _sh(28),
                      align=PP_ALIGN.CENTER)
    _styled_run(tf.paragraphs[0], "TRUST BOUNDARY",
                size_pt=10, bold=True, color=BLUE,
                all_caps=True, letter_spacing="300")

    # ------------------------------------------------------------------
    # TARGET SIDE (right) — air-gapped container
    # ------------------------------------------------------------------
    # SVG group translate(900,100). Eyebrow + headline mirror source side.
    tf = _add_textbox(slide, _sx(900), _sy(116), _sw(600), _sh(32))
    _styled_run(tf.paragraphs[0], "SOVEREIGN TARGET",
                size_pt=14, bold=True, color=BLUE,
                all_caps=True, letter_spacing="300")
    tf = _add_textbox(slide, _sx(900), _sy(140), _sw(600), _sh(48))
    _styled_run(tf.paragraphs[0], "Verify · Deploy",
                size_pt=22, bold=True, color=BLUE_MID)

    # Air-gapped container — rounded rect 600×480 at SVG (900, 240),
    # dashed brand-blue stroke.
    cont_x, cont_y, cont_w, cont_h = 900, 240, 600, 480
    _add_rounded_rect(slide,
                      _sx(cont_x), _sy(cont_y),
                      _sw(cont_w), _sh(cont_h),
                      fill=WHITE, stroke=BLUE, stroke_pt=2.0,
                      corner_pct=0.03, dash=True)

    # "AIR-GAPPED" label inside top-left of the container.
    tf = _add_textbox(slide,
                      _sx(cont_x + 20), _sy(cont_y + 20),
                      _sw(200), _sh(24))
    _styled_run(tf.paragraphs[0], "AIR-GAPPED",
                size_pt=9, bold=True, color=BLUE,
                all_caps=True, letter_spacing="300")

    # Receiving component glyph — IDENTICAL preset to source side.
    # SVG translate(900,100) + translate(0,140) + translate(155,60) →
    # absolute (1055, 300). Glyph is 290×70 → right edge at 1345, vertical
    # centre y = 335.
    rcv_glyph_x, rcv_glyph_y = 1055, 300
    _draw_glyph_in_slot(
        slide, rcv_glyph_x, rcv_glyph_y,
        sx=_sx, sy=_sy, sw=_sw, sh=_sh,
        icons_dir=icons_dir, rasterize_recolored=rasterize_recolored,
        icon_stroke=icon_stroke,
    )

    # Three local-consumer tiles in a row.
    # SVG: translate(900,100)+translate(0,140) puts the container origin at
    # absolute (900, 240). Tiles are at translate(30,170), (216,170),
    # (402,170) inside the container → absolute (930, 410), (1116, 410),
    # (1302, 410). Each 170×120.
    tiles = [
        (930,  "registry.svg",   "Local registry"),
        (1116, "kubernetes.svg", "K8s cluster"),
        (1302, "shield.svg",     "Auditor"),
    ]
    for tile_x, icon_file, label in tiles:
        _draw_tile_in_slot(
            slide, tile_x, 410, icon_file=icon_file, label=label,
            sx=_sx, sy=_sy, sw=_sw, sh=_sh,
            icons_dir=icons_dir, rasterize_recolored=rasterize_recolored,
            icon_stroke=icon_stroke,
        )

    # Tagline + subtitle, centred inside the container.
    # SVG centres at x=300 inside translate(900,240) → absolute x=1200.
    # text y=330 / y=360 within container → absolute 570 / 600. We shift
    # textbox tops up by ~font-size to compensate for SVG baselines.
    tf = _add_textbox(slide,
                      _sx(cont_x + 30), _sy(560),
                      _sw(540), _sh(30),
                      align=PP_ALIGN.CENTER)
    _styled_run(tf.paragraphs[0], "Trust, but verify.",
                size_pt=14, bold=True, color=BLUE_MID)
    tf = _add_textbox(slide,
                      _sx(cont_x + 30), _sy(594),
                      _sw(540), _sh(28),
                      align=PP_ALIGN.CENTER)
    _styled_run(tf.paragraphs[0], "No callback to source. Day-2 ops included.",
                size_pt=11, color=GREY_MID)

    # Bottom accent: light-blue rounded rect with letter-spaced caption.
    # SVG: x=60 + container origin 900 = 960; y=400 + 240 = 640;
    # 480×42; caption centred at x=300 inside container → absolute 1200.
    accent_x, accent_y, accent_w, accent_h = 960, 640, 480, 42
    _add_rounded_rect(slide,
                      _sx(accent_x), _sy(accent_y),
                      _sw(accent_w), _sh(accent_h),
                      fill=BLUE_LIGHT, stroke=None, corner_pct=0.20)
    tf = _add_textbox(slide,
                      _sx(accent_x), _sy(accent_y),
                      _sw(accent_w), _sh(accent_h),
                      anchor="ctr", align=PP_ALIGN.CENTER)
    _styled_run(tf.paragraphs[0],
                "SAME IDENTITY · SAME SIGNATURE · ANY LOCATION",
                size_pt=10, bold=True, color=BLUE_MID,
                all_caps=True, letter_spacing="200")

    # ------------------------------------------------------------------
    # CROSSING ARROW (rendered LAST — z-order requirement)
    # ------------------------------------------------------------------
    # The SVG renders the arrow last so it draws on top of the air-gapped
    # container's dashed border (otherwise the dashed line would visually
    # halt the arrow at x=900 instead of letting it land on the receiving
    # component). We preserve that ordering here.
    #
    # Path: cubic Bézier from (385, 385) to (1055, 335), control points
    # (600, 385) and (850, 335). 3.5pt stroke, brand-blue. SVG paints a
    # gradient (#5cd6ff → #0f6bff); python-pptx supports gradient line via
    # XML, but for simplicity and because the gradient is barely visible at
    # 3.5pt, we use a solid brand-blue stroke. Triangular arrowhead added
    # via the connector's tail-arrow XML.
    _draw_crossing_arrow(slide, _sx, _sy)

    # TRANSPORT label above the arrow's mid-flight. SVG renders "Transport"
    # at (735, 340), 32pt mid-blue bold. The spec asks for ALL-CAPS letter-
    # spaced — we honour the spec since the SVG's intent (per the comment
    # block in the source) was the letter-spaced caps style anyway.
    tf = _add_textbox(slide,
                      _sx(560), _sy(296),
                      _sw(350), _sh(40),
                      align=PP_ALIGN.CENTER)
    _styled_run(tf.paragraphs[0], "TRANSPORT",
                size_pt=18, bold=True, color=BLUE_MID,
                all_caps=True, letter_spacing="400")


# -----------------------------------------------------------------------------
# Internal helpers that need access to the slot-scaled coordinate transforms.
# Defined as module-level functions taking explicit sx/sy/sw/sh so they stay
# pure and testable.
# -----------------------------------------------------------------------------

def _draw_glyph_in_slot(slide, gx, gy, *, sx, sy, sw, sh,
                         icons_dir, rasterize_recolored, icon_stroke=None):
    """Component glyph drawn in slot-coordinates. (gx, gy) is the SVG-space
    top-left of the 290×70 footprint."""
    # Stub fill (drawn first, behind body+outline).
    _add_rect(slide,
              sx(gx + 220), sy(gy), sw(70), sh(70),
              fill=GREY_STUB, stroke=None)
    # Body fill (rounded, white).
    _add_rounded_rect(slide,
                      sx(gx), sy(gy), sw(220), sh(70),
                      fill=WHITE, stroke=None, corner_pct=0.10)
    # Outer enclosing stroke (rounded, brand-blue) — full 290×70 footprint.
    _add_rounded_rect(slide,
                      sx(gx), sy(gy), sw(290), sh(70),
                      fill=None, stroke=BLUE, stroke_pt=2.0,
                      corner_pct=0.10)
    # Hairline divider between body and stub.
    _add_line(slide,
              sx(gx + 220), sy(gy + 6),
              sx(gx + 220), sy(gy + 64),
              color=GREY_HAIR, width_pt=1.4)
    # Lock badge — filled brand-blue circle, centre (gx+30, gy+35), r=18.
    # Use sw for the radius so the circle stays circular under uniform
    # x-scale (the slot scale is uniform when w/h preserves the SVG aspect;
    # if not, we choose sw which matches the badge's natural horizontal
    # diameter).
    _add_oval(slide,
              sx(gx + 30), sy(gy + 35),
              sw(18), fill=BLUE)
    # Lock glyph — rasterized white version of icons/lock.svg.
    lock_path = icons_dir / "lock.svg"
    if lock_path.exists():
        png = rasterize_recolored(lock_path, 96, "FFFFFF", stroke_width=icon_stroke)
        slide.shapes.add_picture(
            str(png),
            px(sx(gx + 14)), px(sy(gy + 19)),
            width=px(sw(32)), height=px(sh(32)),
        )
    # Name + role on body. SVG <text y=...> is the baseline, so we shift
    # textbox top up by font-size to land visually equivalent.
    tf = _add_textbox(slide,
                      sx(gx + 62), sy(gy + 18),
                      sw(158), sh(22))
    _styled_run(tf.paragraphs[0], "github.com/acme/app",
                size_pt=10, bold=True, color=BLUE_MID, font="Consolas")
    tf = _add_textbox(slide,
                      sx(gx + 62), sy(gy + 40),
                      sw(158), sh(20))
    _styled_run(tf.paragraphs[0], "signed component",
                size_pt=9, color=GREY_MID)
    # Version on stub, centred.
    tf = _add_textbox(slide,
                      sx(gx + 220), sy(gy + 25),
                      sw(70), sh(22),
                      anchor="ctr", align=PP_ALIGN.CENTER)
    _styled_run(tf.paragraphs[0], "v1.0.0",
                size_pt=10, bold=True, color=BLUE_MID,
                letter_spacing="100")


def _draw_tile_in_slot(slide, tx, ty, *, icon_file, label,
                        sx, sy, sw, sh,
                        icons_dir, rasterize_recolored, icon_stroke=None):
    """One consumer tile in slot-coordinates. (tx, ty) is the SVG-space
    top-left of the 170×120 footprint."""
    _add_rounded_rect(slide,
                      sx(tx), sy(ty), sw(170), sh(120),
                      fill=WHITE, stroke=BLUE_LIGHT, stroke_pt=2.0,
                      corner_pct=0.08)
    icon_path = icons_dir / icon_file
    if icon_path.exists():
        png = rasterize_recolored(icon_path, 96, "0F6BFF", stroke_width=icon_stroke)
        slide.shapes.add_picture(
            str(png),
            px(sx(tx + 62)), px(sy(ty + 14)),
            width=px(sw(46)), height=px(sh(46)),
        )
    # Label centred horizontally below the icon. SVG baseline at y=80; we
    # use a textbox top at ~y=68 to land visually at the same height.
    tf = _add_textbox(slide,
                      sx(tx), sy(ty + 68),
                      sw(170), sh(22),
                      align=PP_ALIGN.CENTER)
    _styled_run(tf.paragraphs[0], label,
                size_pt=11, bold=True, color=BLUE_MID)
    # Check-badge centred horizontally near tile bottom.
    _draw_check_badge_in_slot(slide, tx + 85, ty + 100,
                               sx=sx, sy=sy, sw=sw, sh=sh)


def _draw_check_badge_in_slot(slide, cx, cy, *, sx, sy, sw, sh, r=11):
    """Filled green circle + white tick, scaled into the slot."""
    _add_oval(slide, sx(cx), sy(cy), sw(r), fill=GREEN_CHK)
    # White tick. SVG path: M -5 0 L -1.5 4 L 5 -4 (relative to badge centre).
    # FreeformBuilder takes ABSOLUTE local-unit coords for start AND each
    # segment; values are multiplied by `scale` to derive EMU bbox.
    # Work in slide pixels (scale=PX). Wrapping the start in px() (EMU) and
    # then scaling again by PX produced a 9525x too-large bbox that
    # PowerPoint flagged as a corrupt file (repair dialog).
    start_x = sx(cx) - sw(5)
    start_y = sy(cy)
    fb = slide.shapes.build_freeform(start_x, start_y, scale=PX)
    fb.add_line_segments([
        (sx(cx) - sw(1.5), sy(cy) + sh(4)),   # to (-1.5, +4)
        (sx(cx) + sw(5),   sy(cy) - sh(4)),   # to (+5, -4)
    ], close=False)
    tick = fb.convert_to_shape()
    tick.fill.background()
    tick.line.color.rgb = WHITE
    tick.line.width = Pt(2.2)


def _draw_crossing_arrow(slide, sx, sy):
    """Brand-blue arrow from source-glyph right edge (385, 385) to
    receiving-glyph left edge (1055, 335). 3.5pt with a triangular
    arrowhead at the target end.

    Earlier revisions used a custGeom freeform with a hand-written
    <a:cubicBezTo> for a curved arc. PowerPoint-Mac rendered the curve
    inconsistently — sometimes the line vanished entirely behind the
    air-gapped container's dashed border — and the resulting shape was
    ungrabbable for an editor (no anchor handles, no resize-as-arrow).

    Use a native MSO straight connector instead. It renders reliably,
    PowerPoint surfaces the begin/end handles for editing, and the
    visual difference from the SVG's gentle curve is negligible at
    3.5pt over a ~670 px run.
    """
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn
    from lxml import etree

    begin_x, begin_y = sx(385), sy(385)
    end_x,   end_y   = sx(1055), sy(335)

    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        px(begin_x), px(begin_y),
        px(end_x),   px(end_y),
    )
    connector.line.color.rgb = BLUE
    connector.line.width = Pt(3.5)

    # Add a triangular arrowhead at the tail (target end).
    ln = connector.line._get_or_add_ln()
    for old in ln.findall(qn("a:tailEnd")):
        ln.remove(old)
    tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
    tailEnd.set("type", "triangle")
    tailEnd.set("w", "med")
    tailEnd.set("len", "med")

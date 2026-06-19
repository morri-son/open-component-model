#!/usr/bin/env python3
"""
Build OCM-Sovereign-Delivery-Exec.pptx by instantiating slides from the
OCM-Master-Template.potx layouts and filling placeholders.

Why this is structured this way: the .potx template owns visual styling
(palette, type scale, eyebrow + title + footer chrome, column rules, tile
backgrounds). This script only supplies content. To restyle the whole deck,
edit the layouts in build_potx.py and rebuild the .potx — every deck that
uses it picks up the change automatically.

Layout-specific extras (banner image on slide 1, brand row on slide 1 + slide
10, diagram images on content slides) are added inline because they're
deck-specific, not template-wide.

Usage:
    .venv/bin/python build_pptx.py
"""
from __future__ import annotations

import shutil
import subprocess
import sys
import zipfile
from pathlib import Path
from tempfile import TemporaryDirectory

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor

from icon_strokes import STROKE_THIN, STROKE_REGULAR, STROKE_BOLD
from pptx.util import Emu, Pt


# -----------------------------------------------------------------------------
# Paths
# -----------------------------------------------------------------------------

SCRIPT_DIR = Path(__file__).resolve().parent
DECK_DIR = SCRIPT_DIR.parent
MARKETING_DIR = DECK_DIR.parent.parent
ASSETS_DIR = MARKETING_DIR / "assets"
DIAGRAMS_DIR = DECK_DIR / "diagrams"
ICONS_DIR = DIAGRAMS_DIR / "icons"
THEME_DIR = DECK_DIR / "theme"
RASTER_DIR = SCRIPT_DIR / "_raster"

POTX_PATH = DECK_DIR / "OCM-Master.potx"
OUTPUT_PPTX = DECK_DIR / "OCM-Sovereign-Delivery-Exec.pptx"

RASTER_DIR.mkdir(exist_ok=True)


# -----------------------------------------------------------------------------
# Slide geometry — 16:9 @ 1920×1080
# -----------------------------------------------------------------------------

SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
PX = 9525  # 1 px in EMU at 96 dpi

def px(n: float) -> Emu:
    return Emu(int(n * PX))


# -----------------------------------------------------------------------------
# OCM brand palette (canonical, mirror of build_potx.py PALETTE)
# -----------------------------------------------------------------------------

class C:
    BLUE       = RGBColor(0x0F, 0x6B, 0xFF)   # accent1 / brand-blue-dark
    BLUE_MID   = RGBColor(0x0A, 0x3A, 0x99)   # accent2 / brand-blue-mid
    CYAN       = RGBColor(0x5C, 0xD6, 0xFF)   # accent3 / brand-cyan
    GREY_MID   = RGBColor(0x6B, 0x72, 0x80)   # accent4
    BLUE_NIGHT = RGBColor(0x0A, 0x15, 0x30)   # accent5
    GREY_SOFT  = RGBColor(0xF3, 0xF4, 0xF6)   # accent6
    BLACK      = RGBColor(0x00, 0x00, 0x00)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)


# -----------------------------------------------------------------------------
# OOXML namespaces (used for the gradient title surgery)
# -----------------------------------------------------------------------------

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# -----------------------------------------------------------------------------
# SVG → PNG rasterization (for diagrams + icons)
# -----------------------------------------------------------------------------

def rasterize_svg(svg_path: Path, target_w_px: int) -> Path:
    if not svg_path.exists():
        raise FileNotFoundError(svg_path)
    out = RASTER_DIR / (svg_path.stem + f"_{target_w_px}.png")
    if out.exists() and out.stat().st_mtime >= svg_path.stat().st_mtime:
        return out
    subprocess.run(
        ["rsvg-convert", "--width", str(target_w_px), "--keep-aspect-ratio",
         str(svg_path), "-o", str(out)],
        check=True, capture_output=True,
    )
    return out


def rasterize_svg_recolored(svg_path: Path, target_w_px: int,
                            color_hex: str,
                            stroke_width: float | None = None) -> Path:
    """Rasterize a `currentColor`-based SVG (Tabler icon family) with an
    explicit stroke/fill colour.

    Resolution order:

      1.  If a prebuilt PNG exists under `diagrams/icons/prebuilt/`
          matching the requested icon × colour × stroke combination,
          return that path directly. Those files are produced by
          `prebuild_icons.py` and constitute the deck's permanent
          icon-asset library — same bytes the build embeds, same
          name a hand-editor finds in Finder.

      2.  Otherwise fall back to on-the-fly rasterisation under
          build-pptx/_raster/. This path is taken for one-off icons
          that aren't on the prebuild manifest, and produces an
          equivalent PNG.

    rsvg-convert defaults `currentColor` to black; we patch a copy of
    the SVG so the icon paints in the brand palette instead. Cached
    per (file, width, colour, stroke).
    """
    if not svg_path.exists():
        raise FileNotFoundError(svg_path)

    # 1. Prebuilt asset shortcut --------------------------------------
    colour = color_hex.lstrip("#").upper()
    colour_name = {"0F6BFF": "brand-blue", "FFFFFF": "white"}.get(colour)
    if colour_name is not None and stroke_width is not None:
        prebuilt_dir = svg_path.parent / "prebuilt"
        # Stroke tag mirrors prebuild_icons.py: integers keep ".0".
        if float(stroke_width).is_integer():
            stroke_tag = f"{int(stroke_width)}.0"
        else:
            stroke_tag = f"{stroke_width:g}"
        prebuilt_png = prebuilt_dir / f"{svg_path.stem}-stroke-{stroke_tag}-{colour_name}.png"
        if prebuilt_png.exists():
            return prebuilt_png

    # 2. Fallback: rasterise on the fly into _raster/ -----------------
    sw_tag = f"_sw{stroke_width:g}".replace(".", "p") if stroke_width is not None else ""
    out = RASTER_DIR / f"{svg_path.stem}_{target_w_px}_{colour}{sw_tag}.png"
    if out.exists() and out.stat().st_mtime >= svg_path.stat().st_mtime:
        return out
    src = svg_path.read_text(encoding="utf-8")
    # Tabler icons declare stroke="currentColor" on the root <svg>; patching
    # the literal token covers both stroke and fill uses.
    patched = src.replace("currentColor", f"#{colour}")
    if stroke_width is not None:
        # Replace stroke-width="2" (or whatever Tabler ships) with the
        # requested override. We only touch the root-svg attribute — child
        # paths inherit, so this lightens every stroke in the icon.
        import re as _re
        patched = _re.sub(
            r'stroke-width="[^"]*"',
            f'stroke-width="{stroke_width:g}"',
            patched,
        )
    tmp = RASTER_DIR / f"{svg_path.stem}_{colour}{sw_tag}.svg"
    tmp.write_text(patched, encoding="utf-8")
    subprocess.run(
        ["rsvg-convert", "--width", str(target_w_px), "--keep-aspect-ratio",
         str(tmp), "-o", str(out)],
        check=True, capture_output=True,
    )
    return out


def first_existing(*candidates: Path) -> Path | None:
    for c in candidates:
        if c.exists():
            return c
    return None


# -----------------------------------------------------------------------------
# Open .potx as .pptx
# -----------------------------------------------------------------------------

def open_template_as_pptx() -> Presentation:
    """python-pptx refuses .potx files, so transparently swap the
    presentation content type to .pptx-flavored when loading. The output we
    save will still be a valid .pptx (we never write the template type back)."""
    if not POTX_PATH.exists():
        raise FileNotFoundError(
            f"{POTX_PATH} not found — run `python build_potx.py` first."
        )
    tmp_pptx = RASTER_DIR / "_potx_loaded.pptx"
    with zipfile.ZipFile(POTX_PATH, "r") as src:
        data = {n: src.read(n) for n in src.namelist()}
    ct = data["[Content_Types].xml"].decode("utf-8")
    ct = ct.replace(
        "presentationml.template.main+xml",
        "presentationml.presentation.main+xml",
    )
    data["[Content_Types].xml"] = ct.encode("utf-8")
    with zipfile.ZipFile(tmp_pptx, "w", zipfile.ZIP_DEFLATED) as out:
        for name, blob in data.items():
            out.writestr(name, blob)
    return Presentation(str(tmp_pptx))


# -----------------------------------------------------------------------------
# Placeholder helpers
# -----------------------------------------------------------------------------

def find_placeholder(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(f"placeholder idx={idx} not found on slide using "
                   f"layout {slide.slide_layout.name!r}")


def delete_placeholder(slide, idx: int):
    """Remove a placeholder shape entirely from the slide. Use when a layout
    supplies a placeholder you don't want on this particular slide (e.g. the
    body box on the Plain layout for slides that draw their own content)."""
    ph = find_placeholder(slide, idx)
    sp = ph._element
    sp.getparent().remove(sp)


def set_text(slide, idx: int, text: str, *, color: RGBColor | None = None,
             align_left: bool = False):
    """Set a placeholder's text. Layout's lstStyle supplies size/font/case;
    we only override color when needed (e.g. hero title white).

    align_left=True forces left alignment via <a:pPr algn="l"/> — useful on
    placeholders that PowerPoint might otherwise center (Hero title, etc.).
    """
    from pptx.enum.text import PP_ALIGN
    ph = find_placeholder(slide, idx)
    tf = ph.text_frame
    tf.clear()
    paragraphs = text.split("\n")
    for i, line in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if align_left:
            p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        if color is not None:
            run.font.color.rgb = color


def set_split_gradient_title(slide, idx: int, prefix: str, noun: str,
                              align_left: bool = True):
    """Hero title second line: prefix in white, noun with the OCM gradient."""
    from pptx.enum.text import PP_ALIGN
    ph = find_placeholder(slide, idx)
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align_left:
        p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = prefix
    r1.font.color.rgb = C.WHITE
    r2 = p.add_run()
    r2.text = noun
    # Replace solid fill with the OCM gradient on this run.
    rPr = r2._r.get_or_add_rPr()
    for tag in ("solidFill", "gradFill", "noFill"):
        existing = rPr.find(f"{{{A_NS}}}{tag}")
        if existing is not None:
            rPr.remove(existing)
    grad_xml = (
        f'<a:gradFill xmlns:a="{A_NS}" flip="none" rotWithShape="1">'
        '<a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="FFFFFF"/></a:gs>'
        '<a:gs pos="35000"><a:srgbClr val="5CD6FF"/></a:gs>'
        '<a:gs pos="75000"><a:srgbClr val="0F6BFF"/></a:gs>'
        '</a:gsLst>'
        '<a:lin ang="0" scaled="1"/>'
        '</a:gradFill>'
    )
    rPr.insert(0, etree.fromstring(grad_xml))


def set_gradient_title(slide, idx: int, text: str, *, align_left: bool = False):
    """Set a title with the OCM gradient applied to the entire run."""
    from pptx.enum.text import PP_ALIGN
    ph = find_placeholder(slide, idx)
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align_left:
        p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    rPr = r._r.get_or_add_rPr()
    for tag in ("solidFill", "gradFill", "noFill"):
        existing = rPr.find(f"{{{A_NS}}}{tag}")
        if existing is not None:
            rPr.remove(existing)
    grad_xml = (
        f'<a:gradFill xmlns:a="{A_NS}" flip="none" rotWithShape="1">'
        '<a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="FFFFFF"/></a:gs>'
        '<a:gs pos="35000"><a:srgbClr val="5CD6FF"/></a:gs>'
        '<a:gs pos="75000"><a:srgbClr val="0F6BFF"/></a:gs>'
        '</a:gsLst>'
        '<a:lin ang="0" scaled="1"/>'
        '</a:gradFill>'
    )
    rPr.insert(0, etree.fromstring(grad_xml))


def set_action_path_lines(slide, idx: int, lines: list[tuple[str, str]],
                            sep: str = " — "):
    """Render a body placeholder where each line is split into ACTION (blue)
    and PATH (white) by `sep`. Used for the CTA slide."""
    ph = find_placeholder(slide, idx)
    tf = ph.text_frame
    tf.clear()
    for i, (action, path) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run()
        r1.text = action
        r1.font.color.rgb = C.CYAN
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = sep + path
        r2.font.color.rgb = C.WHITE


def set_blue_box_bullets(slide, idx: int, items: list[str]):
    """Render a body placeholder as a list with blue square bullets.

    Each item gets a small filled square (▪) prefix in OCM brand blue, then
    the body text in black. We don't use PowerPoint's <a:buChar> machinery —
    Aptos doesn't render the small filled square consistently — so we put a
    blue-coloured glyph as the first run of each paragraph.
    """
    ph = find_placeholder(slide, idx)
    tf = ph.text_frame
    tf.clear()
    for i, body in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        bullet = p.add_run()
        bullet.text = "▪  "
        bullet.font.color.rgb = C.BLUE
        bullet.font.bold = True
        text_run = p.add_run()
        text_run.text = body
        text_run.font.color.rgb = C.BLACK


# -----------------------------------------------------------------------------
# Static decoration helpers (banner, brand row, diagram embeds)
# -----------------------------------------------------------------------------

def add_banner_full_bleed(slide, image_path: Path):
    """Place an image as full-slide background. Sent to back."""
    if not image_path.exists():
        return
    pic = slide.shapes.add_picture(str(image_path), 0, 0,
                                    width=px(SLIDE_W_PX),
                                    height=px(SLIDE_H_PX))
    # Move to bottom of z-order so placeholders sit on top.
    spTree = pic._element.getparent()
    spTree.remove(pic._element)
    # Insert just after grpSpPr (the first non-element header).
    insert_at = 0
    for i, el in enumerate(spTree):
        if el.tag.endswith("}grpSpPr"):
            insert_at = i + 1
            break
    spTree.insert(insert_at, pic._element)


def add_brand_row(slide):
    """Bottom-of-slide brand row: OCM logo left, NeoNephos logo right (white
    variants for use over dark backgrounds)."""
    ocm_svg = ASSETS_DIR / "ocm" / "ocm-horizontal-white.svg"
    nn_svg = ASSETS_DIR / "neonephos" / "neonephos-foundation-horizontal-white.svg"
    ocm_png = rasterize_svg(ocm_svg, target_w_px=400)
    nn_png = rasterize_svg(nn_svg, target_w_px=400)
    ocm_pic = slide.shapes.add_picture(
        str(ocm_png), px(96), px(SLIDE_H_PX - 56 - 76), height=px(76)
    )
    nn_pic = slide.shapes.add_picture(
        str(nn_png), px(96), px(SLIDE_H_PX - 56 - 52), height=px(52)
    )
    nn_pic.left = px(SLIDE_W_PX - 96) - nn_pic.width


def add_diagram(slide, svg_path: Path | None,
                 x_px: int, y_px: int,
                 max_w_px: int, max_h_px: int):
    """Rasterize a diagram SVG and centre it inside (x, y, max_w, max_h).

    The SVG keeps its intrinsic aspect ratio. If it's wider-than-slot, the
    width fills `max_w_px` and the height is whatever falls out, centred
    vertically in the slot. If it's taller-than-slot, the height fills
    `max_h_px` and the width is centred horizontally. Either way the
    diagram never overflows the slot, and the unused slot half is even
    margins instead of slot-top dead space.

    Also drops the layout's empty picture placeholder (idx=10 on the
    Diagram layout) — otherwise PowerPoint shows a dotted outline +
    'Insert picture' prompt next to our embedded picture.
    """
    if svg_path is None or not svg_path.exists():
        return
    try:
        delete_placeholder(slide, 10)
    except KeyError:
        pass
    png = rasterize_svg(svg_path, target_w_px=max_w_px)
    pic = slide.shapes.add_picture(str(png), px(x_px), px(y_px),
                                    width=px(max_w_px))
    if pic.height > px(max_h_px):
        ratio = px(max_h_px) / pic.height
        pic.height = px(max_h_px)
        pic.width = int(pic.width * ratio)
    # Centre horizontally and vertically inside the slot.
    pic.left = px(x_px) + (px(max_w_px) - pic.width) // 2
    pic.top  = px(y_px) + (px(max_h_px) - pic.height) // 2


def add_tile_icon(slide, tile_x_px: int, tile_y_px: int, icon_name: str):
    """Place a brand-blue icon at the top-left of a tile, sitting in the
    same header row as the (right-side) tile label. 48x48 to balance with
    the 18pt bold label baseline; both share the row from y+24 to y+72."""
    icon_path = ICONS_DIR / icon_name
    if not icon_path.exists():
        return
    png = rasterize_svg_recolored(icon_path, target_w_px=96, color_hex="0F6BFF")
    slide.shapes.add_picture(
        str(png),
        px(tile_x_px + 24), px(tile_y_px + 24),
        width=px(48), height=px(48),
    )


def _crop_to_content(png_path: Path) -> Path:
    """Crop trailing transparent / white margins off a logo PNG so that
    `add_logo_row` can normalise on *content* height (not file height).

    Without this step, logos that bake whitespace into the file (e.g. the
    SAP-NS2 PNG, BwI's SVG with .de wordmark padding) render visibly smaller
    next to logos that fill their bbox edge-to-edge (SAP, Platform Mesh).

    Crops are always written into the local `_raster` cache (never back into
    the assets tree, even when the input lives there).
    """
    from PIL import Image
    out = RASTER_DIR / (png_path.stem + "_crop.png")
    if out.exists() and out.stat().st_mtime >= png_path.stat().st_mtime:
        return out
    im = Image.open(png_path).convert("RGBA")
    w, h = im.size
    px_data = im.load()
    left, top, right, bot = w, h, 0, 0
    found = False
    for y in range(h):
        for x in range(w):
            r, g, b, a = px_data[x, y]
            if a > 8 and not (r > 240 and g > 240 and b > 240):
                found = True
                if x < left: left = x
                if x > right: right = x
                if y < top: top = y
                if y > bot: bot = y
    if not found:
        # Defensive fallback: nothing to crop, return original.
        return png_path
    im.crop((left, top, right + 1, bot + 1)).save(out)
    return out


def add_logo_row(slide, logos: list, y_px: int,
                  row_h_px: int = 120,
                  max_logo_w_px: int = 320, max_logo_h_px: int = 80,
                  caption_pt: int = 14):
    """Three (or more) logos centred in a row, normalised on visible content height.

    `logos` accepts entries as:
      - Path                 — no link, no caption
      - (Path, url)          — clickable, no caption
      - (Path, url, caption) — clickable + caption text rendered below

    Captions are useful for icon-only marks (Kyma, OpenControlPlane) where
    the logo alone doesn't carry the project name. Tuples make the picture
    shape clickable in PowerPoint and in PDF export.

    Each input logo is rasterised (or used directly for PNG) and then cropped
    to its visible bounding box, so wordmarks with baked-in whitespace
    (BwI's tall ".de" frame, the SAP-NS2 PNG, Gardener's tagline) render at
    the same optical height as logos that already fill their bbox (SAP,
    Platform Mesh). Final placement uses height-first normalisation with a
    width cap as the safety net for very wide logos.
    """
    from pptx.enum.text import PP_ALIGN
    margin_x = 160
    inner_w = SLIDE_W_PX - 2 * margin_x
    n = len(logos)
    slot_w = inner_w // n
    for i, entry in enumerate(logos):
        caption = None
        if isinstance(entry, tuple):
            if len(entry) == 3:
                path, url, caption = entry
            else:
                path, url = entry
        else:
            path, url = entry, None
        if path is None or not path.exists():
            continue
        if path.suffix.lower() == ".svg":
            img = rasterize_svg(path, target_w_px=max_logo_w_px * 2)
        else:
            img = path
        # Crop to visible content so we normalise on logo height, not file
        # height. Always go through this — even SVG output can carry padding.
        img = _crop_to_content(img)
        slot_x = margin_x + i * slot_w
        # Add at native size, then scale to fit the height constraint first.
        pic = slide.shapes.add_picture(str(img), px(slot_x), px(y_px))
        if pic.height != px(max_logo_h_px):
            ratio = px(max_logo_h_px) / pic.height
            pic.height = px(max_logo_h_px)
            pic.width = int(pic.width * ratio)
        # Then enforce width cap (rare — only triggers for very wide logos).
        if pic.width > px(max_logo_w_px):
            ratio = px(max_logo_w_px) / pic.width
            pic.width = px(max_logo_w_px)
            pic.height = int(pic.height * ratio)
        pic.left = px(slot_x) + (px(slot_w) - pic.width) // 2
        pic.top = px(y_px) + (px(row_h_px) - pic.height) // 2
        if url:
            pic.click_action.hyperlink.address = url
        if caption:
            cap_y = y_px + row_h_px + 6
            tb = slide.shapes.add_textbox(px(slot_x), px(cap_y),
                                           px(slot_w), px(28))
            tf = tb.text_frame
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = caption
            r.font.name = "Aptos"
            r.font.size = Pt(caption_pt)
            r.font.color.rgb = C.GREY_MID
            if url:
                r.hyperlink.address = url


# -----------------------------------------------------------------------------
# Hero / Tile geometry helpers — must mirror what the .potx layouts use
# -----------------------------------------------------------------------------

# These constants drive where the on-slide decoration (banner, brand row,
# tile icons) is placed. They MUST match the corresponding layout positions
# in build_potx.py — if you change one place, change the other.

TILE_X0_PX = 120
TILE_Y0_PX = 520
TILE_W_PX = 544
TILE_H_PX = 230
TILE_GUTTER_PX = 24


def tile_origin(index: int) -> tuple[int, int]:
    col = index % 3
    row = index // 3
    return (TILE_X0_PX + col * (TILE_W_PX + TILE_GUTTER_PX),
            TILE_Y0_PX + row * (TILE_H_PX + TILE_GUTTER_PX))


# -----------------------------------------------------------------------------
# Layout lookup
# -----------------------------------------------------------------------------

def layouts_by_name(prs: Presentation) -> dict[str, object]:
    return {l.name: l for l in prs.slide_masters[0].slide_layouts}


# =============================================================================
# Build the deck
# =============================================================================

def build():
    prs = open_template_as_pptx()
    layouts = layouts_by_name(prs)

    expected = {"Hero", "CTA", "Content / 3-Column",
                "Content / Diagram", "Content / Tiles", "Content / 2-Column",
                "Section Divider", "Plain", "Plain / Compact"}
    missing = expected - set(layouts)
    if missing:
        sys.exit(f"template missing expected layouts: {missing}")

    # ---- SLIDE 1 — HERO (cold-room canonical, revised 2026-06-17) -----------
    # Stake-led title + spannungs-subtitle. The original 11-word "Three minutes
    # from now, you'll know what your supply chain doesn't" was decomposed:
    # the *stake* lands in the title, the *time-bound promise* lands in the
    # subtitle. Title fits cleanly on ONE line at 115pt; the gradient sits on
    # a short second-line noun for visual punch.
    s = prs.slides.add_slide(layouts["Hero"])
    add_banner_full_bleed(s, THEME_DIR / "OCM-Banner.png")
    set_text(s, 1, "Your supply chain has", color=C.WHITE)
    set_split_gradient_title(s, 2, prefix="", noun="blind spots.")
    set_text(s, 3,
             "Three minutes from now, you'll know what they are.",
             color=C.CYAN)
    set_text(s, 4,
             "Open Component Model — open source, NeoNephos Foundation.",
             color=C.WHITE)
    add_brand_row(s)

    # ---- SLIDE 2 — WHY NOW (V1, sovereignty-led) ----------------------------
    s = prs.slides.add_slide(layouts["Content / 3-Column"])
    set_text(s, 1, "WHY NOW")
    set_text(s, 2, "Sovereignty is no longer optional")
    set_text(s, 10, "SOVEREIGNTY PRESSURE")
    set_text(s, 11, "Wherever the law puts the boundary — by jurisdiction, "
                     "sector, or air-gap — software must be deliverable, "
                     "verifiable, and operable inside it.")
    set_text(s, 12, "REGULATION TIGHTENING")
    set_text(s, 13, "EU DORA · NIS2 · CRA. Provable supply-chain control, "
                     "not best effort.")
    set_text(s, 14, "SUPPLY-CHAIN ATTACKS ARE REAL")
    set_text(s, 15, "SolarWinds. xz. log4shell. Signatures must survive the "
                     "journey, or compliance is theatre.")

    # ---- SLIDE 3 — MEET OCM (hub-and-spoke diagram, Option 3 reframe) -------
    # Diagram positioned per user spec 2026-06-17: 50.02 × 15.93 cm,
    # x=-2.4cm (slight bleed left), y=11.65cm.
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "THE ANSWER")
    set_text(s, 2, "Meet OCM. One identity, every boundary.")
    add_diagram(s, DIAGRAMS_DIR / "03-meet-ocm-hub-and-spoke.svg",
                 x_px=60, y_px=240, max_w_px=1800, max_h_px=780)

    # ---- SLIDE 3' — SAME, BUT NATIVE PPT SHAPES -----------------------------
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "THE ANSWER  (NATIVE)")
    set_text(s, 2, "Meet OCM. One identity, every boundary.")
    delete_placeholder(s, 10)
    from slide_3_native import add_hub_and_spoke_native_diagram
    add_hub_and_spoke_native_diagram(s, x=60, y=240, w=1800, h=780,
                                      icons_dir=ICONS_DIR,
                                      rasterize_recolored=rasterize_svg_recolored)

    # ---- SLIDE 4a — THE SHIFT, SBOD (text-only) -----------------------------
    s = prs.slides.add_slide(layouts["Plain / Compact"])
    set_text(s, 1, "THE SHIFT")
    set_text(s, 2, "SBOM lists. SBOD delivers.")
    set_blue_box_bullets(s, 10, [
        "An SBOM tells you what's in your software. It was built for inventory.",
        "A Software Bill of Delivery (SBOD) tells you what you delivered, "
        "how to verify, transport, and operate it. "
        "It was built for delivery.",
        "The SBOD contains the SBOM. OCM doesn't replace your SBOM tooling — "
        "it gives the SBOM an envelope that's compliance-native, signed once, "
        "and travels intact across any boundary.",
    ])

    # ---- SLIDE 4b — THE SHIFT (diagram only, ORIGINAL SVG) -------------------
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "THE SHIFT — SBOM INSIDE SBOD")
    set_text(s, 2, "SBOM lists. SBOD delivers.")
    diagram = first_existing(
        DIAGRAMS_DIR / "04-sbom-inside-sbod.svg",
        DIAGRAMS_DIR / "04-sbom-vs-sbod.svg",
    )
    if diagram:
        add_diagram(s, diagram, x_px=60, y_px=240, max_w_px=1800, max_h_px=780)

    # ---- SLIDE 4b' — SAME, BUT NATIVE PPT SHAPES ----------------------------
    # Side-by-side review version. The native one reframes content (SBOM is
    # one of five elements, not the centrepiece) AND swaps the rendering
    # technique (autoshapes + textboxes + recoloured icons instead of a
    # rasterised SVG). User wants both visible to compare.
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "THE SHIFT — SBOM INSIDE SBOD  (NATIVE A)")
    set_text(s, 2, "SBOM lists. SBOD delivers.")
    delete_placeholder(s, 10)
    from slide_4b_native import add_sbod_native_diagram
    add_sbod_native_diagram(s, x=60, y=240, w=1800, h=780,
                             icons_dir=ICONS_DIR,
                             rasterize_recolored=rasterize_svg_recolored,
                             icon_stroke=STROKE_THIN)

    # ---- SLIDE 4b'' — NATIVE VARIANT B (vertical artifact list + brace) -----
    # Parallel native variant of the SBOM-inside-SBOD slide. Lifts the identity
    # to a header above the artifact list, drops SBOM to one row of five, and
    # uses a curly brace + lock to show that "one digest covers all". Lives
    # alongside Variant A so the user can review them side-by-side.
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "THE SHIFT — SBOM INSIDE SBOD  (NATIVE B)")
    set_text(s, 2, "SBOM lists. SBOD delivers.")
    delete_placeholder(s, 10)
    from slide_4b_native_v2 import add_sbom_inside_sbod_native_v2
    add_sbom_inside_sbod_native_v2(s, x=60, y=240, w=1800, h=780,
                                    icons_dir=ICONS_DIR,
                                    rasterize_recolored=rasterize_svg_recolored,
                                    icon_stroke=STROKE_THIN)

    # ---- SLIDE 5 — HOW OCM COMPOSES (NEW, comparator slide) ----------------
    # Disarms three "we already have this" objections on one slide:
    # signing, transport, compliance. Each column says "what you have today"
    # then "what OCM adds" — OCM doesn't replace, it composes around them.
    s = prs.slides.add_slide(layouts["Content / 3-Column"])
    set_text(s, 1, "HOW OCM COMPOSES")
    set_text(s, 2, "Composes around your existing stack.")
    set_text(s, 10, "SIGNING")
    set_text(s, 11, "Your tools sign artifacts.\n"
                     "OCM signs the whole release — one signature, every digest.")
    set_text(s, 12, "TRANSPORT")
    set_text(s, 13, "Registries differ by type and location.\n"
                     "OCM moves the release across them all.")
    set_text(s, 14, "COMPLIANCE")
    set_text(s, 15, "Your scanners see one artifact at a time.\n"
                     "OCM correlates findings to the release. "
                     "Compliance becomes continuous.")

    # ---- SLIDE 6 — OCM IN ONE PICTURE (was slide 5) -------------------------
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "OCM IN ONE PICTURE")
    set_text(s, 2, "Pack · Sign · Transport · Deploy")
    diagram6 = first_existing(
        DIAGRAMS_DIR / "05-pack-sign-transport-deploy-v2.svg",
        DIAGRAMS_DIR / "05-pack-sign-transport-deploy.svg",
    )
    add_diagram(s, diagram6, x_px=60, y_px=240, max_w_px=1800, max_h_px=780)

    # ---- SLIDE 6' — SAME, BUT NATIVE PPT SHAPES -----------------------------
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "OCM IN ONE PICTURE  (NATIVE)")
    set_text(s, 2, "Pack · Sign · Transport · Deploy")
    delete_placeholder(s, 10)
    from slide_6_native import add_pack_sign_transport_deploy_native
    add_pack_sign_transport_deploy_native(s, x=60, y=240, w=1800, h=780,
                                           icons_dir=ICONS_DIR,
                                           rasterize_recolored=rasterize_svg_recolored,
                                           icon_stroke=STROKE_THIN)

    # ---- SLIDE 7a — SOVEREIGN-READY (text-only, was 6a) --------------------
    s = prs.slides.add_slide(layouts["Plain / Compact"])
    set_text(s, 1, "SOVEREIGN-READY")
    set_text(s, 2, "Trust, but verify.")
    set_blue_box_bullets(s, 10, [
        "Identity is location-independent. A component carries its name "
        "regardless of which registry it lives in.",
        "Signatures are location-independent. Sign once at source; verify at "
        "the destination, or at any hop in between, with no callback upstream.",
        "Day-2 ops happen inside the boundary. Subscribe to the component and "
        "pull upgrades on your schedule, scale across regions, all without "
        "reaching back upstream.",
        "On transfer into a sovereign environment, a component can carry every "
        "artifact it needs along with it. The destination needs nothing more.",
    ])

    # ---- SLIDE 7b — SOVEREIGN-READY (diagram only, was 6b) -----------------
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "SOVEREIGN-READY — AIR-GAP")
    set_text(s, 2, "Trust travels with the component.")
    # Diagram fills the standard diagram slot (x=60 y=240, 1800×780).
    add_diagram(s, DIAGRAMS_DIR / "06-sovereign-airgap.svg",
                 x_px=60, y_px=240, max_w_px=1800, max_h_px=780)

    # ---- SLIDE 7b' — SAME, BUT NATIVE PPT SHAPES ----------------------------
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "SOVEREIGN-READY — AIR-GAP  (NATIVE)")
    set_text(s, 2, "Trust travels with the component.")
    delete_placeholder(s, 10)
    from slide_7b_native import add_sovereign_airgap_native
    add_sovereign_airgap_native(s, x=60, y=240, w=1800, h=780,
                                 icons_dir=ICONS_DIR,
                                 rasterize_recolored=rasterize_svg_recolored,
                                 icon_stroke=STROKE_THIN)

    # ---- SLIDE 8 — SCAN / Compliance-native (was 7) ------------------------
    s = prs.slides.add_slide(layouts["Plain"])
    set_text(s, 1, "SCAN — COMPLIANCE-NATIVE WITH OPEN DELIVERY GEAR")
    set_text(s, 2, "Compliance as a system property —\nnot a quarterly retrofit.")
    set_blue_box_bullets(s, 10, [
        "Open Delivery Gear (ODG) is the OCM compliance automation engine.",
        "The Compliance Dashboard is your entry point: every component, "
        "every finding, every signature in one view.",
        "Continuous scans run asynchronously — even after release.",
        "Findings get rescored against contextual risk, so your team patches "
        "what actually matters.",
        "Every compliance signal correlates by component identity. Auditors "
        "get evidence, not spreadsheets.",
    ])

    # ---- SLIDE 9 — WHAT OCM UNLOCKS (tiles, was 8) -------------------------
    s = prs.slides.add_slide(layouts["Content / Tiles"])
    set_text(s, 1, "WHAT OCM UNLOCKS")
    set_text(s, 2, "One model unlocks all of this.")
    tiles = [
        ("lock.svg", "Code signing across stacks",
         "Sign once at source; verify everywhere, with no per-stack tooling."),
        ("package-export.svg", "Air-gapped delivery",
         "Walk a complete component across an air gap; verify at destination."),
        ("rocket.svg", "Kubernetes-native deployment",
         "OCM controllers deploy components directly into clusters."),
        ("radar.svg", "Asynchronous security scans",
         "Continuous scanning, even after release; findings tied to component identity."),
        ("source-of-truth.svg", "One source of truth",
         "Rebuild any landscape from a single signed descriptor."),
        ("report-analytics.svg", "Automated compliance reporting",
         "Reports composed from SBOD metadata — no spreadsheet drift."),
    ]
    for i, (icon, label, body) in enumerate(tiles):
        # Tile placeholders alternate label/body: idx 20+2i = label, 21+2i = body
        set_text(s, 20 + i * 2, label)
        set_text(s, 21 + i * 2, body)
        x, y = tile_origin(i)
        add_tile_icon(s, x, y, icon)

    # ---- SLIDE 10 — Adopters (two-column logo wall, was 9) ------------------
    # Plain layout + manual logo rows (not enough placeholders for a logo
    # wall; cleaner to draw it inline than add a new layout to the template).
    s = prs.slides.add_slide(layouts["Plain"])
    set_text(s, 1, "TRUSTED IN PRODUCTION")
    set_text(s, 2, "Aligned with NeoNephos.")
    # Plain layout has a body placeholder at idx=10. We want neither text nor
    # an empty container; the logo rows below replace it. Section label
    # ("ADOPTED BY ENTERPRISES…") removed: the eyebrow + title already frame
    # the slide; an extra label on top of the logos read as redundant chrome.
    # Reclaimed vertical space goes to larger logos and a higher first row.
    delete_placeholder(s, 10)
    add_logo_row(s, [
        (ASSETS_DIR / "adopters" / "neonephos" / "neonephos-foundation-horizontal-color.svg",
         "https://neonephos.org", "NeoNephos"),
        (ASSETS_DIR / "adopters" / "sap" / "sap-horizontal-color.svg",
         "https://www.sap.com", "SAP"),
        (ASSETS_DIR / "adopters" / "bwi" / "bwi-horizontal-color.svg",
         "https://www.bwi.de", "BWI"),
        (ASSETS_DIR / "adopters" / "sap-ns2" / "sap-ns2-getlogovector.png",
         "https://sapns2.com", "SAP NS2"),
    ], y_px=550, max_logo_w_px=320, max_logo_h_px=96, caption_pt=20)
    add_logo_row(s, [
        (ASSETS_DIR / "adopters" / "gardener" / "gardener-horizontal-color.svg",
         "https://gardener.cloud", "Gardener"),
        (ASSETS_DIR / "adopters" / "konfidence" / "konfidence-horizontal-light.svg",
         "https://konfidence.cloud", "Konfidence"),
        (ASSETS_DIR / "adopters" / "open-control-plane" / "opencontrolplane-icon-color.svg",
         "https://open-control-plane.io", "OpenControlPlane"),
        (ASSETS_DIR / "adopters" / "platform-mesh" / "platform-mesh-horizontal-color.svg",
         "https://platform-mesh.io", "Platform Mesh"),
    ], y_px=790, max_logo_w_px=320, max_logo_h_px=96, caption_pt=20)

    # ---- SLIDE 11 — CTA (was 10) --------------------------------------------
    s = prs.slides.add_slide(layouts["CTA"])
    set_text(s, 1, "Start delivering with confidence.", color=C.WHITE)
    set_action_path_lines(s, 2, [
        ("Try it",        "ocm.software"),
        ("Build with us", "github.com/open-component-model"),
        ("Talk to us",    "community channels on the website"),
    ])
    add_brand_row(s)

    # ---- HIDDEN — Trademark & licensing notice -----------------------------
    # Marked show="0" so it doesn't appear in slideshow mode but is visible
    # when editing and survives PDF export. Lives at the back of the deck for
    # legal/audit completeness; the speaker need never present it.
    add_appendix_glossary_slide(prs, layouts)
    add_hidden_trademark_slide(prs, layouts)

    prs.save(str(OUTPUT_PPTX))
    print(f"Wrote {OUTPUT_PPTX}")


# -----------------------------------------------------------------------------
# Inline-decoration helpers used by slide 9
# -----------------------------------------------------------------------------

def add_label_at(slide, y_px: int, text: str):
    """Brand-blue ALL-CAPS section label — slide 9 logo wall headers."""
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(px(120), px(y_px),
                                   px(SLIDE_W_PX - 240), px(36))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = "Aptos"
    f.size = Pt(18)
    f.bold = True
    f.color.rgb = C.BLUE
    rPr = r._r.get_or_add_rPr()
    rPr.set("cap", "all")
    rPr.set("spc", "110")


def add_centred_proof(slide, y_px: int, text: str):
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(px(120), px(y_px),
                                   px(SLIDE_W_PX - 240), px(80))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.4
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(18)
    r.font.color.rgb = C.BLUE_MID


def add_source_line(slide, y_px: int, text: str):
    """Small grey source/reference line — for slide footers that name the
    canonical URL of a project (NeoNephos, ODG, OpenControlPlane, etc.).
    Renders at the bottom of the slide as a discreet citation, not a CTA."""
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(px(120), px(y_px),
                                   px(SLIDE_W_PX - 240), px(40))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(13)
    r.font.color.rgb = C.GREY_MID


# Glossary on the appendix slide. Term → expansion. Two-column layout, term
# in brand-blue bold, expansion in black on the same line. Sized so the full
# list fits within the slide body area without overflow. Entries sorted
# alphabetically (case-insensitive) so the audience can scan them as a
# reference, not a curated narrative — the order of introduction across the
# deck is incidental once the slide is in front of them.
GLOSSARY_ENTRIES: list[tuple[str, str]] = [
    ("BSI C5",    "Bundesamt für Sicherheit in der Informationstechnik — Cloud Computing Compliance Criteria Catalogue."),
    ("BTP",       "SAP Business Technology Platform."),
    ("CRA",       "Cyber Resilience Act — EU regulation on cybersecurity for products with digital elements."),
    ("DORA",      "Digital Operational Resilience Act — EU regulation for ICT risk in financial services."),
    ("FedRAMP",   "Federal Risk and Authorization Management Program — US standardised cloud security assessment."),
    ("FISMA",     "Federal Information Security Modernization Act — US federal information security mandate."),
    ("Grype",     "Open-source vulnerability scanner for container images and filesystems (Anchore)."),
    ("Helm",      "Package manager for Kubernetes; reference artifact type for OCM."),
    ("LoB",       "Line of Business — SAP organisational unit owning a product portfolio."),
    ("NeoNephos", "European foundation for sovereign cloud open-source projects, hosted under the Linux Foundation."),
    ("NIS2",      "Network and Information Security Directive 2 — EU baseline for cybersecurity of essential entities."),
    ("OCI",       "Open Container Initiative — open standards for container image format and distribution."),
    ("OCM",       "Open Component Model — vendor-neutral specification for signed, transportable software components."),
    ("ODG",       "Open Delivery Gear — OCM-native compliance automation engine and dashboard."),
    ("OSS",       "Open Source Software."),
    ("PKI",       "Public Key Infrastructure — framework for managing certificates and signing keys."),
    ("SBOD",      "Software Bill of Delivery — the OCM component descriptor, signed and traceable. Containing all artifacts and metadata for delivery and deployment."),
    ("SBOM",      "Software Bill of Materials — inventory of components and dependencies inside a software artifact."),
    ("SecNumCloud", "French cloud security qualification scheme operated by ANSSI."),
    ("Sigstore",  "Open-source project for keyless software signing using OIDC identities."),
    ("SPDX",      "Software Package Data Exchange — ISO/IEC 5962 standard format for SBOM data."),
    ("SWID",      "Software Identification Tags — ISO/IEC 19770-2 standard for software inventory."),
    ("Trivy",     "Open-source security scanner for containers, IaC, and code (Aqua Security)."),
]


def add_glossary_grid(slide, entries: list[tuple[str, str]],
                      y_top_px: int = 360):
    """Two-column grid of term / definition pairs.

    Term in brand-blue bold, definition in black on the same line, separated
    by an em-dash. Two columns side-by-side, entries flowing top-to-bottom
    in column 1 then column 2. Sized for ~22 entries on a 1920x1080 slide.
    """
    n = len(entries)
    per_col = (n + 1) // 2
    col_gap_px = 60
    col_w_px = (SLIDE_W_PX - 240 - col_gap_px) // 2
    row_h_px = 40
    col_h_px = per_col * row_h_px

    for col in range(2):
        col_x = 120 + col * (col_w_px + col_gap_px)
        col_entries = entries[col * per_col : (col + 1) * per_col]
        if not col_entries:
            continue
        tb = slide.shapes.add_textbox(px(col_x), px(y_top_px),
                                       px(col_w_px), px(col_h_px))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.word_wrap = True
        for i, (term, definition) in enumerate(col_entries):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(2)
            p.space_after = Pt(6)
            r1 = p.add_run()
            r1.text = term
            r1.font.name = "Aptos"
            r1.font.size = Pt(14)
            r1.font.bold = True
            r1.font.color.rgb = C.BLUE
            r2 = p.add_run()
            r2.text = f"  —  {definition}"
            r2.font.name = "Aptos"
            r2.font.size = Pt(13)
            r2.font.color.rgb = C.BLACK


def add_appendix_glossary_slide(prs, layouts):
    """Last visible slide: abbreviations and acronyms used in the deck.

    Plain / Compact layout (1-line title) with eyebrow + title in the master,
    body placeholder deleted in favour of a custom two-column term/definition
    grid that starts at the same y-coordinate the Compact layout reserves for
    its body placeholder (y=520) so the grid never bleeds into the title.
    Lives after the CTA slide so it's reachable as a reference but not part
    of the main flow."""
    s = prs.slides.add_slide(layouts["Plain / Compact"])
    set_text(s, 1, "APPENDIX — ABBREVIATIONS")
    set_text(s, 2, "Quick reference for the acronyms used in this deck.")
    delete_placeholder(s, 10)
    add_glossary_grid(s, GLOSSARY_ENTRIES, y_top_px=520)


def add_hidden_trademark_slide(prs, layouts):
    """Append two non-presented slides carrying trademark and licensing
    notices for every third-party logo on the deck. Marked show="0" so they
    are skipped in slideshow mode but visible to editors and survive PDF
    export — keeps the legal acknowledgement attached to the deck without
    burdening the speaker. Two slides used because a single Plain layout
    body bleeds into the footer at this bullet count.
    Canonical record: assets/adopters/LICENSING.md."""
    # ---- Hidden 1/2 — adopter logos -----------------------------------------
    s = prs.slides.add_slide(layouts["Plain"])
    s.element.set("show", "0")
    set_text(s, 1, "TRADEMARK & LICENSE NOTICES (1/2)")
    set_text(s, 2, "Logos and trademarks belong to their respective owners.")
    set_blue_box_bullets(s, 10, [
        "SAP, SAP NS2 — trademarks of SAP SE / SAP National Security "
        "Services. Editorial use only; no endorsement implied. "
        "sap.com · sapns2.com",
        "BWI — trademark of BWI GmbH (Bundeswehr-IT). Editorial use of the "
        "Wikimedia public-domain wordmark; verify against BWI press "
        "conditions before external publication. bwi.de",
        "Gardener, Platform Mesh, NeoNephos Foundation — Linux Foundation "
        "Europe artwork; usage governed by the Linux Foundation trademark "
        "usage guidelines (linuxfoundation.org/legal/trademark-usage). "
        "gardener.cloud · platform-mesh.io · neonephos.org",
        "Konfidence — SAP-supported open project; logo from konfidence.cloud. "
        "Editorial use only; verify with the Konfidence project before "
        "external publication. konfidence.cloud",
        "OpenControlPlane — open-source project at open-control-plane.io. "
        "Editorial use only; verify with the project before external "
        "publication. open-control-plane.io",
    ])

    # ---- Hidden 2/2 — remaining marks + sourcing pointer --------------------
    s = prs.slides.add_slide(layouts["Plain"])
    s.element.set("show", "0")
    set_text(s, 1, "TRADEMARK & LICENSE NOTICES (2/2)")
    set_text(s, 2, "Third-party trademarks named for technical reference.")
    set_blue_box_bullets(s, 10, [
        "Kyma — SAP-originated open-source project at kyma-project.io. "
        "Editorial use only. kyma-project.io",
        "Trivy, Grype, Sigstore, Helm, OCI, Kubernetes, kro, Flux, Argo CD — "
        "third-party trademarks named for technical reference; ownership "
        "remains with their respective projects and organisations.",
    ])


# =============================================================================
if __name__ == "__main__":
    if not shutil.which("rsvg-convert"):
        sys.exit("rsvg-convert not found; install via `brew install librsvg`")
    build()

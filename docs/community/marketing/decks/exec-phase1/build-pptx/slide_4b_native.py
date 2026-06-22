"""
slide_4b_native.py — Slide 4b "An envelope, not a list." rendered as native
PowerPoint shapes (rounded rectangles, textboxes, recoloured icons), not
a rasterized SVG.

Why native: SVG rasterization throws away typography, links, and
edit-in-PowerPoint affordances. For diagrams that are mostly boxes + text
(like this one), drawing them with python-pptx primitives keeps the slide
hand-editable and crisp at any zoom.

Composition:

  Outer SBOD container (brand-blue stroke, white fill)
    ┌─ Header line  "SOFTWARE BILL OF DELIVERY (SBOD)"
    └─ Identity     "github.com/acme/webshop · v1.0.0"

  5 inner boxes in 2-2-1 grid, each carrying:
    icon (top-left, brand-blue) + headline + 1-2 line subtext
      Row 1: SBOM           |  Signature
      Row 2: Artifacts      |  Location-independent identity
      Row 3: Day-1 + Day-2 deployment (full-width)

The five boxes share the same visual weight — SBOM is one of five, not
the centre piece, which is the point of the slide.

Usage:
    from slide_4b_native import add_sbod_native_diagram
    add_sbod_native_diagram(slide, x=60, y=240, w=1800, h=780,
                             icons_dir=ICONS_DIR,
                             rasterize_recolored=rasterize_svg_recolored)
"""
from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


PX = 9525


def px(n: float) -> Emu:
    return Emu(int(n * PX))


# OCM brand palette (mirror of build_pptx C class — kept local so this module
# is self-contained and importable by both deck builders).
BLUE       = RGBColor(0x0F, 0x6B, 0xFF)   # brand-blue
BLUE_MID   = RGBColor(0x0A, 0x3A, 0x99)   # mid blue used for headlines
GREY_MID   = RGBColor(0x6B, 0x72, 0x80)   # secondary text
BLACK      = RGBColor(0x00, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)


def _styled_run(p, text, *, size_pt, bold=False, color=BLACK,
                font="Aptos", italic=False, all_caps=False,
                letter_spacing=None):
    """Add a run with the OCM type system applied."""
    from lxml import etree
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size_pt)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if all_caps or letter_spacing is not None:
        rPr = r._r.get_or_add_rPr()
        if all_caps:
            rPr.set("cap", "all")
        if letter_spacing is not None:
            rPr.set("spc", str(letter_spacing))
    return r


def _add_textbox(slide, x, y, w, h, *, anchor="t", margins=0):
    """Borderless textbox at the given coords, with consistent inset margins."""
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = px(margins)
    tf.margin_top = tf.margin_bottom = px(margins)
    tf.word_wrap = True
    if anchor == "ctr":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "b":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    return tf


def _add_rounded_rect(slide, x, y, w, h, *,
                      fill=WHITE, stroke=BLUE, stroke_pt=1.5,
                      corner_pct=0.04):
    """Rounded rectangle with explicit fill, stroke, and corner radius."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h),
    )
    # corner_pct: 0..0.5 of min(w,h). PowerPoint stores as % * 100000.
    shape.adjustments[0] = corner_pct
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = stroke
        shape.line.width = Pt(stroke_pt)
    # Don't put any default text in the autoshape; we add textboxes on top.
    return shape


# Per-box content. Each entry: (icon_filename, headline, subtext_line1, subtext_line2_optional)
BOXES = [
    ("list-search.svg",
     "SBOM",
     "What's inside the software —",
     "the bill of materials lives here."),
    ("signature.svg",
     "SIGNATURE",
     "One signature covers",
     "every artifact by digest."),
    ("package-export.svg",
     "ARTIFACTS",
     "What gets deployed:",
     "manifests, config, images, charts."),
    ("world-pin.svg",
     "LOCATION-INDEPENDENT IDENTITY",
     "Same component, any registry —",
     "identity and signature travel together."),
    ("rocket.svg",
     "DAY-1 + DAY-2 DEPLOYMENT",
     "Install. Upgrade. Reconfigure.",
     "All from one signed source — at the destination, on its own."),
]


from icon_strokes import STROKE_THIN


def add_sbod_native_diagram(slide, *, x=60, y=240, w=1800, h=780,
                             icons_dir: Path,
                             rasterize_recolored,
                             icon_stroke=STROKE_THIN):
    """Draw the SBOD anatomy diagram natively into the given slot.

    icon_stroke selects the rasterised Tabler-icon weight via
    icon_strokes.STROKE_THIN / STROKE_REGULAR / STROKE_BOLD. Default
    is STROKE_THIN (1.0) — matches the deck's SVG-variant slides.

    Parameters:
        slide: python-pptx Slide object.
        x, y, w, h: diagram slot in slide coordinates (px @ 96 dpi).
        icons_dir: folder containing the Tabler-style SVGs used as box icons.
        rasterize_recolored: callable (svg_path, target_w_px, color_hex) -> Path
            that returns a rasterized PNG of the SVG in the requested colour.
            Passed in (rather than imported) so this module stays free of
            build-script-specific dependencies.

    Layout inside the slot:
        - Outer container: 20px inset on all sides.
        - Header zone:     top 110px of the inner area (title + identity).
        - 5 boxes:         remaining height, split into 2 + 2 + 1 rows with
                            18px gutters between rows and 18px between cols.
    """
    # --- Outer SBOD container ----------------------------------------------
    pad = 20
    outer_x, outer_y = x + pad, y + pad
    outer_w, outer_h = w - 2 * pad, h - 2 * pad
    _add_rounded_rect(slide, outer_x, outer_y, outer_w, outer_h,
                      stroke=BLUE, stroke_pt=4.0, corner_pct=0.025)

    # --- Header text -------------------------------------------------------
    # Title: "SOFTWARE BILL OF DELIVERY (SBOD)" — brand-blue, ALL-CAPS.
    # Identity: "github.com/acme/webshop · v1.0.0" — grey, slightly smaller.
    header_pad_x, header_pad_y = 36, 24
    tf = _add_textbox(slide,
                       outer_x + header_pad_x, outer_y + header_pad_y,
                       outer_w - 2 * header_pad_x, 50)
    p = tf.paragraphs[0]
    _styled_run(p, "SOFTWARE BILL OF DELIVERY (SBOD)",
                size_pt=25, bold=True, color=BLUE,
                all_caps=True, letter_spacing="120")

    tf = _add_textbox(slide,
                       outer_x + header_pad_x, outer_y + header_pad_y + 42,
                       outer_w - 2 * header_pad_x, 32)
    p = tf.paragraphs[0]
    _styled_run(p, "github.com/acme/webshop  ·  v1.0.0",
                size_pt=14, color=GREY_MID, font="Consolas")

    # --- 5 boxes in 2-2-1 grid --------------------------------------------
    box_pad_x, box_pad_y_top = 36, 110
    grid_x = outer_x + box_pad_x
    grid_y = outer_y + box_pad_y_top
    grid_w = outer_w - 2 * box_pad_x
    grid_h = outer_h - box_pad_y_top - 24  # 24px bottom margin

    col_gutter = 18
    row_gutter = 18
    # Two rows of pair-boxes + one full-width row at the bottom. Allocate the
    # vertical space proportionally so all three rows are roughly equal-weight.
    pair_row_h = (grid_h - 2 * row_gutter) * 0.32  # each of rows 1+2
    wide_row_h = (grid_h - 2 * row_gutter) - 2 * pair_row_h  # row 3

    col_w = (grid_w - col_gutter) / 2

    # Layout for each of the 5 boxes: (col_x, col_y, col_w, col_h, content_index)
    layout = [
        (grid_x, grid_y, col_w, pair_row_h, 0),                                          # SBOM
        (grid_x + col_w + col_gutter, grid_y, col_w, pair_row_h, 1),                      # Signature
        (grid_x, grid_y + pair_row_h + row_gutter, col_w, pair_row_h, 2),                 # Artifacts
        (grid_x + col_w + col_gutter, grid_y + pair_row_h + row_gutter,
            col_w, pair_row_h, 3),                                                        # Location identity
        (grid_x, grid_y + 2 * (pair_row_h + row_gutter), grid_w, wide_row_h, 4),          # Day-1 + Day-2
    ]

    icon_size = 44
    icon_pad = 24

    for (bx, by, bw, bh, content_idx) in layout:
        icon_file, headline, sub1, sub2 = BOXES[content_idx]

        # Box itself.
        _add_rounded_rect(slide, bx, by, bw, bh,
                           stroke=BLUE, stroke_pt=1.5, corner_pct=0.06)

        # Icon (top-left, brand-blue). rasterize_recolored is the callable
        # passed in by the deck builder so we don't import its module here.
        icon_path = icons_dir / icon_file
        if icon_path.exists():
            png = rasterize_recolored(icon_path, 96, "0F6BFF",
                                      stroke_width=icon_stroke)
            slide.shapes.add_picture(str(png),
                                       px(bx + icon_pad), px(by + icon_pad),
                                       width=px(icon_size),
                                       height=px(icon_size))

        # Headline text — sits next to the icon, vertically centred to icon.
        text_x = bx + icon_pad + icon_size + 16
        text_w = bw - (icon_pad + icon_size + 16) - icon_pad
        head_h = icon_size  # same as icon, so headline aligns to icon centre
        tf = _add_textbox(slide, text_x, by + icon_pad, text_w, head_h,
                           anchor="ctr")
        p = tf.paragraphs[0]
        _styled_run(p, headline,
                    size_pt=21, bold=True, color=BLUE_MID,
                    all_caps=True, letter_spacing="80")

        # Subtext (1-2 lines) below the headline+icon row.
        sub_y = by + icon_pad + icon_size + 14
        sub_h = bh - (icon_pad + icon_size + 14) - icon_pad
        tf = _add_textbox(slide, bx + icon_pad, sub_y,
                           bw - 2 * icon_pad, sub_h)
        p = tf.paragraphs[0]
        p.line_spacing = 1.2
        _styled_run(p, sub1, size_pt=17, color=BLACK)
        if sub2:
            p2 = tf.add_paragraph()
            p2.line_spacing = 1.2
            _styled_run(p2, sub2, size_pt=17, color=BLACK)

"""
slide_6_native.py — Slide 6 "Pack · Sign · Transport · Deploy → Sovereign Cloud"
rendered as native PowerPoint shapes (rounded rectangles, recoloured icons,
freeform cloud silhouette, textboxes), not a rasterized SVG.

Why native: same rationale as slide_4b_native — preserve typography, stay
crisp at every zoom, and keep the slide hand-editable inside PowerPoint.

Composition (left → right inside the source viewBox 1920×540):

      ┌──── Pack ────┐  →  ┌──── Sign ────┐  →  ┌── Transport ──┐  →  ┌── Deploy ──┐  ⇒  ☁🔒
      grey-soft card     same chrome           same chrome           same chrome     SOVEREIGN
      brand-blue rule    Tabler shield         Tabler world          Tabler rocket   CLOUD
      Tabler package     "SIGN" + sub          "TRANSPORT" + sub     "DEPLOY" + sub  caption
      "PACK" + sub

Source SVG: diagrams/05-pack-sign-transport-deploy-v2.svg (1920×540 viewBox).

Cloud-lock target glyph
-----------------------
The Sovereign Cloud target is a bespoke composition: brand-blue cloud
silhouette, white lock shackle (stroked) + white lock body (filled), and a
brand-blue keyhole inside the white lock body. Because the colours are not
uniform (currentColor would only give us one), this is drawn natively as
MSO_SHAPE.FREEFORM (cloud) plus shapes for shackle / body / keyhole, taking
option (a) from the task brief — cleaner than dragging a multi-colour SVG
through the recolour rasterizer, which only knows how to swap currentColor.

Public entry point
------------------
    add_pack_sign_transport_deploy_native(
        slide, *, x, y, w, h, icons_dir, rasterize_recolored,
    )

`x`, `y`, `w`, `h` describe the slot in slide coordinates (px @ 96 dpi);
the SVG's 1920×540 viewBox is mapped uniformly into that slot and centred.
`rasterize_recolored` is injected (not imported) so this module stays free
of build-script-specific dependencies, mirroring slide_4b_native.
"""
from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


PX = 9525


def px(n: float) -> Emu:
    return Emu(int(n * PX))


# OCM brand palette (mirror of slide_4b_native — kept local for self-containment).
BLUE       = RGBColor(0x0F, 0x6B, 0xFF)   # brand-blue (icons, cloud, keyhole)
BLUE_MID   = RGBColor(0x0A, 0x3A, 0x99)   # mid blue (labels)
GREY_SOFT  = RGBColor(0xF3, 0xF4, 0xF6)   # card fill
GREY_MID   = RGBColor(0x6B, 0x72, 0x80)   # arrows
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)   # body text
BLACK      = RGBColor(0x00, 0x00, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)


# -----------------------------------------------------------------------------
# Shared helpers (same shape as slide_4b_native)
# -----------------------------------------------------------------------------

def _styled_run(p, text, *, size_pt, bold=False, color=BLACK,
                font="Inter", italic=False, all_caps=False,
                letter_spacing=None):
    """Add a run with the OCM type system applied."""
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size_pt)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if all_caps or letter_spacing is not None:
        rPr = r._r.get_or_add_rPr()
        if all_caps:
            rPr.set("cap", "all")
        if letter_spacing is not None:
            rPr.set("spc", str(letter_spacing))
    return r


def _add_textbox(slide, x, y, w, h, *, anchor="t", margins=0, align=None):
    """Borderless textbox at the given coords with consistent inset margins."""
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = px(margins)
    tf.margin_top = tf.margin_bottom = px(margins)
    tf.word_wrap = True
    if anchor == "ctr":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "b":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    if align is not None:
        tf.paragraphs[0].alignment = align
    return tf


def _add_rounded_rect(slide, x, y, w, h, *,
                      fill=WHITE, stroke=BLUE, stroke_pt=1.5,
                      corner_pct=0.04, shadow=False):
    """Rounded rectangle with explicit fill, stroke, and corner radius.

    shadow=True attaches a soft outer drop-shadow matching the SVG's
    `card-shadow` filter (blur 3 px, dy 3 px, 30% black). PowerPoint's
    native <a:outerShdw> survives editing and selection-grouping, so the
    user can move the card and the shadow follows.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h),
    )
    shape.adjustments[0] = corner_pct
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = stroke
        shape.line.width = Pt(stroke_pt)
    if shadow:
        from lxml import etree
        from pptx.oxml.ns import qn
        spPr = shape._element.spPr
        # Drop any existing effectLst (autoshape defaults usually have none),
        # then attach an outerShdw matching the SVG: blur 3 px = 28575 EMU,
        # offset (0, 3 px) = (dir 5400000 = 90°, dist 28575 EMU), 30% black.
        for old in spPr.findall(qn("a:effectLst")):
            spPr.remove(old)
        effectLst = etree.SubElement(spPr, qn("a:effectLst"))
        outerShdw = etree.SubElement(effectLst, qn("a:outerShdw"))
        outerShdw.set("blurRad", "28575")     # 3 px
        outerShdw.set("dist",    "28575")     # 3 px
        outerShdw.set("dir",     "5400000")   # 90° = straight down
        outerShdw.set("rotWithShape", "0")
        clr = etree.SubElement(outerShdw, qn("a:srgbClr"))
        clr.set("val", "000000")
        alpha = etree.SubElement(clr, qn("a:alpha"))
        alpha.set("val", "30000")             # 30 %
    return shape


def _add_rect(slide, x, y, w, h, *, fill=BLUE, stroke=None, stroke_pt=0.0):
    """Plain (non-rounded) rectangle — used for the brand-blue top rule."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = stroke
        shape.line.width = Pt(stroke_pt)
    return shape


def _add_arrow(slide, x1, y1, x2, y2, *, stroke=GREY_MID, stroke_pt=2.5,
               head_w=14, head_h=10):
    """Short straight horizontal arrow: thin rectangle stem + triangular head.

    We avoid MSO_SHAPE.RIGHT_ARROW because its head/stem ratio bulks the
    short inter-card arrows. Drawing stem + head separately gives us the
    same proportions as the SVG (line-with-marker) and lets us match
    stroke_pt exactly.
    """
    # Stem: thin rectangle, vertically centred on the arrow line.
    stem_h = stroke_pt * 0.9          # visual stroke width in px terms
    stem_w = max(0, (x2 - x1) - head_w + 1)  # leave room for the head
    stem_y = y1 - stem_h / 2.0
    if stem_w > 0:
        stem = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, px(x1), px(stem_y), px(stem_w), px(stem_h),
        )
        stem.fill.solid()
        stem.fill.fore_color.rgb = stroke
        stem.line.fill.background()

    # Head: triangle pointing right, anchored at (x2, y1).
    head_x = x2 - head_w
    head_y = y1 - head_h / 2.0
    head = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE, px(head_x), px(head_y),
        px(head_w), px(head_h),
    )
    # MSO right-triangle has its right-angle at bottom-left; we want a
    # symmetric isoceles triangle pointing right. Replace it with a freeform.
    sp = head._element
    sp.getparent().remove(sp)

    # python-pptx FreeformBuilder treats vertices as ABSOLUTE coordinates in
    # the freeform's local coordinate system (with start_x/start_y as the
    # origin), not deltas. _dx/_dy compute max-min over the absolute set,
    # then _width = _dx * scale. Passing deltas in EMU on top of an EMU
    # start point makes the bbox span from 0..start_x_emu — i.e. across half
    # the slide. Pass absolute EMU coordinates instead.
    apex_x = px(head_x + head_w)
    apex_top_y = px(head_y)
    apex_bot_y = px(head_y + head_h)
    builder = slide.shapes.build_freeform(
        px(head_x), px(head_y + head_h / 2.0), scale=1.0,
    )
    builder.add_line_segments([
        (apex_x, apex_top_y),
        (apex_x, apex_bot_y),
    ], close=True)
    head_shape = builder.convert_to_shape()
    head_shape.fill.solid()
    head_shape.fill.fore_color.rgb = stroke
    head_shape.line.fill.background()


# -----------------------------------------------------------------------------
# Cloud-lock glyph (option (a) from the task brief)
# -----------------------------------------------------------------------------

def _add_cloud_lock(slide, ox, oy, scale, *,
                     icons_dir=None, rasterize_recolored=None,
                     icon_stroke=None):
    """Draw the bespoke cloud-lock at (ox, oy) using `scale` to convert
    the SVG's 226-px-wide local space (cloud bounding box ≈ 218×130) into
    slide-local px.

    Cloud silhouette is a custGeom Freeform (matches the SVG's bespoke
    blue blob — no Tabler equivalent). The lock that sits inside it is
    overlaid as a recoloured Tabler lock.svg when `icons_dir` and
    `rasterize_recolored` are provided; otherwise we fall back to the
    hand-drawn body+shackle+keyhole composition (kept for callers that
    can't supply the rasteriser).

    Cloud silhouette path (filled #0F6BFF):
        M 56 130 C 28 130 4 110 4 80 C 4 54 26 32 54 36
              C 64 14 92 6 112 22 C 132 6 162 14 174 38
              C 200 36 222 56 222 80 C 222 108 198 130 174 130 Z
    """
    def s(v):  # local-px → px
        return v * scale

    # --- Cloud silhouette (filled brand-blue, freeform with cubic beziers) ---
    # python-pptx has no public bezier builder on freeform, but we can drop
    # to lxml on the underlying <a:path> after building the move-to. The
    # cleanest route is to construct the whole path via a custGeom XML.
    from lxml import etree
    from pptx.oxml.ns import qn

    # Bounding box for the freeform: the cloud spans x∈[4,222], y∈[6,130]
    # in local space. Translate so the freeform's local origin is (0,0).
    cloud_min_x, cloud_min_y = 4.0, 6.0
    cloud_max_x, cloud_max_y = 222.0, 130.0
    cloud_local_w = cloud_max_x - cloud_min_x   # 218
    cloud_local_h = cloud_max_y - cloud_min_y   # 124

    cloud_x = ox + s(cloud_min_x)
    cloud_y = oy + s(cloud_min_y)
    cloud_w = s(cloud_local_w)
    cloud_h = s(cloud_local_h)

    # Build a freeform that starts at the move-to, then we patch the
    # underlying <a:path> with the full bezier sequence so we don't have
    # to approximate the curves with line segments.
    builder = slide.shapes.build_freeform(
        px(cloud_x), px(cloud_y), scale=1.0,
    )
    # Add a tiny throwaway segment so build_freeform produces a valid path
    # element we can rewrite.
    builder.add_line_segments([(1, 0)], close=True)
    cloud_shape = builder.convert_to_shape()
    cloud_shape.fill.solid()
    cloud_shape.fill.fore_color.rgb = BLUE
    cloud_shape.line.fill.background()

    # Resize the freeform's bounding box to match the cloud's local extent.
    cloud_shape.left = px(cloud_x)
    cloud_shape.top = px(cloud_y)
    cloud_shape.width = px(cloud_w)
    cloud_shape.height = px(cloud_h)

    # Rewrite the path: the freeform's <a:custGeom><a:pathLst><a:path> needs
    # to describe the cloud in path-local coordinates (0..pathW, 0..pathH).
    # python-pptx's freeform path is sized in EMU; we mirror that here so
    # the existing bbox + custGeom stay in sync.
    sp_tree = cloud_shape._element.spPr
    cust = sp_tree.find(qn("a:custGeom"))
    path_lst = cust.find(qn("a:pathLst"))
    # Wipe old <a:path> and inject our own.
    for old in list(path_lst):
        path_lst.remove(old)

    path_w_emu = int(cloud_w * PX)
    path_h_emu = int(cloud_h * PX)
    path_el = etree.SubElement(
        path_lst, qn("a:path"),
        {"w": str(path_w_emu), "h": str(path_h_emu)},
    )

    # Translate the SVG path's absolute coords into path-local (subtract
    # cloud_min_x, cloud_min_y) and convert to EMU.
    def pt_emu(svg_x, svg_y):
        local_x = (svg_x - cloud_min_x) * scale
        local_y = (svg_y - cloud_min_y) * scale
        return int(local_x * PX), int(local_y * PX)

    def add_pt(parent, tag, svg_x, svg_y):
        ex, ey = pt_emu(svg_x, svg_y)
        pt = etree.SubElement(parent, qn(tag))
        etree.SubElement(pt, qn("a:pt"), {"x": str(ex), "y": str(ey)})
        return pt

    def add_moveto(svg_x, svg_y):
        m = etree.SubElement(path_el, qn("a:moveTo"))
        ex, ey = pt_emu(svg_x, svg_y)
        etree.SubElement(m, qn("a:pt"), {"x": str(ex), "y": str(ey)})

    def add_cubic(c1x, c1y, c2x, c2y, ex_, ey_):
        cb = etree.SubElement(path_el, qn("a:cubicBezTo"))
        for sx, sy in [(c1x, c1y), (c2x, c2y), (ex_, ey_)]:
            x_emu, y_emu = pt_emu(sx, sy)
            pt = etree.SubElement(cb, qn("a:pt"),
                                  {"x": str(x_emu), "y": str(y_emu)})

    # SVG path (verbatim):
    #   M 56 130
    #   C 28 130 4 110 4 80
    #   C 4 54 26 32 54 36
    #   C 64 14 92 6 112 22
    #   C 132 6 162 14 174 38
    #   C 200 36 222 56 222 80
    #   C 222 108 198 130 174 130
    #   Z
    add_moveto(56, 130)
    add_cubic(28, 130,  4, 110,   4,  80)
    add_cubic( 4,  54, 26,  32,  54,  36)
    add_cubic(64,  14, 92,   6, 112,  22)
    add_cubic(132,  6, 162, 14, 174,  38)
    add_cubic(200, 36, 222, 56, 222,  80)
    add_cubic(222, 108, 198, 130, 174, 130)
    etree.SubElement(path_el, qn("a:close"))

    # --- Lock glyph -------------------------------------------------------
    # Two paths: a Tabler lock.svg overlaid on the cloud (preferred — keeps
    # the deck's icon vocabulary consistent) or the hand-drawn body+
    # shackle+keyhole composition (kept as a fallback for older callers
    # that don't pass icons_dir / rasterize_recolored).
    lock_path = (icons_dir / "lock.svg") if icons_dir else None
    if lock_path and lock_path.exists() and rasterize_recolored is not None:
        # The hand-drawn lock occupies x∈[91, 135], y∈[36, 98] in cloud-
        # local coords — a 44×62 box. Tabler lock.svg has a 24×24 viewBox
        # with the body filling roughly the bottom half — so we centre a
        # square box on the same midpoint and size it to match the
        # hand-drawn glyph's height.
        lock_size_local = 56               # tweak: a touch wider than the
                                            # 44-wide hand glyph because
                                            # the Tabler lock has internal
                                            # padding inside its viewBox
        lock_cx_local = (91 + 135) / 2     # 113 — middle of the body
        lock_cy_local = (36 + 98) / 2      # 67  — middle of the lock
        lock_x = ox + s(lock_cx_local - lock_size_local / 2)
        lock_y = oy + s(lock_cy_local - lock_size_local / 2)
        lock_size = s(lock_size_local)
        png = rasterize_recolored(lock_path, 192, "FFFFFF",
                                  stroke_width=icon_stroke)
        slide.shapes.add_picture(
            str(png), px(lock_x), px(lock_y),
            width=px(lock_size), height=px(lock_size),
        )
    else:
        # Fallback: hand-drawn body + shackle + keyhole.
        # --- Lock body (filled white, rounded rect rx=4) ------------------
        body_x = ox + s(91)
        body_y = oy + s(66)
        body_w = s(44)
        body_h = s(32)
        # rx=4 in a 44-wide rect ≈ 9.1% of the shorter side (32 high) → use
        # PowerPoint's rounded-rect adjustment. corner_pct uses fraction of
        # min(w,h); 4/32 = 0.125.
        _add_rounded_rect(slide, body_x, body_y, body_w, body_h,
                          fill=WHITE, stroke=None, corner_pct=0.125)

        # --- Lock shackle (white stroke, no fill, rounded U) --------------
        # SVG path: M 99 66 L 99 54 C 99 42 105 36 113 36 C 121 36 127 42 127 54 L 127 66
        shackle_min_x, shackle_min_y = 99.0, 36.0
        shackle_max_x, shackle_max_y = 127.0, 66.0
        shackle_local_w = shackle_max_x - shackle_min_x   # 28
        shackle_local_h = shackle_max_y - shackle_min_y   # 30

        shackle_x = ox + s(shackle_min_x)
        shackle_y = oy + s(shackle_min_y)
        shackle_w = s(shackle_local_w)
        shackle_h = s(shackle_local_h)

        builder = slide.shapes.build_freeform(
            px(shackle_x), px(shackle_y), scale=1.0,
        )
        builder.add_line_segments([(1, 0)], close=False)
        shackle_shape = builder.convert_to_shape()
        shackle_shape.fill.background()
        shackle_shape.line.color.rgb = WHITE
        shackle_shape.line.width = Pt(4.0 * scale)
        shackle_shape.left = px(shackle_x)
        shackle_shape.top = px(shackle_y)
        shackle_shape.width = px(shackle_w)
        shackle_shape.height = px(shackle_h)

        sp_tree = shackle_shape._element.spPr
        cust = sp_tree.find(qn("a:custGeom"))
        path_lst = cust.find(qn("a:pathLst"))
        for old in list(path_lst):
            path_lst.remove(old)

        path_w_emu = int(shackle_w * PX)
        path_h_emu = int(shackle_h * PX)
        path_el = etree.SubElement(
            path_lst, qn("a:path"),
            {"w": str(path_w_emu), "h": str(path_h_emu), "fill": "none"},
        )

        def shackle_pt_emu(svg_x, svg_y):
            local_x = (svg_x - shackle_min_x) * scale
            local_y = (svg_y - shackle_min_y) * scale
            return int(local_x * PX), int(local_y * PX)

        def shackle_moveto(svg_x, svg_y):
            m = etree.SubElement(path_el, qn("a:moveTo"))
            ex, ey = shackle_pt_emu(svg_x, svg_y)
            etree.SubElement(m, qn("a:pt"), {"x": str(ex), "y": str(ey)})

        def shackle_lineto(svg_x, svg_y):
            ln = etree.SubElement(path_el, qn("a:lnTo"))
            ex, ey = shackle_pt_emu(svg_x, svg_y)
            etree.SubElement(ln, qn("a:pt"), {"x": str(ex), "y": str(ey)})

        def shackle_cubic(c1x, c1y, c2x, c2y, ex_, ey_):
            cb = etree.SubElement(path_el, qn("a:cubicBezTo"))
            for cx_pt, cy_pt in [(c1x, c1y), (c2x, c2y), (ex_, ey_)]:
                x_emu, y_emu = shackle_pt_emu(cx_pt, cy_pt)
                etree.SubElement(cb, qn("a:pt"),
                                  {"x": str(x_emu), "y": str(y_emu)})

        shackle_moveto(99, 66)
        shackle_lineto(99, 54)
        shackle_cubic(99, 42, 105, 36, 113, 36)
        shackle_cubic(121, 36, 127, 42, 127, 54)
        shackle_lineto(127, 66)

        # --- Keyhole (brand-blue circle + stem) ----------------------------
        kc_x = ox + s(113 - 3)
        kc_y = oy + s(78 - 3)
        kc_d = s(6)
        keyhole_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, px(kc_x), px(kc_y), px(kc_d), px(kc_d),
        )
        keyhole_circle.fill.solid()
        keyhole_circle.fill.fore_color.rgb = BLUE
        keyhole_circle.line.fill.background()

        # Stem rect: x=111 y=78 w=4 h=10
        ks_x = ox + s(111)
        ks_y = oy + s(78)
        ks_w = s(4)
        ks_h = s(10)
        _add_rect(slide, ks_x, ks_y, ks_w, ks_h, fill=BLUE)


# -----------------------------------------------------------------------------
# Per-card content (each: icon_filename, label, sub1, sub2)
# -----------------------------------------------------------------------------

CARDS = [
    ("package.svg",      "PACK",      "Bundle your software",       "One source of truth."),
    ("shield-check.svg", "SIGN",      "One signature covers",       "all artifacts."),
    ("world.svg",        "TRANSPORT", "Across any boundary.",       "Even air-gapped."),
    ("rocket.svg",       "DEPLOY",    "Bring your own GitOps.",     "Use OCM's K8s Controllers"),
]


# Icon stroke-width presets are centralised in icon_strokes.py so all native
# slides agree on the 1.0 / 1.5 / 2.0 scale.
from icon_strokes import STROKE_THIN, STROKE_REGULAR, STROKE_BOLD


# -----------------------------------------------------------------------------
# Public entry point
# -----------------------------------------------------------------------------

def add_pack_sign_transport_deploy_native(slide, *, x, y, w, h,
                                           icons_dir: Path,
                                           rasterize_recolored,
                                           icon_stroke: float = STROKE_THIN):
    """Draw the Pack · Sign · Transport · Deploy → Sovereign Cloud diagram
    natively into the (x, y, w, h) slot.

    icon_stroke selects the rasterised icon weight — STROKE_THIN matches
    the SVG variant's outline weight; STROKE_REGULAR is slightly heavier;
    STROKE_BOLD is Tabler's own default (visually heavy at 60-px display).

    Parameters mirror slide_4b_native.add_sbod_native_diagram so the deck
    builder can call both with the same callable wiring.

    The source SVG's 1920×540 viewBox is mapped uniformly into the slot,
    centred on whichever axis has the slack.
    """
    # ---- Map the 1920×540 viewBox into the slot, uniformly + centred. ----
    SVG_W, SVG_H = 1920.0, 540.0
    scale = min(w / SVG_W, h / SVG_H)
    drawn_w = SVG_W * scale
    drawn_h = SVG_H * scale
    ox = x + (w - drawn_w) / 2.0    # slot-local origin for SVG coords
    oy = y + (h - drawn_h) / 2.0

    def s(v):  # SVG-local px → slide px
        return v * scale

    # ---- Step cards (4 × 360-wide cards, 30-px gaps) ----------------------
    # SVG geometry: cards at SVG-x = 10, 420, 830, 1240; y=70; 360×400.
    # Brand-blue top rule = 4 px tall sitting on the top edge of the card.
    card_top_rule_h = 4
    card_w_svg = 360
    card_h_svg = 400
    card_y_svg = 70

    # Inner-card padding (SVG): icon at (32,36) size 60×60, label baseline
    # at SVG-y=170, body at y=216 with 36-px line height. We re-flow these
    # into native textboxes; SVG y-coordinates correspond to the *baseline*
    # of <text>, while textboxes anchor to the top — adjust accordingly.
    card_icon_x_svg = 32
    card_icon_y_svg = 36
    card_icon_size_svg = 60

    # Label "PACK" etc.: SVG text baseline y=170 with font-size 34 → top of
    # the cap-height roughly y=170-34=136. Allow some breathing room above.
    card_label_y_svg = 136
    card_label_h_svg = 44

    # Body text: first line baseline y=216 → top ≈ 216-26=190. Two lines
    # 36 px apart → block height 26 + 36 = 62 (top → just under the second
    # baseline) → use ~80 px to give us margin.
    card_body_y_svg = 190
    card_body_h_svg = 100

    card_x_svgs = [10, 420, 830, 1240]

    for (cx_svg, (icon_file, label, sub1, sub2)) in zip(card_x_svgs, CARDS):
        # --- Card chrome ---------------------------------------------------
        card_x = ox + s(cx_svg)
        card_y = oy + s(card_y_svg)
        card_w = s(card_w_svg)
        card_h = s(card_h_svg)

        # Grey-soft rounded card (rx=14 in a 360-wide card → 14/360 ≈ 0.039
        # of the longer side; corner_pct uses min(w,h)=360 so 0.039 it is).
        # shadow=True paints the soft drop-shadow the SVG variant has via
        # its <filter id="card-shadow"> definition.
        _add_rounded_rect(slide, card_x, card_y, card_w, card_h,
                          fill=GREY_SOFT, stroke=None,
                          corner_pct=14.0 / min(card_w_svg, card_h_svg),
                          shadow=True)

        # Brand-blue 4-px top rule sitting flush on the top edge.
        _add_rect(slide, card_x, card_y, card_w, s(card_top_rule_h),
                  fill=BLUE)

        # --- Icon (top-left, brand-blue) -----------------------------------
        icon_path = icons_dir / icon_file
        if icon_path.exists():
            icon_x = card_x + s(card_icon_x_svg)
            icon_y = card_y + s(card_icon_y_svg)
            icon_size = s(card_icon_size_svg)
            # Rasterise at 192 px (3.2× supersample for the 60-px display).
            # stroke_width comes from icon_stroke; STROKE_THIN matches the
            # SVG variant of this slide.
            png = rasterize_recolored(icon_path, 192, "0F6BFF",
                                      stroke_width=icon_stroke)
            slide.shapes.add_picture(
                str(png), px(icon_x), px(icon_y),
                width=px(icon_size), height=px(icon_size),
            )

        # --- Label "PACK" / "SIGN" / "TRANSPORT" / "DEPLOY" ----------------
        label_x = card_x + s(32)
        label_y = card_y + s(card_label_y_svg)
        label_w = card_w - s(64)
        label_h = s(card_label_h_svg)
        tf = _add_textbox(slide, label_x, label_y, label_w, label_h,
                          anchor="ctr")
        p = tf.paragraphs[0]
        # Source SVG: font-size 34px, font-weight 700, fill #0A3A99 (BLUE_MID),
        # letter-spacing 0.6px (SVG attribute) → ~60 in PPT spc units.
        # 34 px @ 96 dpi mathematically rounds to 25.5 pt, but the SVG
        # variant on this slide is rasterised at 1.0× into the diagram slot
        # which makes the text read closer to 28 pt at PowerPoint's display
        # zoom. User confirmed 28 pt visually matches the SVG render.
        # No letter_spacing — Inter at 28pt with 0.6pt spc reads visibly
        # looser than the SVG render at 0.6 SVG-units; Inter's metrics
        # already include enough tracking.
        _styled_run(p, label,
                    size_pt=26, bold=True, color=BLUE_MID)

        # --- Body sub-text (2 lines) ---------------------------------------
        body_x = card_x + s(32)
        body_y = card_y + s(card_body_y_svg)
        body_w = card_w - s(64)
        body_h = s(card_body_h_svg)
        tf = _add_textbox(slide, body_x, body_y, body_w, body_h)
        p = tf.paragraphs[0]
        p.line_spacing = 1.25
        # 26 px @ 96 dpi ≈ 19.5 pt → use 19 pt to leave room for the
        # second line on tighter slot heights.
        _styled_run(p, sub1, size_pt=19, color=DARK_GREY)
        p2 = tf.add_paragraph()
        p2.line_spacing = 1.25
        _styled_run(p2, sub2, size_pt=19, color=DARK_GREY)

    # ---- Inter-card arrows (3) at SVG y=270, x = 380→406, 790→816,
    #      1200→1226. These sit in the 30-px gaps between cards. ---------
    arrow_y_svg = 270
    inter_arrows = [(380, 406), (790, 816), (1200, 1226)]
    for (a1, a2) in inter_arrows:
        _add_arrow(slide,
                   ox + s(a1), oy + s(arrow_y_svg),
                   ox + s(a2), oy + s(arrow_y_svg),
                   stroke=GREY_MID, stroke_pt=2.5,
                   head_w=s(12), head_h=s(10))

    # ---- Big arrow from Deploy card to target (3-px stroke, slightly
    #      thicker per the source SVG line @ stroke-width=3). -------------
    _add_arrow(slide,
               ox + s(1610), oy + s(arrow_y_svg),
               ox + s(1660), oy + s(arrow_y_svg),
               stroke=GREY_MID, stroke_pt=3.0,
               head_w=s(14), head_h=s(12))

    # ---- Sovereign Cloud target (no card chrome) -------------------------
    # Source group: translate(1670, 60). Inside it:
    #   - cloud-lock at translate(8, 50)        → glyph origin (1678, 110)
    #   - "SOVEREIGN" centred at x=124 y=252    → SVG (1794, 312)
    #   - "CLOUD"               x=124 y=284    → SVG (1794, 344)
    #   - "Verify at destination."  y=320       → SVG (1794, 380)
    #   - "No callback upstream."   y=358       → SVG (1794, 418)
    target_origin_svg = (1670, 60)
    glyph_local_origin_svg = (8, 50)
    glyph_origin_svg = (
        target_origin_svg[0] + glyph_local_origin_svg[0],
        target_origin_svg[1] + glyph_local_origin_svg[1],
    )

    _add_cloud_lock(slide,
                    ox + s(glyph_origin_svg[0]),
                    oy + s(glyph_origin_svg[1]),
                    scale=scale,
                    icons_dir=icons_dir,
                    rasterize_recolored=rasterize_recolored,
                    icon_stroke=icon_stroke)

    # Centre x for the target labels, in slide-px.
    target_centre_x = ox + s(target_origin_svg[0] + 124)

    # SOVEREIGN / CLOUD label block — two-line, centred, mid-blue, 34 px
    # SVG → ~26 pt. Anchor textbox so its centre x equals target_centre_x.
    label_block_w = s(280)
    label_block_x = target_centre_x - label_block_w / 2.0
    label_block_y = oy + s(target_origin_svg[1] + 252 - 34)  # baseline → top
    label_block_h = s(80)
    tf = _add_textbox(slide, label_block_x, label_block_y,
                      label_block_w, label_block_h,
                      align=PP_ALIGN.CENTER)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _styled_run(p, "SOVEREIGN",
                size_pt=26, bold=True, color=BLUE_MID)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    _styled_run(p2, "CLOUD",
                size_pt=26, bold=True, color=BLUE_MID)

    # Caption: "Verify at destination." / "No callback upstream." — dark
    # grey, 26 px SVG → ~19 pt, centred under the label block.
    cap_block_w = s(340)
    cap_block_x = target_centre_x - cap_block_w / 2.0
    cap_block_y = oy + s(target_origin_svg[1] + 320 - 26)
    cap_block_h = s(90)
    tf = _add_textbox(slide, cap_block_x, cap_block_y,
                      cap_block_w, cap_block_h,
                      align=PP_ALIGN.CENTER)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.25
    _styled_run(p, "Verify at destination.", size_pt=19, color=DARK_GREY)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.line_spacing = 1.25
    _styled_run(p2, "No callback upstream.", size_pt=19, color=DARK_GREY)

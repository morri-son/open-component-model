"""
slide_4b_native_v2.py — Slide 4b "An envelope, not a list." VARIANT B.

Variant B reframes the same SBOD-anatomy story as Variant A but in a more
linear, less grid-y composition:

    ┌──────────────────────────────────────────────────────────────────┐
    │   github.com/acme/webshop:v1.0.0                                 │
    │   Location-independent. Same identity, every registry.           │   ← Identity header (full-width tinted card)
    └──────────────────────────────────────────────────────────────────┘
            ┌────────────────────────────────────┐
            │ 📦  Docker Image                    │ ⎫
            ├────────────────────────────────────┤ ⎬
            │ ⎈  Helm Chart                       │ ⎬
            ├────────────────────────────────────┤ ⎬     🔒 SIGNATURE
            │ ☸  Kubernetes Deployment Manifests  │ ⎬     One digest covers all.
            ├────────────────────────────────────┤ ⎬
            │ 📄  Configuration Files             │ ⎬
            ├────────────────────────────────────┤ ⎬
            │ 📋  SBOM  ·  CycloneDX              │ ⎭
            └────────────────────────────────────┘

The identity sits ABOVE the artifacts (it's a property of the whole envelope,
not a sibling). The artifact list reads top-to-bottom — Docker Image, Helm
Chart, K8s Manifests, Configuration Files, SBOM. SBOM is one of five rows,
deliberately last (counter-narrative cue: SBOM is just inventory; the rest is
what gets delivered). The signature curly-brace visually wraps the whole
artifact list — one digest covers every row.

Public entry point mirrors slide_4b_native.add_sbod_native_diagram:

    from slide_4b_native_v2 import add_sbom_inside_sbod_native_v2
    add_sbom_inside_sbod_native_v2(
        slide, x=60, y=240, w=1800, h=780,
        icons_dir=ICONS_DIR, rasterize_recolored=rasterize_svg_recolored,
    )
"""
from __future__ import annotations

from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


PX = 9525  # EMU per px @ 96 dpi


def px(n: float) -> Emu:
    return Emu(int(n * PX))


# OCM brand palette (mirror of the Variant A constants — kept local so this
# module is self-contained and importable by both deck builders).
BLUE       = RGBColor(0x0F, 0x6B, 0xFF)   # brand-blue
BLUE_MID   = RGBColor(0x0A, 0x3A, 0x99)   # mid blue used for headlines
BLUE_TINT  = RGBColor(0xEA, 0xF2, 0xFF)   # very pale blue card tint (local-only)
GREY_SOFT  = RGBColor(0xF3, 0xF4, 0xF6)   # secondary surface
GREY_MID   = RGBColor(0x6B, 0x72, 0x80)   # secondary text
GREY_LINE  = RGBColor(0xE5, 0xE7, 0xEB)   # row dividers
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)   # primary body text
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)


def _styled_run(p, text, *, size_pt, bold=False, color=BLACK,
                font="Inter", italic=False, all_caps=False,
                letter_spacing=None):
    """Add a run with the OCM type system applied."""
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size_pt)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if all_caps or letter_spacing is not None:
        rPr = r._r.get_or_add_rPr()
        if all_caps:
            rPr.set("cap", "all")
        if letter_spacing is not None:
            rPr.set("spc", str(letter_spacing))
    return r


def _add_textbox(slide, x, y, w, h, *, anchor="t", align=None, margins=0):
    """Borderless textbox at the given coords, with consistent inset margins."""
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = px(margins)
    tf.margin_top = tf.margin_bottom = px(margins)
    tf.word_wrap = True
    if anchor == "ctr":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "b":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    if align is not None:
        tf.paragraphs[0].alignment = align
    return tf


def _add_rounded_rect(slide, x, y, w, h, *,
                      fill=WHITE, stroke=BLUE, stroke_pt=1.5,
                      corner_pct=0.04):
    """Rounded rectangle with explicit fill, stroke, and corner radius."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h),
    )
    shape.adjustments[0] = corner_pct
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = stroke
        shape.line.width = Pt(stroke_pt)
    return shape


def _add_line(slide, x1, y1, x2, y2, *, color=GREY_LINE, width_pt=0.75):
    """Plain straight connector (used for the row dividers in the list)."""
    from pptx.enum.shapes import MSO_CONNECTOR
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, px(x1), px(y1), px(x2), px(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    return line


def _add_right_brace(slide, x, y, w, h, *, color=BLUE, width_pt=2.5):
    """Curly right brace `}` spanning (x, y, w, h). Stroke only, no fill."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_BRACE, px(x), px(y), px(w), px(h),
    )
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)
    return shape


# Per-row content. icon filename + label (+ optional secondary inline label).
ROWS = [
    ("docker.svg",     "Docker Image",                         None),
    ("helm.svg",       "Helm Chart",                           None),
    ("kubernetes.svg", "Kubernetes Deployment Manifests",      None),
    ("file-text.svg",  "Configuration Files",                  None),
    ("shield.svg",     "SBOM",                                 "CycloneDX"),
]


from icon_strokes import STROKE_THIN


def add_sbom_inside_sbod_native_v2(slide, *, x=60, y=240, w=1800, h=780,
                                    icons_dir: Path,
                                    rasterize_recolored,
                                    icon_stroke=STROKE_THIN):
    """Render Variant B of the SBOM-inside-SBOD diagram into (x, y, w, h).

    icon_stroke selects Tabler-icon weight (icon_strokes.STROKE_THIN /
    STROKE_REGULAR / STROKE_BOLD). Default STROKE_THIN.

    Layout inside the slot:
        - Top band       (~110 px): identity header card, centred ~65% width.
        - Middle band    (~520 px): vertical artifact list (left-of-centre)
                                     plus brace + signature label on the right.
        - Slot is 1800 × 780; this module assumes the standard Content/Diagram
          slot but accepts arbitrary (x, y, w, h) for reuse.
    """
    # ------------------------------------------------------------------ Header
    header_card_w = int(w * 0.66)
    header_card_h = 110
    header_card_x = x + (w - header_card_w) / 2
    header_card_y = y + 24

    _add_rounded_rect(slide,
                       header_card_x, header_card_y,
                       header_card_w, header_card_h,
                       fill=BLUE_TINT, stroke=BLUE, stroke_pt=1.5,
                       corner_pct=0.12)

    # Identity line 1 — monospace-ish, large, bold-mid blue.
    pad_inner_x, pad_inner_y = 32, 18
    line1_h = 38
    tf = _add_textbox(slide,
                       header_card_x + pad_inner_x,
                       header_card_y + pad_inner_y,
                       header_card_w - 2 * pad_inner_x,
                       line1_h,
                       align=PP_ALIGN.CENTER)
    p = tf.paragraphs[0]
    _styled_run(p, "github.com/acme/webshop:v1.0.0",
                size_pt=22, bold=True, color=BLUE_MID,
                font="Consolas")

    # Identity line 2 — regular sans, smaller, grey.
    line2_h = 28
    tf = _add_textbox(slide,
                       header_card_x + pad_inner_x,
                       header_card_y + pad_inner_y + line1_h + 4,
                       header_card_w - 2 * pad_inner_x,
                       line2_h,
                       align=PP_ALIGN.CENTER)
    p = tf.paragraphs[0]
    _styled_run(p, "Location-independent. Same identity, every registry.",
                size_pt=15, color=GREY_MID)

    # ------------------------------------------------------------- Artifact list
    # The list + brace + signature copy form one composition. Centre it
    # horizontally inside the slot; the list sits left-of-centre, the brace
    # immediately right of it, the signature label right of the brace.

    list_w = 720
    row_h  = 84
    n_rows = len(ROWS)
    list_h = row_h * n_rows

    brace_w = 36
    brace_gap_left = 12   # gap between list and brace
    sig_gap_left   = 24   # gap between brace and signature copy
    sig_w = 320

    composition_w = list_w + brace_gap_left + brace_w + sig_gap_left + sig_w
    comp_x = x + (w - composition_w) / 2

    # Vertical positioning: place the list under the header with a comfortable
    # gap, but never overflow the slot — clamp if necessary.
    list_y_target = header_card_y + header_card_h + 42
    list_y_max    = y + h - list_h - 24
    list_y        = min(list_y_target, list_y_max)
    list_x        = comp_x

    # The single rounded-rect that encloses the artifact list.
    _add_rounded_rect(slide, list_x, list_y, list_w, list_h,
                       fill=WHITE, stroke=BLUE, stroke_pt=1.5,
                       corner_pct=0.04)

    # Row dividers (between consecutive rows — n_rows-1 lines).
    for i in range(1, n_rows):
        ly = list_y + i * row_h
        _add_line(slide,
                   list_x + 16, ly,
                   list_x + list_w - 16, ly,
                   color=GREY_LINE, width_pt=0.75)

    # Render each row: icon (left) + label (right of icon), vertically centred.
    icon_size = 36
    icon_pad_x = 24
    label_pad_x = 18

    for i, (icon_file, label, secondary) in enumerate(ROWS):
        row_y = list_y + i * row_h

        # Icon
        icon_path = icons_dir / icon_file
        if icon_path.exists():
            png = rasterize_recolored(icon_path, 96, "0F6BFF",
                                      stroke_width=icon_stroke)
            icon_y = row_y + (row_h - icon_size) / 2
            slide.shapes.add_picture(
                str(png),
                px(list_x + icon_pad_x), px(icon_y),
                width=px(icon_size), height=px(icon_size),
            )

        # Label (+ secondary inline label) — vertically centred to icon.
        text_x = list_x + icon_pad_x + icon_size + label_pad_x
        text_w = list_w - (text_x - list_x) - icon_pad_x
        tf = _add_textbox(slide, text_x, row_y, text_w, row_h,
                           anchor="ctr")
        p = tf.paragraphs[0]
        _styled_run(p, label, size_pt=18, bold=True, color=DARK_GREY)
        if secondary:
            _styled_run(p, "  ·  ", size_pt=18, color=GREY_MID)
            _styled_run(p, secondary, size_pt=15, color=GREY_MID,
                        italic=True)

    # ---------------------------------------------------------------- Brace
    brace_x = list_x + list_w + brace_gap_left
    brace_y = list_y
    brace_h = list_h
    _add_right_brace(slide, brace_x, brace_y, brace_w, brace_h,
                      color=BLUE, width_pt=2.5)

    # --------------------------------------------------------- Signature copy
    sig_x = brace_x + brace_w + sig_gap_left
    sig_y_top = list_y + (list_h / 2) - 70  # vertically centre the block on the brace

    # Lock icon (uses lock.svg, recoloured brand-blue).
    lock_path = icons_dir / "lock.svg"
    lock_size = 44
    if lock_path.exists():
        png = rasterize_recolored(lock_path, 96, "0F6BFF",
                                  stroke_width=icon_stroke)
        slide.shapes.add_picture(
            str(png),
            px(sig_x), px(sig_y_top),
            width=px(lock_size), height=px(lock_size),
        )

    # Eyebrow "SIGNATURE" — small caps, brand-blue, letter-spaced.
    eyebrow_y = sig_y_top + lock_size + 10
    tf = _add_textbox(slide, sig_x, eyebrow_y, sig_w, 24)
    p = tf.paragraphs[0]
    _styled_run(p, "SIGNATURE",
                size_pt=13, bold=True, color=BLUE,
                all_caps=True, letter_spacing="160")

    # Caption "One digest covers all."
    cap_y = eyebrow_y + 28
    tf = _add_textbox(slide, sig_x, cap_y, sig_w, 36)
    p = tf.paragraphs[0]
    _styled_run(p, "One digest covers all.",
                size_pt=18, bold=True, color=BLUE_MID)

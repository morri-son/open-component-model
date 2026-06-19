#!/usr/bin/env python3
"""
Build OCM-Master-Template.potx — a PowerPoint template with the OCM brand
theme and 8 named slide layouts ready for any deck.

Layouts (ordered as they'll appear in the Slide Master):

    1. Hero               — full-bleed banner + title (1 line + gradient noun line) + subtitle + brand row
    2. CTA                — dark backdrop + title + 3-item bullet list + brand row
    3. Content / 3-Column — eyebrow + title + three columns w/ blue rules
    4. Content / Diagram  — eyebrow + title + image area
    5. Content / Tiles    — eyebrow + title + 3x2 grid of grey tiles
    6. Content / 2-Column — eyebrow + title + two-column body
    7. Section Divider    — solid blue background + centered large title
    8. Plain              — eyebrow + title + free body (body y=580, for 2-line titles)
    9. Plain / Compact    — like Plain, but body y=520 (for 1-line titles)

Theme colors (OCM canonical, deck-spec palette confirmed against the Marp
reference render):

    accent1   #0F6BFF (--brand-blue-dark)   — primary brand accent
    accent2   #0A3A99 (--brand-blue-mid)    — secondary deeper blue
    accent3   #5CD6FF (--brand-cyan)        — header cyan
    accent4   #6B7280 (grey-mid)            — secondary text
    accent5   #0A1530 (--brand-blue-night)  — hero/CTA backdrop
    accent6   #F3F4F6 (grey-soft)           — tile fill
    dk1       #000000                       — body text
    lt1       #FFFFFF                       — slide background
    dk2       #0A3A99                       — subheadings (mirror of accent2)
    lt2       #F3F4F6                       — light fill (mirror of accent6)
    hlink     #0F6BFF
    folHlink  #0A3A99

Output: decks/exec-phase1/OCM-Master-Template.potx

Usage:
    .venv/bin/python build_potx.py
"""
from __future__ import annotations

import os
import re
import shutil
import sys
import tempfile
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Emu


# -----------------------------------------------------------------------------
# Paths
# -----------------------------------------------------------------------------

SCRIPT_DIR = Path(__file__).resolve().parent
DECK_DIR = SCRIPT_DIR.parent
MARKETING_DIR = DECK_DIR.parent.parent
ASSETS_DIR = MARKETING_DIR / "assets"
THEME_DIR = DECK_DIR / "theme"

OUTPUT_POTX = DECK_DIR / "OCM-Master.potx"

# -----------------------------------------------------------------------------
# Slide geometry (16:9, 1920x1080)
# -----------------------------------------------------------------------------

SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
PX = 9525  # 1 px in EMU at 96 dpi

def emu(v: float) -> int:
    return int(v * PX)

def px(v: float) -> int:
    return int(v * PX)


# -----------------------------------------------------------------------------
# OCM brand palette (canonical)
# -----------------------------------------------------------------------------

PALETTE = {
    "accent1":  "0F6BFF",   # brand blue dark — primary
    "accent2":  "0A3A99",   # brand blue mid
    "accent3":  "5CD6FF",   # brand cyan
    "accent4":  "6B7280",   # grey mid
    "accent5":  "0A1530",   # brand blue night
    "accent6":  "F3F4F6",   # grey soft
    "dk1":      "000000",
    "lt1":      "FFFFFF",
    "dk2":      "0A3A99",
    "lt2":      "F3F4F6",
    "hlink":    "0F6BFF",
    "folHlink": "0A3A99",
}


# -----------------------------------------------------------------------------
# OOXML namespaces
# -----------------------------------------------------------------------------

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NSMAP = {"a": A, "p": P, "r": R}

def qa(tag): return f"{{{A}}}{tag}"
def qp(tag): return f"{{{P}}}{tag}"
def qr(tag): return f"{{{R}}}{tag}"


# -----------------------------------------------------------------------------
# OOXML helpers
# -----------------------------------------------------------------------------

def element(tag, **attrs):
    """Make an element with the given tag (use qa/qp helpers) and attributes."""
    el = etree.Element(tag)
    for k, v in attrs.items():
        el.set(k, str(v))
    return el


def srgb(hex_color: str):
    """Make <a:srgbClr val="..."/>"""
    el = etree.Element(qa("srgbClr"))
    el.set("val", hex_color)
    return el


def text_pr(size_pt=None, bold=None, color_hex=None, font="Inter",
            all_caps=False, letter_spacing_pt=None):
    """Make <a:rPr ...> with common attributes."""
    rPr = etree.Element(qa("rPr"))
    rPr.set("lang", "en-US")
    if size_pt is not None:
        rPr.set("sz", str(int(size_pt * 100)))
    if bold is not None:
        rPr.set("b", "1" if bold else "0")
    if all_caps:
        rPr.set("cap", "all")
    if letter_spacing_pt is not None:
        rPr.set("spc", str(int(letter_spacing_pt * 100)))
    if color_hex:
        fill = etree.SubElement(rPr, qa("solidFill"))
        fill.append(srgb(color_hex))
    if font:
        latin = etree.SubElement(rPr, qa("latin"))
        latin.set("typeface", font)
    return rPr


def make_textbox(name, idx, x, y, w, h, *, placeholder_type=None,
                 placeholder_idx=None, default_text="",
                 size_pt=18, bold=False, color_hex="000000", font="Inter",
                 all_caps=False, letter_spacing_pt=None,
                 anchor="t", algn=None, line_spacing_pct=None,
                 no_autofit=False):
    """Build a <p:sp> shape that's a placeholder (provide placeholder_type
    like 'title', 'body', 'ctrTitle'…) or a static textbox (placeholder_type=None).

    Critical: paragraph-level styling (size, color, bold, all-caps, no-bullet)
    is set on <a:lstStyle><a:lvl1pPr><a:defRPr> so PowerPoint applies it to
    new slides spawned from the layout. Setting only the run rPr is NOT
    enough — PowerPoint inherits paragraph defaults from the master body
    style otherwise (which includes bullets and the wrong font size).
    """
    sp = etree.Element(qp("sp"))

    # nvSpPr
    nvSpPr = etree.SubElement(sp, qp("nvSpPr"))
    cNvPr = etree.SubElement(nvSpPr, qp("cNvPr"))
    cNvPr.set("id", str(idx))
    cNvPr.set("name", name)
    cNvSpPr = etree.SubElement(nvSpPr, qp("cNvSpPr"))
    etree.SubElement(cNvSpPr, qa("spLocks")).set("noGrp", "1")
    nvPr = etree.SubElement(nvSpPr, qp("nvPr"))
    if placeholder_type is not None:
        ph = etree.SubElement(nvPr, qp("ph"))
        ph.set("type", placeholder_type)
        if placeholder_idx is not None:
            ph.set("idx", str(placeholder_idx))

    # spPr
    spPr = etree.SubElement(sp, qp("spPr"))
    xfrm = etree.SubElement(spPr, qa("xfrm"))
    off = etree.SubElement(xfrm, qa("off"))
    off.set("x", str(emu(x))); off.set("y", str(emu(y)))
    ext = etree.SubElement(xfrm, qa("ext"))
    ext.set("cx", str(emu(w))); ext.set("cy", str(emu(h)))
    prstGeom = etree.SubElement(spPr, qa("prstGeom"))
    prstGeom.set("prst", "rect")
    etree.SubElement(prstGeom, qa("avLst"))

    # txBody
    txBody = etree.SubElement(sp, qp("txBody"))
    bodyPr = etree.SubElement(txBody, qa("bodyPr"))
    bodyPr.set("wrap", "square")
    bodyPr.set("rtlCol", "0")
    bodyPr.set("anchor", anchor)
    bodyPr.set("lIns", "0"); bodyPr.set("tIns", "0")
    bodyPr.set("rIns", "0"); bodyPr.set("bIns", "0")
    if no_autofit:
        # Disable PowerPoint's auto-shrink-to-fit. Without this, long text in
        # column headers shrinks the font and produces inconsistent type sizes
        # across columns when one wraps and another doesn't.
        etree.SubElement(bodyPr, qa("noAutofit"))

    # lstStyle — this is where the real styling lives. PowerPoint inherits
    # these defaults when authors type into the placeholder.
    lstStyle = etree.SubElement(txBody, qa("lstStyle"))
    lvl1 = etree.SubElement(lstStyle, qa("lvl1pPr"))
    lvl1.set("marL", "0")
    lvl1.set("indent", "0")
    # Default to left alignment unless explicitly overridden. The master's
    # titleStyle inherits as centered for title placeholders, so without an
    # explicit algn here every title would render centered.
    lvl1.set("algn", algn or "l")
    # Tighter line spacing (titles wrapping to 2 lines) — must come BEFORE
    # buNone per OOXML schema ordering. line_spacing_pct is a ratio: 0.9 = 90%.
    if line_spacing_pct is not None:
        lnSpc = etree.SubElement(lvl1, qa("lnSpc"))
        spcPct = etree.SubElement(lnSpc, qa("spcPct"))
        spcPct.set("val", str(int(line_spacing_pct * 100000)))
    # No bullet — critical, otherwise master body style adds one.
    etree.SubElement(lvl1, qa("buNone"))
    defRPr = etree.SubElement(lvl1, qa("defRPr"))
    defRPr.set("sz", str(int(size_pt * 100)))
    if bold:
        defRPr.set("b", "1")
    if all_caps:
        defRPr.set("cap", "all")
    if letter_spacing_pt is not None:
        defRPr.set("spc", str(int(letter_spacing_pt * 100)))
    if color_hex:
        fill = etree.SubElement(defRPr, qa("solidFill"))
        fill.append(srgb(color_hex))
    if font:
        latin = etree.SubElement(defRPr, qa("latin"))
        latin.set("typeface", font)

    # Placeholder text (the prompt user sees in Slide Master view; replaced
    # when they type into the placeholder on a new slide).
    p = etree.SubElement(txBody, qa("p"))
    if default_text:
        r = etree.SubElement(p, qa("r"))
        rPr = etree.SubElement(r, qa("rPr"))
        rPr.set("lang", "en-US")
        t = etree.SubElement(r, qa("t"))
        t.text = default_text
    endParaRPr = etree.SubElement(p, qa("endParaRPr"))
    endParaRPr.set("lang", "en-US")
    return sp


def make_rect(name, idx, x, y, w, h, fill_hex):
    """A solid-fill rectangle (no line). Used for column rules and tile chrome."""
    sp = etree.Element(qp("sp"))
    nvSpPr = etree.SubElement(sp, qp("nvSpPr"))
    cNvPr = etree.SubElement(nvSpPr, qp("cNvPr"))
    cNvPr.set("id", str(idx))
    cNvPr.set("name", name)
    etree.SubElement(nvSpPr, qp("cNvSpPr"))
    etree.SubElement(nvSpPr, qp("nvPr"))
    spPr = etree.SubElement(sp, qp("spPr"))
    xfrm = etree.SubElement(spPr, qa("xfrm"))
    off = etree.SubElement(xfrm, qa("off"))
    off.set("x", str(emu(x))); off.set("y", str(emu(y)))
    ext = etree.SubElement(xfrm, qa("ext"))
    ext.set("cx", str(emu(w))); ext.set("cy", str(emu(h)))
    prstGeom = etree.SubElement(spPr, qa("prstGeom"))
    prstGeom.set("prst", "rect")
    etree.SubElement(prstGeom, qa("avLst"))
    fill = etree.SubElement(spPr, qa("solidFill"))
    fill.append(srgb(fill_hex))
    ln = etree.SubElement(spPr, qa("ln"))
    etree.SubElement(ln, qa("noFill"))
    # Empty txBody (required for spec compliance).
    txBody = etree.SubElement(sp, qp("txBody"))
    bp = etree.SubElement(txBody, qa("bodyPr"))
    bp.set("rtlCol", "0"); bp.set("anchor", "ctr")
    etree.SubElement(txBody, qa("lstStyle"))
    p = etree.SubElement(txBody, qa("p"))
    etree.SubElement(p, qa("endParaRPr")).set("lang", "en-US")
    return sp


# -----------------------------------------------------------------------------
# Theme XML — replaces the default Office theme with OCM brand colors
# -----------------------------------------------------------------------------

def build_theme_xml() -> bytes:
    """Build a complete theme1.xml with OCM colors and Aptos as default font."""
    # Use string template — most of theme1.xml is boilerplate we want unchanged.
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:theme xmlns:a="{A}" name="OCM">
  <a:themeElements>
    <a:clrScheme name="OCM">
      <a:dk1><a:srgbClr val="{PALETTE['dk1']}"/></a:dk1>
      <a:lt1><a:srgbClr val="{PALETTE['lt1']}"/></a:lt1>
      <a:dk2><a:srgbClr val="{PALETTE['dk2']}"/></a:dk2>
      <a:lt2><a:srgbClr val="{PALETTE['lt2']}"/></a:lt2>
      <a:accent1><a:srgbClr val="{PALETTE['accent1']}"/></a:accent1>
      <a:accent2><a:srgbClr val="{PALETTE['accent2']}"/></a:accent2>
      <a:accent3><a:srgbClr val="{PALETTE['accent3']}"/></a:accent3>
      <a:accent4><a:srgbClr val="{PALETTE['accent4']}"/></a:accent4>
      <a:accent5><a:srgbClr val="{PALETTE['accent5']}"/></a:accent5>
      <a:accent6><a:srgbClr val="{PALETTE['accent6']}"/></a:accent6>
      <a:hlink><a:srgbClr val="{PALETTE['hlink']}"/></a:hlink>
      <a:folHlink><a:srgbClr val="{PALETTE['folHlink']}"/></a:folHlink>
    </a:clrScheme>
    <a:fontScheme name="OCM">
      <a:majorFont>
        <a:latin typeface="Inter Display"/>
        <a:ea typeface=""/>
        <a:cs typeface=""/>
      </a:majorFont>
      <a:minorFont>
        <a:latin typeface="Inter"/>
        <a:ea typeface=""/>
        <a:cs typeface=""/>
      </a:minorFont>
    </a:fontScheme>
    <a:fmtScheme name="Office">
      <a:fillStyleLst>
        <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
        <a:gradFill rotWithShape="1">
          <a:gsLst>
            <a:gs pos="0"><a:schemeClr val="phClr"><a:tint val="50000"/><a:satMod val="300000"/></a:schemeClr></a:gs>
            <a:gs pos="35000"><a:schemeClr val="phClr"><a:tint val="37000"/><a:satMod val="300000"/></a:schemeClr></a:gs>
            <a:gs pos="100000"><a:schemeClr val="phClr"><a:tint val="15000"/><a:satMod val="350000"/></a:schemeClr></a:gs>
          </a:gsLst>
          <a:lin ang="16200000" scaled="1"/>
        </a:gradFill>
        <a:gradFill rotWithShape="1">
          <a:gsLst>
            <a:gs pos="0"><a:schemeClr val="phClr"><a:shade val="51000"/><a:satMod val="130000"/></a:schemeClr></a:gs>
            <a:gs pos="80000"><a:schemeClr val="phClr"><a:shade val="93000"/><a:satMod val="130000"/></a:schemeClr></a:gs>
            <a:gs pos="100000"><a:schemeClr val="phClr"><a:shade val="94000"/><a:satMod val="135000"/></a:schemeClr></a:gs>
          </a:gsLst>
          <a:lin ang="16200000" scaled="0"/>
        </a:gradFill>
      </a:fillStyleLst>
      <a:lnStyleLst>
        <a:ln w="9525" cap="flat" cmpd="sng" algn="ctr"><a:solidFill><a:schemeClr val="phClr"><a:shade val="95000"/><a:satMod val="105000"/></a:schemeClr></a:solidFill><a:prstDash val="solid"/></a:ln>
        <a:ln w="25400" cap="flat" cmpd="sng" algn="ctr"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill><a:prstDash val="solid"/></a:ln>
        <a:ln w="38100" cap="flat" cmpd="sng" algn="ctr"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill><a:prstDash val="solid"/></a:ln>
      </a:lnStyleLst>
      <a:effectStyleLst>
        <a:effectStyle><a:effectLst/></a:effectStyle>
        <a:effectStyle><a:effectLst/></a:effectStyle>
        <a:effectStyle><a:effectLst><a:outerShdw blurRad="40000" dist="20000" dir="5400000" rotWithShape="0"><a:srgbClr val="000000"><a:alpha val="38000"/></a:srgbClr></a:outerShdw></a:effectLst></a:effectStyle>
      </a:effectStyleLst>
      <a:bgFillStyleLst>
        <a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
        <a:solidFill><a:schemeClr val="phClr"><a:tint val="95000"/><a:satMod val="170000"/></a:schemeClr></a:solidFill>
        <a:gradFill rotWithShape="1"><a:gsLst><a:gs pos="0"><a:schemeClr val="phClr"><a:tint val="93000"/><a:satMod val="150000"/><a:shade val="98000"/></a:schemeClr></a:gs><a:gs pos="50000"><a:schemeClr val="phClr"><a:tint val="98000"/><a:satMod val="130000"/><a:shade val="90000"/></a:schemeClr></a:gs><a:gs pos="100000"><a:schemeClr val="phClr"><a:shade val="63000"/><a:satMod val="120000"/></a:schemeClr></a:gs></a:gsLst><a:lin ang="5400000" scaled="0"/></a:gradFill>
      </a:bgFillStyleLst>
    </a:fmtScheme>
  </a:themeElements>
  <a:objectDefaults/>
  <a:extraClrSchemeLst/>
</a:theme>'''.encode("utf-8")


# -----------------------------------------------------------------------------
# Slide-layout XML builders (one per layout)
# -----------------------------------------------------------------------------

LAYOUT_HEADER = '''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:a="{A}" xmlns:r="{R}" xmlns:p="{P}"
             type="{type}" preserve="1" userDrawn="1">
  <p:cSld name="{name}">'''.format(A=A, R=R, P=P, type="{type}", name="{name}")


def wrap_layout(layout_type: str, layout_name: str, sp_tree_xml: str,
                 bg_hex: str = "FFFFFF") -> bytes:
    """Wrap a sequence of <p:sp>... in the boilerplate of a slideLayout."""
    raw = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:a="{A}" xmlns:r="{R}" xmlns:p="{P}"
             type="{layout_type}" preserve="1" userDrawn="1">
  <p:cSld name="{layout_name}">
    <p:bg>
      <p:bgPr>
        <a:solidFill><a:srgbClr val="{bg_hex}"/></a:solidFill>
        <a:effectLst/>
      </p:bgPr>
    </p:bg>
    <p:spTree>
      <p:nvGrpSpPr>
        <p:cNvPr id="1" name=""/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="0" y="0"/>
          <a:ext cx="0" cy="0"/>
          <a:chOff x="0" y="0"/>
          <a:chExt cx="0" cy="0"/>
        </a:xfrm>
      </p:grpSpPr>
      {sp_tree_xml}
    </p:spTree>
  </p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sldLayout>'''.encode("utf-8")
    # Canonicalize: collapse the per-<p:sp> xmlns redeclarations that
    # lxml's per-element tostring leaves embedded in sp_tree_xml. Without
    # this step a slideLayout5 hits 84 xmlns: declarations and PowerPoint
    # flags it for repair.
    return canonicalize_xml(raw)


def serialize(el) -> str:
    # Re-parse the element under a stable nsmap so lxml uses canonical p:/a:/r:
    # prefixes instead of inventing ns0:/ns1:/... per top-level shape. PowerPoint
    # parses both fine, but the auto-prefix variant trips its consistency check
    # and triggers the "file needs repair" dialog on every open.
    raw = etree.tostring(el, pretty_print=False)
    reparsed = etree.fromstring(raw)
    canonical = etree.Element(reparsed.tag, nsmap=NSMAP)
    for k, v in reparsed.attrib.items():
        canonical.set(k, v)
    for child in reparsed:
        canonical.append(child)
    canonical.text = reparsed.text
    return etree.tostring(canonical, pretty_print=False).decode("utf-8")


def shapes_xml(*shapes) -> str:
    return "\n".join(serialize(sh) for sh in shapes)


def canonicalize_xml(xml_bytes: bytes) -> bytes:
    """Re-parse a complete OOXML document and re-emit it with namespace
    declarations only at the document root.

    Concatenating string-serialised lxml elements (as wrap_layout does) means
    every <p:sp> child carries its own xmlns:p / xmlns:a / xmlns:r declarations.
    PowerPoint-for-Mac's parser flags this as malformed and triggers the
    "file is damaged, repair?" dialog.

    lxml preserves redundant nsmaps it sees on parse, even when they're
    identical to the parent's. We walk the tree manually and rebuild every
    element with no nsmap of its own — children inherit the root's
    declarations and re-serialization emits a single root-level xmlns block."""
    src = etree.fromstring(xml_bytes)

    def _clone(node):
        if not isinstance(node.tag, str):
            # Comment / PI — pass through.
            return node
        new = etree.Element(node.tag)
        for k, v in node.attrib.items():
            new.set(k, v)
        new.text = node.text
        new.tail = node.tail
        for child in node:
            new.append(_clone(child))
        return new

    new_root = etree.Element(src.tag, nsmap=NSMAP)
    for k, v in src.attrib.items():
        new_root.set(k, v)
    new_root.text = src.text
    for child in src:
        new_root.append(_clone(child))

    return (b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
            + etree.tostring(new_root, pretty_print=False))


# -- Layout 1: Hero -----------------------------------------------------------

def layout_hero() -> bytes:
    """Hero layout: full-bleed banner image (locked) + title/subtitle/orgline
    placeholders + brand row (locked).

    Title is split across two placeholders (line 1 and line 2) so the deck
    can apply a gradient text fill to just the operative noun on line 2.
    Both placeholders use type='body' with explicit algn='l' to avoid
    PowerPoint's automatic centering for ctrTitle/title placeholders.

    Vertical layout (slide 1080 tall, banner is full-bleed background):
      - Title line 1   y=160  h=200   (single row, 115pt)
      - Title line 2   y=380  h=160   (gradient line, 115pt)
      - Subtitle       y=600  h=120
      - Org line       y=860  h=60

    Use this layout when the hero title fits cleanly on TWO lines at 115pt.
    Convention: hero titles are at most TWO lines, never three. If a title
    needs three lines, the title is too long — shorten it and push the rest
    into the subtitle.
    """
    shapes = []
    shapes.append(make_textbox(
        "Hero Title Line 1", 10, 96, 180, 1700, 160,
        placeholder_type="title", placeholder_idx=1,
        default_text="Secure Delivery for",
        size_pt=115, bold=True, color_hex="FFFFFF", algn="l",
        font="Inter Display",
        no_autofit=True,
    ))
    shapes.append(make_textbox(
        "Hero Title Line 2", 11, 96, 345, 1700, 160,
        placeholder_type="body", placeholder_idx=2,
        default_text="Sovereign Clouds",
        size_pt=115, bold=True, color_hex="5CD6FF", algn="l",
        font="Inter Display",
        no_autofit=True,
    ))
    shapes.append(make_textbox(
        "Hero Subtitle", 12, 96, 560, 1700, 90,
        placeholder_type="body", placeholder_idx=3,
        default_text="Subtitle — one sentence describing the deck.",
        size_pt=30, color_hex="5CD6FF", algn="l",
        line_spacing_pct=1.15,
        no_autofit=True,
    ))
    shapes.append(make_textbox(
        "Hero Org Line", 13, 96, 690, 1700, 60,
        placeholder_type="body", placeholder_idx=4,
        default_text="Open Component Model — open source, NeoNephos Foundation.",
        size_pt=24, color_hex="FFFFFF", algn="l",
        no_autofit=True,
    ))
    return wrap_layout("title", "Hero", shapes_xml(*shapes), bg_hex="0A1530")


# -- Layout 1b: (removed) Hero / 3-Line --------------------------------------
#
# A three-line hero layout was prototyped in 2026-06-17 to support an 11-word
# cold-room opener title that wouldn't fit in the standard two-line Hero. It
# was abandoned the same day: the visual mass of three 115pt title lines on
# slide 1 reads as overwhelming, not crisp. Hero convention going forward:
# **maximum two lines, never three.** If a title needs three lines, the
# title is too long — shorten the title, push the rest into the subtitle.


# -- Layout 2: CTA ------------------------------------------------------------

def layout_cta() -> bytes:
    shapes = []
    shapes.append(make_textbox(
        "CTA Title", 10, 96, 320, 1700, 100,
        placeholder_type="title", placeholder_idx=1,
        default_text="Call to action",
        size_pt=56, bold=True, color_hex="FFFFFF",
        font="Inter Display",
    ))
    shapes.append(make_textbox(
        "CTA Body", 11, 96, 460, 1700, 400,
        placeholder_type="body", placeholder_idx=2,
        default_text="Three-line CTA bullets go here.",
        size_pt=22, color_hex="FFFFFF",
    ))
    return wrap_layout("obj", "CTA", shapes_xml(*shapes), bg_hex="0A1530")


# -- Layout 3: Content / 3-Column --------------------------------------------

def layout_three_column() -> bytes:
    shapes = []
    # Eyebrow
    shapes.append(make_textbox(
        "Eyebrow", 10, 120, 255, 1680, 48,
        placeholder_type="body", placeholder_idx=1,
        default_text="EYEBROW",
        size_pt=28, bold=True, color_hex="0F6BFF",
        all_caps=True, letter_spacing_pt=1.4,
    ))
    # Title
    shapes.append(make_textbox(
        "Title", 11, 120, 308, 1680, 200,
        placeholder_type="title", placeholder_idx=2,
        default_text="Section title goes here.",
        size_pt=64, bold=True, color_hex="000000",
        font="Inter Display", line_spacing_pct=0.9,
    ))
    # Three columns
    margin_x = 120
    gutter = 56
    inner_w = SLIDE_W_PX - 2 * margin_x
    col_w = (inner_w - 2 * gutter) // 3
    col_y = 520
    next_id = 12
    for i in range(3):
        cx = margin_x + i * (col_w + gutter)
        # 4px Brand Blue rule
        shapes.append(make_rect(f"Col{i+1} Rule", next_id,
                                  cx, col_y, col_w, 4, "0F6BFF"))
        next_id += 1
        # Header. no_autofit ensures all three column headers render at the
        # same fixed size — long headers wrap rather than auto-shrinking, so
        # neighbouring columns stay visually consistent.
        shapes.append(make_textbox(
            f"Col{i+1} Header", next_id, cx, col_y + 16, col_w, 56,
            placeholder_type="body", placeholder_idx=10 + i * 2,
            default_text=f"COLUMN {i+1} HEADER",
            size_pt=20, bold=True, color_hex="0F6BFF",
            all_caps=True, letter_spacing_pt=1.3,
            no_autofit=True,
        ))
        next_id += 1
        # Body. Sits below a 1-line header (16 top padding + 56 box = 72) with
        # a 12px gap before the body. Convention: column headers stay on
        # ONE line; 2-line headers force a tradeoff between eyebrow-header
        # and header-body spacing that no single layout can satisfy.
        shapes.append(make_textbox(
            f"Col{i+1} Body", next_id, cx, col_y + 84, col_w, 460,
            placeholder_type="body", placeholder_idx=11 + i * 2,
            default_text=f"Column {i+1} body. Replace with one or two short sentences.",
            size_pt=22, color_hex="000000",
        ))
        next_id += 1
    # Footer
    shapes.append(_footer_shape(next_id))
    return wrap_layout("obj", "Content / 3-Column", shapes_xml(*shapes))


# -- Layout 4: Content / Diagram ---------------------------------------------

def layout_diagram() -> bytes:
    """Diagram-first layout. Eyebrow + title sit high (variant C from the
    layout test) so the diagram gets the bulk of the slide. Title slot is
    short (80px ≈ one line at 64pt) — diagram-slide titles are short by
    convention ("One identity, every boundary.", "Trust travels with the
    component.") so we don't need the 2-line slot the text-heavy Plain
    layouts reserve. Slimming it lets the diagram start 40px earlier.

    Geometry (1920×1080 canvas):
      - Eyebrow   y= 75  h= 48   (28pt, brand-blue, ALL-CAPS)
      - Title     y=128  h= 80   (64pt, 1 line; 2-line titles overflow into the diagram)
      - Diagram   y=240  h=780   (1800 wide, 60px gutter both sides)
    """
    shapes = [
        make_textbox("Eyebrow", 10, 120, 75, 1680, 48,
                     placeholder_type="body", placeholder_idx=1,
                     default_text="EYEBROW", size_pt=28, bold=True,
                     color_hex="0F6BFF", all_caps=True, letter_spacing_pt=1.4),
        make_textbox("Title", 11, 120, 128, 1680, 80,
                     placeholder_type="title", placeholder_idx=2,
                     default_text="Section title goes here.",
                     size_pt=64, bold=True, color_hex="000000",
                     font="Inter Display", line_spacing_pct=0.9),
        _picture_placeholder("Diagram", 12, 60, 240, 1800, 780, ph_idx=10),
        _footer_shape(13),
    ]
    return wrap_layout("obj", "Content / Diagram", shapes_xml(*shapes))


# -- Layout 5: Content / Tiles -----------------------------------------------

def layout_tiles() -> bytes:
    shapes = [
        make_textbox("Eyebrow", 10, 120, 255, 1680, 48,
                     placeholder_type="body", placeholder_idx=1,
                     default_text="EYEBROW", size_pt=28, bold=True,
                     color_hex="0F6BFF", all_caps=True, letter_spacing_pt=1.4),
        make_textbox("Title", 11, 120, 308, 1680, 200,
                     placeholder_type="title", placeholder_idx=2,
                     default_text="Section title goes here.",
                     size_pt=64, bold=True, color_hex="000000",
                     font="Inter Display", line_spacing_pct=0.9),
    ]
    # 3x2 grid
    x0, y0 = 120, 520
    tile_w, tile_h = 544, 230
    gutter = 24
    next_id = 12
    for i in range(6):
        col = i % 3
        row = i // 3
        x = x0 + col * (tile_w + gutter)
        y = y0 + row * (tile_h + gutter)
        # Tile background (grey-soft)
        shapes.append(make_rect(f"Tile{i+1} Bg", next_id, x, y, tile_w, tile_h, "F3F4F6"))
        next_id += 1
        # Top rule (brand blue)
        shapes.append(make_rect(f"Tile{i+1} Rule", next_id, x, y, tile_w, 3, "0F6BFF"))
        next_id += 1
        # Tile label placeholder — sits to the RIGHT of the icon (icon-on-left
        # variant). Icon occupies x..x+48 of the tile padding zone (added
        # inline at build time, not by the layout). Label starts at x+24+48+16
        # = x+88, vertically centred against the 48-tall icon row at y+24.
        # anchor="ctr" puts the label text mid-row so the icon and the label
        # baseline align visually instead of the label sitting on the icon's top.
        shapes.append(make_textbox(
            f"Tile{i+1} Label", next_id, x + 88, y + 24, tile_w - 88 - 24, 48,
            placeholder_type="body", placeholder_idx=20 + i * 2,
            default_text=f"Tile {i+1} label",
            size_pt=18, bold=True, color_hex="0F6BFF",
            anchor="ctr",
        ))
        next_id += 1
        # Tile body placeholder — full tile width below the icon/label row.
        shapes.append(make_textbox(
            f"Tile{i+1} Body", next_id, x + 24, y + 96, tile_w - 48, tile_h - 116,
            placeholder_type="body", placeholder_idx=21 + i * 2,
            default_text="One short sentence describing this tile.",
            size_pt=18, color_hex="000000",
        ))
        next_id += 1
    shapes.append(_footer_shape(next_id))
    return wrap_layout("obj", "Content / Tiles", shapes_xml(*shapes))


# -- Layout 6: Content / 2-Column --------------------------------------------

def layout_two_column() -> bytes:
    shapes = [
        make_textbox("Eyebrow", 10, 120, 255, 1680, 48,
                     placeholder_type="body", placeholder_idx=1,
                     default_text="EYEBROW", size_pt=28, bold=True,
                     color_hex="0F6BFF", all_caps=True, letter_spacing_pt=1.4),
        make_textbox("Title", 11, 120, 308, 1680, 200,
                     placeholder_type="title", placeholder_idx=2,
                     default_text="Section title goes here.",
                     size_pt=64, bold=True, color_hex="000000",
                     font="Inter Display", line_spacing_pct=0.9),
        make_textbox("Left Body", 12, 120, 520, 820, 460,
                     placeholder_type="body", placeholder_idx=10,
                     default_text="Left column body.",
                     size_pt=22, color_hex="000000"),
        make_textbox("Right Body", 13, 980, 520, 820, 460,
                     placeholder_type="body", placeholder_idx=11,
                     default_text="Right column body.",
                     size_pt=22, color_hex="000000"),
        _footer_shape(14),
    ]
    return wrap_layout("twoObj", "Content / 2-Column", shapes_xml(*shapes))


# -- Layout 7: Section Divider ------------------------------------------------

def layout_section() -> bytes:
    shapes = [
        make_textbox("Section Title", 10, 80, 460, 1760, 200,
                     placeholder_type="title", placeholder_idx=1,
                     default_text="Section",
                     size_pt=72, bold=True, color_hex="FFFFFF",
                     anchor="ctr", algn="ctr",
                     font="Inter Display"),
    ]
    return wrap_layout("secHead", "Section Divider", shapes_xml(*shapes),
                        bg_hex="0F6BFF")


# -- Layout 8: Plain ----------------------------------------------------------

def layout_plain() -> bytes:
    """Plain body, generous gap between title and body — for slides whose
    title wraps to 2 lines (e.g. SCAN — Compliance-native)."""
    shapes = [
        make_textbox("Eyebrow", 10, 120, 255, 1680, 48,
                     placeholder_type="body", placeholder_idx=1,
                     default_text="EYEBROW", size_pt=28, bold=True,
                     color_hex="0F6BFF", all_caps=True, letter_spacing_pt=1.4),
        make_textbox("Title", 11, 120, 308, 1680, 200,
                     placeholder_type="title", placeholder_idx=2,
                     default_text="Section title goes here.",
                     size_pt=64, bold=True, color_hex="000000",
                     font="Inter Display", line_spacing_pct=0.9),
        make_textbox("Body", 12, 120, 580, 1680, 400,
                     placeholder_type="body", placeholder_idx=10,
                     default_text="Body text.",
                     size_pt=22, color_hex="000000"),
        _footer_shape(13),
    ]
    return wrap_layout("obj", "Plain", shapes_xml(*shapes))


def layout_plain_compact() -> bytes:
    """Plain body sitting close under the title — for slides with 1-line
    titles (e.g. THE SHIFT, SOVEREIGN-READY) where the y=580 layout leaves
    too much air between title and content."""
    shapes = [
        make_textbox("Eyebrow", 10, 120, 255, 1680, 48,
                     placeholder_type="body", placeholder_idx=1,
                     default_text="EYEBROW", size_pt=28, bold=True,
                     color_hex="0F6BFF", all_caps=True, letter_spacing_pt=1.4),
        make_textbox("Title", 11, 120, 308, 1680, 200,
                     placeholder_type="title", placeholder_idx=2,
                     default_text="Section title goes here.",
                     size_pt=64, bold=True, color_hex="000000",
                     font="Inter Display", line_spacing_pct=0.9),
        make_textbox("Body", 12, 120, 520, 1680, 460,
                     placeholder_type="body", placeholder_idx=10,
                     default_text="Body text.",
                     size_pt=22, color_hex="000000"),
        _footer_shape(13),
    ]
    return wrap_layout("obj", "Plain / Compact", shapes_xml(*shapes))


# Helpers shared across layouts
def _footer_shape(idx):
    """Footer text — locked, ALL-CAPS, Grey Mid."""
    return make_textbox(
        "Footer", idx, 120, SLIDE_H_PX - 32, SLIDE_W_PX - 240, 24,
        default_text="OPEN COMPONENT MODEL · OCM.SOFTWARE",
        size_pt=9, color_hex="6B7280", all_caps=True, letter_spacing_pt=0.5,
    )


def _picture_placeholder(name, idx, x, y, w, h, *, ph_idx=10):
    """A picture placeholder — the author can drop an image into it."""
    sp = etree.Element(qp("sp"))
    nvSpPr = etree.SubElement(sp, qp("nvSpPr"))
    cNvPr = etree.SubElement(nvSpPr, qp("cNvPr"))
    cNvPr.set("id", str(idx))
    cNvPr.set("name", name)
    cNvSpPr = etree.SubElement(nvSpPr, qp("cNvSpPr"))
    etree.SubElement(cNvSpPr, qa("spLocks")).set("noGrp", "1")
    nvPr = etree.SubElement(nvSpPr, qp("nvPr"))
    ph = etree.SubElement(nvPr, qp("ph"))
    ph.set("type", "pic")
    ph.set("idx", str(ph_idx))
    spPr = etree.SubElement(sp, qp("spPr"))
    xfrm = etree.SubElement(spPr, qa("xfrm"))
    off = etree.SubElement(xfrm, qa("off"))
    off.set("x", str(emu(x))); off.set("y", str(emu(y)))
    ext = etree.SubElement(xfrm, qa("ext"))
    ext.set("cx", str(emu(w))); ext.set("cy", str(emu(h)))
    prstGeom = etree.SubElement(spPr, qa("prstGeom"))
    prstGeom.set("prst", "rect")
    etree.SubElement(prstGeom, qa("avLst"))
    txBody = etree.SubElement(sp, qp("txBody"))
    bp = etree.SubElement(txBody, qa("bodyPr"))
    bp.set("rtlCol", "0"); bp.set("anchor", "ctr")
    etree.SubElement(txBody, qa("lstStyle"))
    p = etree.SubElement(txBody, qa("p"))
    etree.SubElement(p, qa("endParaRPr")).set("lang", "en-US")
    return sp


_picture_placeholder.__name__ = "_picture_placeholder"


# -----------------------------------------------------------------------------
# Embedded fonts (Inter)
# -----------------------------------------------------------------------------

# Resolve user font dir: TTFs were placed under ~/Library/Fonts/Inter-*.ttf
# during setup. Each <p:embeddedFont> entry maps a typeface name to up to
# four style variants (regular, bold, italic, boldItalic). Inter Display is
# the headline-tuned cut used for hero/section titles via the theme's
# majorFont.
USER_FONT_DIR = Path.home() / "Library" / "Fonts"

EMBEDDED_FONTS = [
    # (typeface name, {style: TTF filename})
    ("Inter", {
        "regular":     "Inter-Regular.ttf",
        "bold":        "Inter-Bold.ttf",
        "italic":      "Inter-Italic.ttf",
        "boldItalic":  "Inter-BoldItalic.ttf",
    }),
    ("Inter Display", {
        "regular":    "InterDisplay-Regular.ttf",
        "bold":       "InterDisplay-Bold.ttf",
        "italic":     "InterDisplay-Italic.ttf",
        "boldItalic": "InterDisplay-BoldItalic.ttf",
    }),
]


def embed_fonts(archive: dict[str, bytes]) -> None:
    """Embed Inter / Inter Display TTFs into the .potx so PowerPoint renders
    them on every machine, regardless of local installation. Modifies
    `archive` in place: adds ppt/fonts/font*.fntdata parts, registers them
    in presentation.xml.rels, and inserts <p:embeddedFontLst> into
    presentation.xml.

    The font byte stream stored in font*.fntdata is the raw TTF/OTF — no
    obfuscation, no header, no compression. PowerPoint Mac/Windows both
    accept this; the older "ObfuscatedFont" encoding is an Office-2007 era
    quirk that's optional and skipped here for simplicity.
    """
    if not USER_FONT_DIR.exists():
        print(f"[embed_fonts] {USER_FONT_DIR} not found — skipping embed.")
        return

    # Flatten the (font, style) → file mapping into a numbered list of parts.
    parts = []  # [(part_path, ttf_path, font_idx, style)]
    for font_idx, (typeface, styles) in enumerate(EMBEDDED_FONTS, start=1):
        for style, filename in styles.items():
            ttf_path = USER_FONT_DIR / filename
            if not ttf_path.exists():
                print(f"[embed_fonts] missing {ttf_path} — skipping {typeface} {style}")
                continue
            part_idx = len(parts) + 1
            parts.append({
                "part": f"ppt/fonts/font{part_idx}.fntdata",
                "ttf":  ttf_path,
                "font_idx": font_idx,
                "typeface": typeface,
                "style": style,
                "rid": None,  # filled in below once we know rId numbers
            })

    if not parts:
        print("[embed_fonts] no font files found — skipping.")
        return

    # 1. Add the font part data to the archive.
    for p in parts:
        archive[p["part"]] = p["ttf"].read_bytes()

    # 2. Append relationships to presentation.xml.rels — pick rId numbers
    #    above whatever's already in the file so we don't collide.
    rels_key = "ppt/_rels/presentation.xml.rels"
    rels_xml = archive[rels_key].decode("utf-8")
    existing_ids = re.findall(r'Id="rId(\d+)"', rels_xml)
    next_rid = max((int(i) for i in existing_ids), default=0) + 1
    rel_inserts = []
    for p in parts:
        p["rid"] = f"rId{next_rid}"
        rel_inserts.append(
            f'<Relationship Id="{p["rid"]}" '
            f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/font" '
            f'Target="fonts/font{parts.index(p) + 1}.fntdata"/>'
        )
        next_rid += 1
    rels_xml = rels_xml.replace(
        "</Relationships>",
        "".join(rel_inserts) + "</Relationships>",
    )
    archive[rels_key] = rels_xml.encode("utf-8")

    # 3. Build <p:embeddedFontLst> grouping parts by typeface.
    by_font = {}
    for p in parts:
        by_font.setdefault(p["font_idx"], {"typeface": p["typeface"],
                                              "styles": {}})
        by_font[p["font_idx"]]["styles"][p["style"]] = p["rid"]

    embedded_xml = "<p:embeddedFontLst>"
    for font_idx in sorted(by_font.keys()):
        info = by_font[font_idx]
        embedded_xml += "<p:embeddedFont>"
        embedded_xml += (
            f'<p:font typeface="{info["typeface"]}" '
            f'panose="020F0502020204030204" pitchFamily="34" charset="0"/>'
        )
        for style in ("regular", "bold", "italic", "boldItalic"):
            if style in info["styles"]:
                embedded_xml += (
                    f'<p:{style} r:id="{info["styles"][style]}"/>'
                )
        embedded_xml += "</p:embeddedFont>"
    embedded_xml += "</p:embeddedFontLst>"

    # 4. Inject embeddedFontLst into presentation.xml. The schema places it
    #    just before <p:defaultTextStyle>; if that element isn't present we
    #    drop it just before the closing </p:presentation>.
    pres_xml = archive["ppt/presentation.xml"].decode("utf-8")
    # Remove any prior embeddedFontLst (idempotent rebuild).
    pres_xml = re.sub(r'<p:embeddedFontLst>.*?</p:embeddedFontLst>',
                       "", pres_xml, flags=re.DOTALL)
    if "<p:defaultTextStyle" in pres_xml:
        pres_xml = pres_xml.replace("<p:defaultTextStyle",
                                      embedded_xml + "<p:defaultTextStyle",
                                      1)
    else:
        pres_xml = pres_xml.replace("</p:presentation>",
                                      embedded_xml + "</p:presentation>", 1)
    archive["ppt/presentation.xml"] = pres_xml.encode("utf-8")

    print(f"[embed_fonts] embedded {len(parts)} font parts "
          f"({sum(len(v['styles']) for v in by_font.values())} styles "
          f"across {len(by_font)} typefaces)")


# -----------------------------------------------------------------------------
# Build pipeline
# -----------------------------------------------------------------------------

LAYOUTS = [
    ("Hero",                "title",   layout_hero),
    ("CTA",                 "obj",     layout_cta),
    ("Content / 3-Column",  "obj",     layout_three_column),
    ("Content / Diagram",   "obj",     layout_diagram),
    ("Content / Tiles",     "obj",     layout_tiles),
    ("Content / 2-Column",  "twoObj",  layout_two_column),
    ("Section Divider",     "secHead", layout_section),
    ("Plain",               "obj",     layout_plain),
    ("Plain / Compact",     "obj",     layout_plain_compact),
]


def build():
    # 1. Start from a blank Presentation to get a valid skeleton.
    prs = Presentation()
    prs.slide_width = Emu(emu(SLIDE_W_PX))
    prs.slide_height = Emu(emu(SLIDE_H_PX))
    tmp_pptx = SCRIPT_DIR / "_raster" / "_potx_skeleton.pptx"
    tmp_pptx.parent.mkdir(exist_ok=True)
    prs.save(str(tmp_pptx))

    # 2. Open the resulting zip and surgically replace theme + layouts +
    #    update content types.
    out_archive_data = {}
    with zipfile.ZipFile(tmp_pptx, "r") as src:
        for name in src.namelist():
            out_archive_data[name] = src.read(name)

    # 2a. Replace theme1.xml
    out_archive_data["ppt/theme/theme1.xml"] = build_theme_xml()

    # 2b. Drop existing layouts + their rels, write our 8.
    layout_paths_to_drop = [n for n in out_archive_data
                              if n.startswith("ppt/slideLayouts/")]
    for p in layout_paths_to_drop:
        del out_archive_data[p]
    layout_files = []
    for i, (name, layout_type, builder) in enumerate(LAYOUTS, 1):
        xml_bytes = builder()
        layout_files.append((i, name, layout_type, xml_bytes))
        out_archive_data[f"ppt/slideLayouts/slideLayout{i}.xml"] = xml_bytes
        out_archive_data[f"ppt/slideLayouts/_rels/slideLayout{i}.xml.rels"] = (
            f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/>
</Relationships>'''.encode("utf-8")
        )

    # 2c. Update slideMaster1.xml.rels to point at the 8 layouts (not 11).
    master_rels_xml = ('''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>'''
        + "".join(
            f'\n<Relationship Id="rId{i+1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout{i}.xml"/>'
            for i in range(1, len(LAYOUTS) + 1)
        )
        + "\n</Relationships>").encode("utf-8")
    out_archive_data["ppt/slideMasters/_rels/slideMaster1.xml.rels"] = master_rels_xml

    # 2d. Update slideMaster1.xml — replace the <p:sldLayoutIdLst> contents to
    #    reference our 8 layouts via rId2..rId9 (rId1 = theme).
    master_xml = out_archive_data["ppt/slideMasters/slideMaster1.xml"].decode("utf-8")
    new_id_list = '<p:sldLayoutIdLst>' + "".join(
        f'<p:sldLayoutId id="{2147483649 + i}" r:id="rId{i + 2}"/>'
        for i in range(len(LAYOUTS))
    ) + '</p:sldLayoutIdLst>'
    import re
    master_xml = re.sub(r'<p:sldLayoutIdLst>.*?</p:sldLayoutIdLst>',
                          new_id_list, master_xml, count=1, flags=re.DOTALL)
    out_archive_data["ppt/slideMasters/slideMaster1.xml"] = master_xml.encode("utf-8")

    # 2d-bis. Fix presentation.xml — set slide size to 1920x1080 (16:9), and
    #    ensure no <p:sldIdLst> (templates have no slides).
    pres_xml = out_archive_data["ppt/presentation.xml"].decode("utf-8")
    # Replace whatever sldSz is there with the correct widescreen dims.
    pres_xml = re.sub(
        r'<p:sldSz[^/]*/>',
        f'<p:sldSz cx="{emu(SLIDE_W_PX)}" cy="{emu(SLIDE_H_PX)}"/>',
        pres_xml, count=1,
    )
    # Drop any <p:sldIdLst>...</p:sldIdLst> if it slipped in.
    pres_xml = re.sub(r'<p:sldIdLst>.*?</p:sldIdLst>', '', pres_xml,
                       flags=re.DOTALL)
    out_archive_data["ppt/presentation.xml"] = pres_xml.encode("utf-8")

    # 2e. Rewrite [Content_Types].xml: point the layout overrides at our 8,
    #    and change presentation.xml's content type to potx.
    content_types_xml = build_content_types_xml(len(LAYOUTS))
    out_archive_data["[Content_Types].xml"] = content_types_xml.encode("utf-8")

    # 2f. Strip macOS-generated printerSettings. python-pptx on macOS embeds an
    #    Apple plist (instead of a Windows DEVMODE) in printerSettings1.bin —
    #    PowerPoint considers this format corrupt and triggers a repair dialog
    #    on every open. Verified by diff against PowerPoint's repaired file:
    #    repair removes printerSettings1.bin, its Relationship, and the
    #    Default Extension="bin" content-type — and only that. Removing them
    #    here pre-emptively eliminates the repair dialog.
    for key in list(out_archive_data.keys()):
        if "printerSettings" in key:
            del out_archive_data[key]
    prs_rels_key = "ppt/_rels/presentation.xml.rels"
    if prs_rels_key in out_archive_data:
        rels_xml = out_archive_data[prs_rels_key].decode("utf-8")
        rels_xml = re.sub(
            r'<Relationship[^/]*/officeDocument/2006/relationships/printerSettings[^/]*/>', "",
            rels_xml,
        )
        out_archive_data[prs_rels_key] = rels_xml.encode("utf-8")

    # 2g. Embed Inter font files so the deck renders identically on machines
    #     where Inter is not installed (Windows users who only have Calibri,
    #     fresh macOS without manual font install). PowerPoint stores embedded
    #     fonts as ppt/fonts/font*.fntdata (raw TTF/OTF bytes — no special
    #     wrapper), referenced from <p:embeddedFontLst> in presentation.xml
    #     with one <p:embeddedFont> per typeface and per-style relationships
    #     to the part. The Default Extension="fntdata" content-type maps the
    #     part to application/x-fontdata.
    embed_fonts(out_archive_data)

    # 3. Write the final .potx archive.
    OUTPUT_POTX.unlink(missing_ok=True)
    with zipfile.ZipFile(OUTPUT_POTX, "w", zipfile.ZIP_DEFLATED) as out:
        for name, data in out_archive_data.items():
            out.writestr(name, data)
    print(f"Wrote {OUTPUT_POTX}")


def build_content_types_xml(n_layouts: int) -> str:
    """Build [Content_Types].xml with proper potx content type and N layouts."""
    layout_overrides = "".join(
        f'<Override PartName="/ppt/slideLayouts/slideLayout{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>'
        for i in range(1, n_layouts + 1)
    )
    return f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Default Extension="jpeg" ContentType="image/jpeg"/>
  <Default Extension="jpg" ContentType="image/jpeg"/>
  <Default Extension="png" ContentType="image/png"/>
  <Default Extension="fntdata" ContentType="application/x-fontdata"/>
  <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"/>
  <Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>
  {layout_overrides}
  <Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>
  <Override PartName="/ppt/presProps.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presProps+xml"/>
  <Override PartName="/ppt/viewProps.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.viewProps+xml"/>
  <Override PartName="/ppt/tableStyles.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.tableStyles+xml"/>
  <Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
  <Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>
</Types>'''


if __name__ == "__main__":
    build()

#!/usr/bin/env python3
"""
Build OCM-Sovereign-Delivery-Exec.pptx by instantiating slides from the
OCM-Master-Template.potx layouts and filling placeholders.

Why this is structured this way: the .potx template owns visual styling
(palette, type scale, eyebrow + title + footer chrome, column rules, tile
backgrounds). This script only supplies content. To restyle the whole deck,
edit the layouts in build_potx.py and rebuild the .potx — every deck that
uses it picks up the change automatically.

Layout-specific extras (banner image on slide 1, brand row on slide 1 + slide
10, diagram images on content slides) are added inline because they're
deck-specific, not template-wide.

Usage:
    .venv/bin/python build_pptx.py
"""
from __future__ import annotations

import shutil
import subprocess
import sys
import zipfile
from pathlib import Path
from tempfile import TemporaryDirectory

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

from icon_strokes import STROKE_THIN, STROKE_REGULAR, STROKE_BOLD

# -----------------------------------------------------------------------------
# Paths
# -----------------------------------------------------------------------------

SCRIPT_DIR = Path(__file__).resolve().parent
DECK_DIR = SCRIPT_DIR.parent
MARKETING_DIR = DECK_DIR.parent.parent
ASSETS_DIR = MARKETING_DIR / "assets"
DIAGRAMS_DIR = DECK_DIR / "diagrams"
ICONS_DIR = DIAGRAMS_DIR / "icons"
THEME_DIR = DECK_DIR / "theme"
RASTER_DIR = SCRIPT_DIR / "_raster"

POTX_PATH = DECK_DIR / "OCM-Master.potx"
OUTPUT_PPTX = DECK_DIR / "OCM-Sovereign-Delivery-Internal-Sponsor.pptx"

RASTER_DIR.mkdir(exist_ok=True)


# -----------------------------------------------------------------------------
# Slide geometry — 16:9 @ 1920×1080
# -----------------------------------------------------------------------------

SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
PX = 9525  # 1 px in EMU at 96 dpi

def px(n: float) -> Emu:
    return Emu(int(n * PX))


# -----------------------------------------------------------------------------
# OCM brand palette (canonical, mirror of build_potx.py PALETTE)
# -----------------------------------------------------------------------------

class C:
    BLUE       = RGBColor(0x0F, 0x6B, 0xFF)   # accent1 / brand-blue-dark
    BLUE_MID   = RGBColor(0x0A, 0x3A, 0x99)   # accent2 / brand-blue-mid
    CYAN       = RGBColor(0x5C, 0xD6, 0xFF)   # accent3 / brand-cyan
    GREY_MID   = RGBColor(0x6B, 0x72, 0x80)   # accent4
    BLUE_NIGHT = RGBColor(0x0A, 0x15, 0x30)   # accent5
    GREY_SOFT  = RGBColor(0xF3, 0xF4, 0xF6)   # accent6
    BLACK      = RGBColor(0x00, 0x00, 0x00)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)


# -----------------------------------------------------------------------------
# OOXML namespaces (used for the gradient title surgery)
# -----------------------------------------------------------------------------

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# -----------------------------------------------------------------------------
# SVG → PNG rasterization (for diagrams + icons)
# -----------------------------------------------------------------------------

def rasterize_svg(svg_path: Path, target_w_px: int) -> Path:
    if not svg_path.exists():
        raise FileNotFoundError(svg_path)
    out = RASTER_DIR / (svg_path.stem + f"_{target_w_px}.png")
    if out.exists() and out.stat().st_mtime >= svg_path.stat().st_mtime:
        return out
    subprocess.run(
        ["rsvg-convert", "--width", str(target_w_px), "--keep-aspect-ratio",
         str(svg_path), "-o", str(out)],
        check=True, capture_output=True,
    )
    return out


def rasterize_svg_recolored(svg_path: Path, target_w_px: int,
                            color_hex: str,
                            stroke_width: float | None = None) -> Path:
    """Rasterize a `currentColor`-based SVG (Tabler icon family) with an
    explicit stroke/fill colour. See build_pptx.py for the rationale.

    Resolution order: prebuilt asset under diagrams/icons/prebuilt/
    first; on-the-fly _raster/ cache as fallback. Same scheme as the
    Exec-deck builder."""
    if not svg_path.exists():
        raise FileNotFoundError(svg_path)
    colour = color_hex.lstrip("#").upper()

    # 1. Prebuilt asset shortcut --------------------------------------
    colour_name = {"0F6BFF": "brand-blue", "FFFFFF": "white"}.get(colour)
    if colour_name is not None and stroke_width is not None:
        prebuilt_dir = svg_path.parent / "prebuilt"
        if float(stroke_width).is_integer():
            stroke_tag = f"{int(stroke_width)}.0"
        else:
            stroke_tag = f"{stroke_width:g}"
        prebuilt_png = prebuilt_dir / f"{svg_path.stem}-stroke-{stroke_tag}-{colour_name}.png"
        if prebuilt_png.exists():
            return prebuilt_png

    # 2. Fallback ------------------------------------------------------
    sw_tag = f"_sw{stroke_width:g}".replace(".", "p") if stroke_width is not None else ""
    out = RASTER_DIR / f"{svg_path.stem}_{target_w_px}_{colour}{sw_tag}.png"
    if out.exists() and out.stat().st_mtime >= svg_path.stat().st_mtime:
        return out
    src = svg_path.read_text(encoding="utf-8")
    patched = src.replace("currentColor", f"#{colour}")
    if stroke_width is not None:
        import re as _re
        patched = _re.sub(
            r'stroke-width="[^"]*"',
            f'stroke-width="{stroke_width:g}"',
            patched,
        )
    tmp = RASTER_DIR / f"{svg_path.stem}_{colour}{sw_tag}.svg"
    tmp.write_text(patched, encoding="utf-8")
    subprocess.run(
        ["rsvg-convert", "--width", str(target_w_px), "--keep-aspect-ratio",
         str(tmp), "-o", str(out)],
        check=True, capture_output=True,
    )
    return out


def first_existing(*candidates: Path) -> Path | None:
    for c in candidates:
        if c.exists():
            return c
    return None


# -----------------------------------------------------------------------------
# Open .potx as .pptx
# -----------------------------------------------------------------------------

def open_template_as_pptx() -> Presentation:
    """python-pptx refuses .potx files, so transparently swap the
    presentation content type to .pptx-flavored when loading. The output we
    save will still be a valid .pptx (we never write the template type back)."""
    if not POTX_PATH.exists():
        raise FileNotFoundError(
            f"{POTX_PATH} not found — run `python build_potx.py` first."
        )
    tmp_pptx = RASTER_DIR / "_potx_loaded.pptx"
    with zipfile.ZipFile(POTX_PATH, "r") as src:
        data = {n: src.read(n) for n in src.namelist()}
    ct = data["[Content_Types].xml"].decode("utf-8")
    ct = ct.replace(
        "presentationml.template.main+xml",
        "presentationml.presentation.main+xml",
    )
    data["[Content_Types].xml"] = ct.encode("utf-8")
    with zipfile.ZipFile(tmp_pptx, "w", zipfile.ZIP_DEFLATED) as out:
        for name, blob in data.items():
            out.writestr(name, blob)
    return Presentation(str(tmp_pptx))


# -----------------------------------------------------------------------------
# Placeholder helpers
# -----------------------------------------------------------------------------

def find_placeholder(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(f"placeholder idx={idx} not found on slide using "
                   f"layout {slide.slide_layout.name!r}")


def delete_placeholder(slide, idx: int):
    """Remove a placeholder shape entirely from the slide. Use when a layout
    supplies a placeholder you don't want on this particular slide (e.g. the
    body box on the Plain layout for slides that draw their own content)."""
    ph = find_placeholder(slide, idx)
    sp = ph._element
    sp.getparent().remove(sp)


def set_text(slide, idx: int, text: str, *, color: RGBColor | None = None,
             align_left: bool = False):
    """Set a placeholder's text. Layout's lstStyle supplies size/font/case;
    we only override color when needed (e.g. hero title white).

    align_left=True forces left alignment via <a:pPr algn="l"/> — useful on
    placeholders that PowerPoint might otherwise center (Hero title, etc.).
    """
    from pptx.enum.text import PP_ALIGN
    ph = find_placeholder(slide, idx)
    tf = ph.text_frame
    tf.clear()
    paragraphs = text.split("\n")
    for i, line in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if align_left:
            p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        if color is not None:
            run.font.color.rgb = color


def set_split_gradient_title(slide, idx: int, prefix: str, noun: str,
                              align_left: bool = True):
    """Hero title second line: prefix in white, noun with the OCM gradient."""
    from pptx.enum.text import PP_ALIGN
    ph = find_placeholder(slide, idx)
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align_left:
        p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = prefix
    r1.font.color.rgb = C.WHITE
    r2 = p.add_run()
    r2.text = noun
    # Replace solid fill with the OCM gradient on this run.
    rPr = r2._r.get_or_add_rPr()
    for tag in ("solidFill", "gradFill", "noFill"):
        existing = rPr.find(f"{{{A_NS}}}{tag}")
        if existing is not None:
            rPr.remove(existing)
    grad_xml = (
        f'<a:gradFill xmlns:a="{A_NS}" flip="none" rotWithShape="1">'
        '<a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="FFFFFF"/></a:gs>'
        '<a:gs pos="35000"><a:srgbClr val="5CD6FF"/></a:gs>'
        '<a:gs pos="75000"><a:srgbClr val="0F6BFF"/></a:gs>'
        '</a:gsLst>'
        '<a:lin ang="0" scaled="1"/>'
        '</a:gradFill>'
    )
    rPr.insert(0, etree.fromstring(grad_xml))


def set_gradient_title(slide, idx: int, text: str, *, align_left: bool = False):
    """Set a title with the OCM gradient applied to the entire run."""
    from pptx.enum.text import PP_ALIGN
    ph = find_placeholder(slide, idx)
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if align_left:
        p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    rPr = r._r.get_or_add_rPr()
    for tag in ("solidFill", "gradFill", "noFill"):
        existing = rPr.find(f"{{{A_NS}}}{tag}")
        if existing is not None:
            rPr.remove(existing)
    grad_xml = (
        f'<a:gradFill xmlns:a="{A_NS}" flip="none" rotWithShape="1">'
        '<a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="FFFFFF"/></a:gs>'
        '<a:gs pos="35000"><a:srgbClr val="5CD6FF"/></a:gs>'
        '<a:gs pos="75000"><a:srgbClr val="0F6BFF"/></a:gs>'
        '</a:gsLst>'
        '<a:lin ang="0" scaled="1"/>'
        '</a:gradFill>'
    )
    rPr.insert(0, etree.fromstring(grad_xml))


def set_action_path_lines(slide, idx: int, lines: list[tuple[str, str]],
                            sep: str = " — "):
    """Render a body placeholder where each line is split into ACTION (blue)
    and PATH (white) by `sep`. Used for the CTA slide."""
    ph = find_placeholder(slide, idx)
    tf = ph.text_frame
    tf.clear()
    for i, (action, path) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run()
        r1.text = action
        r1.font.color.rgb = C.CYAN
        r1.font.bold = True
        r2 = p.add_run()
        r2.text = sep + path
        r2.font.color.rgb = C.WHITE


def set_blue_box_bullets(slide, idx: int, items: list[str]):
    """Render a body placeholder as a list with blue square bullets.

    Each item gets a small filled square (▪) prefix in OCM brand blue, then
    the body text in black. We don't use PowerPoint's <a:buChar> machinery —
    Aptos doesn't render the small filled square consistently — so we put a
    blue-coloured glyph as the first run of each paragraph.
    """
    ph = find_placeholder(slide, idx)
    tf = ph.text_frame
    tf.clear()
    for i, body in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        bullet = p.add_run()
        bullet.text = "▪  "
        bullet.font.color.rgb = C.BLUE
        bullet.font.bold = True
        text_run = p.add_run()
        text_run.text = body
        text_run.font.color.rgb = C.BLACK


# -----------------------------------------------------------------------------
# Static decoration helpers (banner, brand row, diagram embeds)
# -----------------------------------------------------------------------------

def add_banner_full_bleed(slide, image_path: Path):
    """Place an image as full-slide background. Sent to back."""
    if not image_path.exists():
        return
    pic = slide.shapes.add_picture(str(image_path), 0, 0,
                                    width=px(SLIDE_W_PX),
                                    height=px(SLIDE_H_PX))
    # Move to bottom of z-order so placeholders sit on top.
    spTree = pic._element.getparent()
    spTree.remove(pic._element)
    # Insert just after grpSpPr (the first non-element header).
    insert_at = 0
    for i, el in enumerate(spTree):
        if el.tag.endswith("}grpSpPr"):
            insert_at = i + 1
            break
    spTree.insert(insert_at, pic._element)


def add_brand_row(slide):
    """Bottom-of-slide brand row: OCM logo left, NeoNephos logo right (white
    variants for use over dark backgrounds)."""
    ocm_svg = ASSETS_DIR / "ocm" / "ocm-horizontal-white.svg"
    nn_svg = ASSETS_DIR / "neonephos" / "neonephos-foundation-horizontal-white.svg"
    ocm_png = rasterize_svg(ocm_svg, target_w_px=400)
    nn_png = rasterize_svg(nn_svg, target_w_px=400)
    ocm_pic = slide.shapes.add_picture(
        str(ocm_png), px(96), px(SLIDE_H_PX - 56 - 76), height=px(76)
    )
    nn_pic = slide.shapes.add_picture(
        str(nn_png), px(96), px(SLIDE_H_PX - 56 - 52), height=px(52)
    )
    nn_pic.left = px(SLIDE_W_PX - 96) - nn_pic.width


def add_diagram(slide, svg_path: Path | None,
                 x_px: int, y_px: int,
                 max_w_px: int, max_h_px: int):
    """Rasterize a diagram SVG and centre it inside the slot. See
    build_pptx.py for full rationale."""
    if svg_path is None or not svg_path.exists():
        return
    try:
        delete_placeholder(slide, 10)
    except KeyError:
        pass
    png = rasterize_svg(svg_path, target_w_px=max_w_px)
    pic = slide.shapes.add_picture(str(png), px(x_px), px(y_px),
                                    width=px(max_w_px))
    if pic.height > px(max_h_px):
        ratio = px(max_h_px) / pic.height
        pic.height = px(max_h_px)
        pic.width = int(pic.width * ratio)
    pic.left = px(x_px) + (px(max_w_px) - pic.width) // 2
    pic.top  = px(y_px) + (px(max_h_px) - pic.height) // 2


def add_tile_icon(slide, tile_x_px: int, tile_y_px: int, icon_name: str):
    """Place a brand-blue icon at the top-left of a tile (label sits to the
    right in the same header row). See build_pptx.py for the rationale."""
    icon_path = ICONS_DIR / icon_name
    if not icon_path.exists():
        return
    png = rasterize_svg_recolored(icon_path, target_w_px=96, color_hex="0F6BFF")
    slide.shapes.add_picture(
        str(png),
        px(tile_x_px + 24), px(tile_y_px + 24),
        width=px(48), height=px(48),
    )


def add_logo_row(slide, logos: list, y_px: int,
                  row_h_px: int = 120,
                  max_logo_w_px: int = 320, max_logo_h_px: int = 80,
                  caption_pt: int = 14):
    """Three (or more) logos centred in a row, sized to a uniform height
    (max_logo_h_px) so the row reads as visually consistent.

    `logos` accepts entries as:
      - Path                — no link, no caption
      - (Path, url)         — clickable, no caption
      - (Path, url, caption) — clickable + caption text rendered below

    Captions are useful for icon-only marks (Kyma, OpenControlPlane) where
    the logo alone doesn't carry the project name.
    """
    from pptx.enum.text import PP_ALIGN
    margin_x = 160
    inner_w = SLIDE_W_PX - 2 * margin_x
    n = len(logos)
    slot_w = inner_w // n
    for i, entry in enumerate(logos):
        caption = None
        if isinstance(entry, tuple):
            if len(entry) == 3:
                path, url, caption = entry
            else:
                path, url = entry
        else:
            path, url = entry, None
        if path is None or not path.exists():
            continue
        if path.suffix.lower() == ".svg":
            img = rasterize_svg(path, target_w_px=max_logo_w_px * 2)
        else:
            img = path
        slot_x = margin_x + i * slot_w
        # Add at native size, then scale to fit the height constraint first.
        pic = slide.shapes.add_picture(str(img), px(slot_x), px(y_px))
        if pic.height != px(max_logo_h_px):
            ratio = px(max_logo_h_px) / pic.height
            pic.height = px(max_logo_h_px)
            pic.width = int(pic.width * ratio)
        # Then enforce width cap (rare — only triggers for very wide logos).
        if pic.width > px(max_logo_w_px):
            ratio = px(max_logo_w_px) / pic.width
            pic.width = px(max_logo_w_px)
            pic.height = int(pic.height * ratio)
        pic.left = px(slot_x) + (px(slot_w) - pic.width) // 2
        pic.top = px(y_px) + (px(row_h_px) - pic.height) // 2
        if url:
            pic.click_action.hyperlink.address = url
        if caption:
            cap_y = y_px + row_h_px + 6
            tb = slide.shapes.add_textbox(px(slot_x), px(cap_y),
                                           px(slot_w), px(28))
            tf = tb.text_frame
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = caption
            r.font.name = "Aptos"
            r.font.size = Pt(caption_pt)
            r.font.color.rgb = C.GREY_MID
            # Caption is plain descriptive text (e.g. "Managed Kubernetes"),
            # not a link — the logo above carries the URL via pic.click_action.
            # Earlier iterations had captions duplicate the logo name and act
            # as hyperlinks; both were dropped (see slide 10a comment).


# -----------------------------------------------------------------------------
# Hero / Tile geometry helpers — must mirror what the .potx layouts use
# -----------------------------------------------------------------------------

# These constants drive where the on-slide decoration (banner, brand row,
# tile icons) is placed. They MUST match the corresponding layout positions
# in build_potx.py — if you change one place, change the other.

TILE_X0_PX = 120
TILE_Y0_PX = 520
TILE_W_PX = 544
TILE_H_PX = 230
TILE_GUTTER_PX = 24


def tile_origin(index: int) -> tuple[int, int]:
    col = index % 3
    row = index // 3
    return (TILE_X0_PX + col * (TILE_W_PX + TILE_GUTTER_PX),
            TILE_Y0_PX + row * (TILE_H_PX + TILE_GUTTER_PX))


# -----------------------------------------------------------------------------
# Layout lookup
# -----------------------------------------------------------------------------

def layouts_by_name(prs: Presentation) -> dict[str, object]:
    return {l.name: l for l in prs.slide_masters[0].slide_layouts}


# =============================================================================
# Build the deck
# =============================================================================

def build():
    prs = open_template_as_pptx()
    layouts = layouts_by_name(prs)

    expected = {"Hero", "CTA", "Content / 3-Column",
                "Content / Diagram", "Content / Tiles", "Content / 2-Column",
                "Section Divider", "Plain", "Plain / Compact"}
    missing = expected - set(layouts)
    if missing:
        sys.exit(f"template missing expected layouts: {missing}")

    # ---- SLIDE 1 — HERO (internal-sponsor, observation-frame) --------------
    # Hero opens with two parallel stop-sentences: a concrete observation about
    # the SAP-internal status quo (every LoB delivers independently) and the
    # punchline (and does so every release cycle). Subtitle resolves the
    # observation: each LoB still ships its own artifacts — what changes is
    # that they ship "on the same model" (shared concept, shared vocabulary,
    # shared signing/transport/compliance mechanics), not on parallel
    # bespoke pipelines.
    s = prs.slides.add_slide(layouts["Hero"])
    add_banner_full_bleed(s, THEME_DIR / "OCM-Banner.png")
    set_text(s, 1, "Every LoB ships.", color=C.WHITE)
    set_gradient_title(s, 2, "Separately, every time.", align_left=True)
    set_text(s, 3,
             "OCM is the shared standard. Each LoB still ships — "
             "on the same model.",
             color=C.CYAN)
    set_text(s, 4,
             "Open Component Model — open source, NeoNephos Foundation. "
             "Stewarded by SAP.",
             color=C.WHITE)
    add_brand_row(s)

    # ---- SLIDE 2 — WHY NOW (internal lens) ---------------------------------
    # Three columns, each parallel [observation. consequence-if-no-action.] form.
    # Columns 2 and 3 imply action by stating the cost of inaction; column 1
    # was reworded to follow the same pattern — "biggest contributor shapes
    # the standard" makes the implicit point that staying biggest is an active
    # choice, not a given.
    s = prs.slides.add_slide(layouts["Content / 3-Column"])
    set_text(s, 1, "WHY NOW")
    set_text(s, 2, "Compliance and sovereignty are given.\nOur strategic position is a choice.")
    set_text(s, 10, "ECOSYSTEM VELOCITY")
    set_text(s, 11, "The peer ecosystem is converging.\n"
                     "The biggest contributor shapes the standard.")
    set_text(s, 12, "THE WINDOW")
    set_text(s, 13, "The rails are being laid now.\n"
                     "Late entrants pay migration cost.")
    set_text(s, 14, "DISINVESTMENT COST")
    set_text(s, 15, "Walking away costs more than staying.\n"
                     "The standard gets shaped without us.")

    # ---- SLIDE 3 — MEET OCM (hub-and-spoke diagram, Option 3 reframe) -------
    # Diagram positioned per user spec 2026-06-17: 50.02 × 15.93 cm,
    # x=-2.4cm (slight bleed left), y=11.65cm.
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "THE ANSWER")
    set_text(s, 2, "Meet OCM. One identity, every boundary.")
    add_diagram(s, DIAGRAMS_DIR / "03-meet-ocm-hub-and-spoke.svg",
                 x_px=60, y_px=240, max_w_px=1800, max_h_px=780)

    # ---- SLIDE 4a — THE SHIFT (diagram, Option A) --------------------------
    # Hand-editing in PowerPoint settled on this ordering: SBOD diagram first
    # (shows the picture), bullets after (explains it). The two SBOD diagram
    # variants from earlier iterations (5-tile grid and 3-property layout)
    # were dropped — only the artifact-list + signature-bracket layout
    # remains, because it shows SBOM as one line in a list.
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "THE SHIFT — SBOM INSIDE SBOD")
    set_text(s, 2, "SBOM lists. SBOD delivers.")
    diagram = first_existing(
        DIAGRAMS_DIR / "04-sbom-inside-sbod.svg",
        DIAGRAMS_DIR / "04-sbom-vs-sbod.svg",
    )
    if diagram:
        add_diagram(s, diagram, x_px=60, y_px=240, max_w_px=1800, max_h_px=780)

    # ---- SLIDE 4b — THE SHIFT, SBOD (text-only, internal-sponsor) ----------
    s = prs.slides.add_slide(layouts["Plain / Compact"])
    set_text(s, 1, "THE SHIFT")
    set_text(s, 2, "SBOM lists. SBOD delivers.")
    set_blue_box_bullets(s, 10, [
        "SBOM — what's inside your software. Built for inventory.",
        "A Software Bill of Delivery (SBOD) — what you delivered, "
        "how to verify, transport, operate. Built for delivery.",
        "SBOD contains SBOM. OCM doesn't replace your SBOM tooling — "
        "OCM gives the SBOM an envelope.",
        "SBOD is the category SAP defined. "
        "Now standardised through NeoNephos.",
    ])

    # ---- SLIDE 5 — HOW OCM COMPOSES (NEW, comparator slide) ----------------
    # Two-line columns: [status quo today] / [OCM contribution]. Each column
    # disarms a "we already have this" objection: signing, transport,
    # compliance. The earlier multi-clause version had its punchlines
    # consumed by the slide; the speaker now delivers them verbally.
    s = prs.slides.add_slide(layouts["Content / 3-Column"])
    set_text(s, 1, "HOW OCM COMPOSES")
    set_text(s, 2, "Composes around your existing stack.")
    set_text(s, 10, "SIGNING")
    set_text(s, 11, "You sign artifacts.\nOCM signs the release.")
    set_text(s, 12, "TRANSPORT")
    set_text(s, 13, "Your registries differ.\nOCM moves the release across them.")
    set_text(s, 14, "COMPLIANCE")
    set_text(s, 15, "Your scanners see one artifact at a time.\n"
                     "OCM correlates findings to the release.")

    # ---- SLIDE 6 — OCM IN ONE PICTURE --------------------------------------
    # Pack · Sign · Transport · Deploy → Sovereign Cloud. The DEPLOY tile's
    # second line was rewritten "OCM K8s Controllers." (with period, no
    # imperative) so all four tiles read with the same grammatical rhythm.
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "OCM IN ONE PICTURE")
    set_text(s, 2, "Pack · Sign · Transport · Deploy")
    diagram6 = first_existing(
        DIAGRAMS_DIR / "05-pack-sign-transport-deploy-v2.svg",
        DIAGRAMS_DIR / "05-pack-sign-transport-deploy.svg",
    )
    add_diagram(s, diagram6, x_px=60, y_px=240, max_w_px=1800, max_h_px=780)

    # ---- SLIDE 7a — SOVEREIGN-READY (text-only, anchor + halfsentence) -----
    # Internal-sponsor format: anchor word + characterisation half-sentence +
    # consequence half-sentence. Shorter than exec, longer than a one-liner,
    # carries enough for hand-out reading while still leaving punchlines for
    # the speaker.
    s = prs.slides.add_slide(layouts["Plain / Compact"])
    set_text(s, 1, "SOVEREIGN-READY")
    set_text(s, 2, "Trust, but verify.")
    set_blue_box_bullets(s, 10, [
        "Identity — location-independent. "
        "The component carries its name regardless of registry.",
        "Signatures — location-independent. "
        "Sign once at source, verify anywhere downstream. No callback upstream.",
        "Day-2 ops — happen inside the boundary. "
        "Subscribe, pull upgrades, scale across regions. Still no callback.",
        "Transfer — self-contained. "
        "Every artifact travels with the component.",
    ])

    # ---- SLIDE 7b — SOVEREIGN-READY (diagram only) -------------------------
    s = prs.slides.add_slide(layouts["Content / Diagram"])
    set_text(s, 1, "SOVEREIGN-READY — AIR-GAP")
    set_text(s, 2, "Trust travels with the component.")
    add_diagram(s, DIAGRAMS_DIR / "06-sovereign-airgap.svg",
                 x_px=60, y_px=240, max_w_px=1800, max_h_px=780)

    # ---- SLIDE 8 — SCAN (shortened eyebrow, 4 bullets) ---------------------
    # Eyebrow was "SCAN — COMPLIANCE-NATIVE WITH OPEN DELIVERY GEAR" — too
    # long, doubled with the subtitle. Trimmed to just "SCAN" (parallel to the
    # PACK/SIGN/TRANSPORT/DEPLOY verb chain on slide 6, positioning Scan as
    # the 5th step). Internal audience knows ODG, so the definition bullet is
    # dropped — the four substantive bullets remain.
    s = prs.slides.add_slide(layouts["Plain"])
    set_text(s, 1, "SCAN")
    set_text(s, 2, "Compliance as a system property —\nnot a quarterly retrofit.")
    set_blue_box_bullets(s, 10, [
        "The Compliance Dashboard — every component, every finding, one view.",
        "Continuous scans — asynchronous, even post-release.",
        "Contextual rescoring — patch what matters, not the noise.",
        "Identity-correlated evidence — auditors get answers, not spreadsheets.",
    ])

    # ---- SLIDE 9 — WHAT OCM UNLOCKS FOR SAP (tiles, internal outcomes) -----
    s = prs.slides.add_slide(layouts["Content / Tiles"])
    set_text(s, 1, "WHAT OCM UNLOCKS FOR SAP")
    set_text(s, 2, "Six outcomes from one shared primitive.")
    tiles = [
        ("package-export.svg", "Faster sovereign delivery",
         "Pack once, ship everywhere.\n"
         "Sovereign Cloud for all products."),
        ("report-analytics.svg", "Compliance leverage across LoBs",
         "Report from one shared primitive —\n"
         "ODG correlates all findings."),
        ("git-merge.svg", "Integration after acquisition",
         "Acquired companies\n"
         "converge onto one model."),
        ("radar.svg", "Cross-LoB security correlation",
         "Blast radius is one query —\n"
         "answered via the OCM coordinate system."),
        ("source-of-truth.svg", "One source of truth",
         "One signed descriptor per delivery.\n"
         "Rebuild any landscape."),
        ("heart-handshake.svg", "Ecosystem stewardship",
         "SAP investment compounds with\n"
         "the open-peer ecosystem."),
    ]
    for i, (icon, label, body) in enumerate(tiles):
        # Tile placeholders alternate label/body: idx 20+2i = label, 21+2i = body
        set_text(s, 20 + i * 2, label)
        set_text(s, 21 + i * 2, body)
        x, y = tile_origin(i)
        add_tile_icon(s, x, y, icon)

    # ---- SLIDE 10a — WHERE OCM IS SHIPPING — OPEN ECOSYSTEM ---------------
    # Internal-sponsor: open-peer wall, logo-based. Captions changed from
    # logo-name duplication to substantive characterisations (drawn from
    # each project's own homepage tagline) so the slide teaches what each
    # project does, not just names them. Logos are clickable; the captions
    # are now plain grey labels, not hyperlinks.
    s = prs.slides.add_slide(layouts["Plain / Compact"])
    set_text(s, 1, "WHERE OCM IS SHIPPING — OPEN ECOSYSTEM")
    set_text(s, 2, "Peer in the open ecosystem.")
    delete_placeholder(s, 10)
    add_logo_row(s, [
        (ASSETS_DIR / "adopters" / "gardener" / "gardener-horizontal-color.svg",
         "https://gardener.cloud", "Managed Kubernetes"),
        (ASSETS_DIR / "adopters" / "kyma" / "kyma-icon-color.svg",
         "https://kyma-project.io", "Cloud-native runtime"),
        (ASSETS_DIR / "adopters" / "open-control-plane" / "opencontrolplane-icon-color.svg",
         "https://open-control-plane.io", "Control-plane framework"),
        (ASSETS_DIR / "adopters" / "konfidence" / "konfidence-horizontal-light.svg",
         "https://konfidence.cloud", "Reproducible delivery"),
    ], y_px=560, max_logo_w_px=360, max_logo_h_px=120, caption_pt=20)
    # "Aligned with [NeoNephos logo]" — single, larger, clickable NeoNephos
    # logo carries the alignment claim.
    add_centred_proof_with_logo(
        s, 820,
        "Aligned with ",
        ASSETS_DIR / "adopters" / "neonephos" / "neonephos-foundation-horizontal-color.svg",
        "",
        logo_url="https://neonephos.org",
        logo_caption="NeoNephos")

    # ---- SLIDE 10b — WHERE OCM IS SHIPPING — SAP -------------------------
    # Internal-only delivery infrastructure converging on OCM. Bullets shortened
    # to parallel "Name — short characterisation" form so the slide reads as
    # proof ("running on five platforms already") rather than a teaching map.
    s = prs.slides.add_slide(layouts["Plain / Compact"])
    set_text(s, 1, "WHERE OCM IS SHIPPING — SAP")
    set_text(s, 2, "Backbone of internal SAP delivery.")
    set_blue_box_bullets(s, 10, [
        "Hyperspace — internal Dev Portal & product delivery.",
        "Release-Based Shipment Channel (RBSC) — customer shipment channel.",
        "Common Service Infrastructure (CSI) — shared internal services platform.",
        "Steampunk — ABAP Development PaaS.",
        "Greenhouse — Cloud ops platform.",
    ])

    # ---- SLIDE 11 — CTA (sponsor / scale / standardize) -------------------
    s = prs.slides.add_slide(layouts["CTA"])
    set_text(s, 1, "Sponsor. Scale. Standardize.", color=C.WHITE)
    set_action_path_lines(s, 2, [
        ("Sponsor",     "Allocate engineering capacity to OCM stewardship in your LoB."),
        ("Scale",       "Pack one regulated component as an OCM component this quarter."),
        ("Standardize", "Bring your LoB into the OCM steering conversation — SAP Slack #sap-tech-ocm."),
    ])
    add_brand_row(s)

    # ---- HIDDEN — Trademark & licensing notice -----------------------------
    add_appendix_glossary_slide(prs, layouts)
    add_hidden_trademark_slide(prs, layouts)

    prs.save(str(OUTPUT_PPTX))
    print(f"Wrote {OUTPUT_PPTX}")


# -----------------------------------------------------------------------------
# Inline-decoration helpers used by slide 9
# -----------------------------------------------------------------------------

def add_label_at(slide, y_px: int, text: str):
    """Brand-blue ALL-CAPS section label — slide 9 logo wall headers."""
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(px(120), px(y_px),
                                   px(SLIDE_W_PX - 240), px(36))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = "Aptos"
    f.size = Pt(18)
    f.bold = True
    f.color.rgb = C.BLUE
    rPr = r._r.get_or_add_rPr()
    rPr.set("cap", "all")
    rPr.set("spc", "110")


def add_centred_proof(slide, y_px: int, text: str):
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(px(120), px(y_px),
                                   px(SLIDE_W_PX - 240), px(80))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.3
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(16)
    r.font.italic = True
    r.font.color.rgb = C.BLUE_MID


def add_centred_proof_with_logo(slide, y_px: int, text_before: str,
                                 logo_path: Path, text_after: str,
                                 logo_h_px: int = 40,
                                 logo_url: str | None = None,
                                 logo_caption: str | None = None,
                                 caption_pt: int = 20):
    """Centred proof line with an inline image substituted for one word.

    Renders text_before + <logo> + text_after as three siblings on the same
    baseline, all horizontally centred as a group. Approximates text widths
    from character counts (python-pptx has no font-metrics API); good enough
    for a single static line. Style mirrors `add_centred_proof` minus the
    italic — "Aligned with [LOGO]" reads as a regular caption next to a
    branded logo, not as a quoted aside.

    If `logo_url` is provided the logo picture becomes clickable. If
    `logo_caption` is provided, a small caption (e.g. "neonephos.org") is
    rendered under the logo, matching the caption_pt size used by
    add_logo_row so the visual hierarchy stays consistent.
    """
    from pptx.enum.text import PP_ALIGN
    # Approx width in px for 20pt regular Aptos: ~11.0 px/char average. Was
    # 10.5 when the line was italic (italic glyphs run a touch narrower);
    # bumped slightly so the composite stays visually centred at regular weight.
    char_w_px = 11.0
    # Horizontal padding between text segments and the inline logo, so the
    # logo doesn't sit flush against the italic text baseline edge.
    GAP_PX = 16
    tb_before_w_px = int(len(text_before) * char_w_px) + 8
    tb_after_w_px = int(len(text_after) * char_w_px) + 8

    # Rasterize SVG so we can read its natural size, then scale to logo_h_px.
    if logo_path.suffix.lower() == ".svg":
        img = rasterize_svg(logo_path, target_w_px=600)
    else:
        img = logo_path
    # Add picture off-slide first to read intrinsic size, then move/resize.
    pic = slide.shapes.add_picture(str(img), px(0), px(0))
    ratio = px(logo_h_px) / pic.height
    pic.height = px(logo_h_px)
    pic.width = int(pic.width * ratio)
    logo_w_px = pic.width / PX

    total_w_px = tb_before_w_px + GAP_PX + logo_w_px + GAP_PX + tb_after_w_px
    start_x_px = (SLIDE_W_PX - total_w_px) / 2

    # Vertical: text-frame top at y_px; logo nudged to sit on the same
    # visual baseline as 20pt italic text (~24px line height, cap-height ~22).
    logo_y_px = y_px + 2
    text_y_px = y_px

    # ---- text_before (right-aligned in its slot) ---------------------------
    tb1 = slide.shapes.add_textbox(px(start_x_px), px(text_y_px),
                                    px(tb_before_w_px), px(50))
    tf1 = tb1.text_frame
    tf1.margin_left = tf1.margin_right = 0
    tf1.margin_top = tf1.margin_bottom = 0
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.RIGHT
    r1 = p1.add_run()
    r1.text = text_before
    r1.font.name = "Aptos"
    r1.font.size = Pt(20)
    r1.font.italic = False
    r1.font.color.rgb = C.BLUE_MID

    # ---- inline logo -------------------------------------------------------
    pic.left = px(start_x_px + tb_before_w_px + GAP_PX)
    pic.top = px(logo_y_px)
    if logo_url:
        pic.click_action.hyperlink.address = logo_url

    # ---- caption under the whole composite (text + logo) -----------------
    # Centred under the entire line so it sits under the visual mid-point of
    # "Aligned with [LOGO]", not just under the logo (which would read as
    # off-centre). Caption styling matches add_logo_row's caption convention.
    if logo_caption:
        cap_y = logo_y_px + logo_h_px + 6
        cap_w = max(int(total_w_px), 200)
        cap_x = start_x_px + (total_w_px - cap_w) / 2
        cb = slide.shapes.add_textbox(px(cap_x), px(cap_y), px(cap_w), px(28))
        cf = cb.text_frame
        cf.margin_left = cf.margin_right = 0
        cf.margin_top = cf.margin_bottom = 0
        cp = cf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = logo_caption
        cr.font.name = "Aptos"
        cr.font.size = Pt(caption_pt)
        cr.font.color.rgb = C.GREY_MID
        if logo_url:
            cr.hyperlink.address = logo_url

    # ---- text_after (left-aligned in its slot) -----------------------------
    if text_after:
        tb2 = slide.shapes.add_textbox(
            px(start_x_px + tb_before_w_px + GAP_PX + logo_w_px + GAP_PX),
            px(text_y_px), px(tb_after_w_px), px(50))
        tf2 = tb2.text_frame
        tf2.margin_left = tf2.margin_right = 0
        tf2.margin_top = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = text_after
        r2.font.name = "Aptos"
        r2.font.size = Pt(20)
        r2.font.italic = False
        r2.font.color.rgb = C.BLUE_MID


def add_source_line(slide, y_px: int, text: str):
    """Small grey source/reference line — for slide footers that name the
    canonical URL of a project (NeoNephos, ODG, OpenControlPlane, etc.)."""
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(px(120), px(y_px),
                                   px(SLIDE_W_PX - 240), px(40))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(16)
    r.font.color.rgb = C.GREY_MID


# Glossary on the appendix slide. Term → expansion. Two-column layout, term
# in brand-blue bold, expansion in black. Mirrors the external deck list so
# the abbreviations have one definition surface across both decks. Sorted
# alphabetically (case-insensitive) for scannable reference.
GLOSSARY_ENTRIES: list[tuple[str, str]] = [
    ("BSI C5",    "Bundesamt für Sicherheit in der Informationstechnik — Cloud Computing Compliance Criteria Catalogue."),
    ("BTP",       "SAP Business Technology Platform."),
    ("CRA",       "Cyber Resilience Act — EU regulation on cybersecurity for products with digital elements."),
    ("DORA",      "Digital Operational Resilience Act — EU regulation for ICT risk in financial services."),
    ("FedRAMP",   "Federal Risk and Authorization Management Program — US standardised cloud security assessment."),
    ("FISMA",     "Federal Information Security Modernization Act — US federal information security mandate."),
    ("Grype",     "Open-source vulnerability scanner for container images and filesystems (Anchore)."),
    ("Helm",      "Package manager for Kubernetes; reference artifact type for OCM."),
    ("LoB",       "Line of Business — SAP organisational unit owning a product portfolio."),
    ("NeoNephos", "European foundation for sovereign cloud open-source projects, hosted under the Linux Foundation."),
    ("NIS2",      "Network and Information Security Directive 2 — EU baseline for cybersecurity of essential entities."),
    ("OCI",       "Open Container Initiative — open standards for container image format and distribution."),
    ("OCM",       "Open Component Model — vendor-neutral specification for signed, transportable software components."),
    ("ODG",       "Open Delivery Gear — OCM-native compliance automation engine and dashboard."),
    ("OSS",       "Open Source Software."),
    ("PKI",       "Public Key Infrastructure — framework for managing certificates and signing keys."),
    ("SBOD",      "Software Bill of Delivery — the OCM component descriptor, signed and traceable. Containing all artifacts and metadata for delivery and deployment."),
    ("SBOM",      "Software Bill of Materials — inventory of components and dependencies inside a software artifact."),
    ("SecNumCloud", "French cloud security qualification scheme operated by ANSSI."),
    ("Sigstore",  "Open-source project for keyless software signing using OIDC identities."),
    ("SPDX",      "Software Package Data Exchange — ISO/IEC 5962 standard format for SBOM data."),
    ("SWID",      "Software Identification Tags — ISO/IEC 19770-2 standard for software inventory."),
    ("Trivy",     "Open-source security scanner for containers, IaC, and code (Aqua Security)."),
]


def add_glossary_grid(slide, entries: list[tuple[str, str]],
                      y_top_px: int = 360):
    """Two-column grid of term / definition pairs. See build_pptx.py for
    full rationale; this is a verbatim mirror so the internal-sponsor
    deck has its own appendix without cross-importing."""
    n = len(entries)
    per_col = (n + 1) // 2
    col_gap_px = 60
    col_w_px = (SLIDE_W_PX - 240 - col_gap_px) // 2
    row_h_px = 40
    col_h_px = per_col * row_h_px

    for col in range(2):
        col_x = 120 + col * (col_w_px + col_gap_px)
        col_entries = entries[col * per_col : (col + 1) * per_col]
        if not col_entries:
            continue
        tb = slide.shapes.add_textbox(px(col_x), px(y_top_px),
                                       px(col_w_px), px(col_h_px))
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.word_wrap = True
        for i, (term, definition) in enumerate(col_entries):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(2)
            p.space_after = Pt(6)
            r1 = p.add_run()
            r1.text = term
            r1.font.name = "Aptos"
            r1.font.size = Pt(14)
            r1.font.bold = True
            r1.font.color.rgb = C.BLUE
            r2 = p.add_run()
            r2.text = f"  —  {definition}"
            r2.font.name = "Aptos"
            r2.font.size = Pt(13)
            r2.font.color.rgb = C.BLACK


def add_appendix_glossary_slide(prs, layouts):
    """Last visible slide: abbreviations and acronyms used in the deck.
    Plain / Compact layout — see external builder for rationale."""
    s = prs.slides.add_slide(layouts["Plain / Compact"])
    set_text(s, 1, "APPENDIX — ABBREVIATIONS")
    set_text(s, 2, "Quick reference for the acronyms used in this deck.")
    delete_placeholder(s, 10)
    add_glossary_grid(s, GLOSSARY_ENTRIES, y_top_px=520)


def add_hidden_trademark_slide(prs, layouts):
    """Append two non-presented slides with trademark/licensing notices.
    show="0" hides them in slideshow but they stay in the file. Mirrors the
    external deck's helper. Canonical record: assets/adopters/LICENSING.md."""
    # ---- Hidden 1/2 — adopter logos -----------------------------------------
    s = prs.slides.add_slide(layouts["Plain"])
    s.element.set("show", "0")
    set_text(s, 1, "TRADEMARK & LICENSE NOTICES (1/2)")
    set_text(s, 2, "Logos and trademarks belong to their respective owners.")
    set_blue_box_bullets(s, 10, [
        "SAP, SAP NS2 — trademarks of SAP SE / SAP National Security "
        "Services. Editorial use only; no endorsement implied. "
        "sap.com · sapns2.com",
        "BWI — trademark of BWI GmbH (Bundeswehr-IT). Editorial use of the "
        "Wikimedia public-domain wordmark; verify against BWI press "
        "conditions before external publication. bwi.de",
        "Gardener, Platform Mesh, NeoNephos Foundation — Linux Foundation "
        "Europe artwork; usage governed by the Linux Foundation trademark "
        "usage guidelines (linuxfoundation.org/legal/trademark-usage). "
        "gardener.cloud · platform-mesh.io · neonephos.org",
        "Konfidence — SAP-supported open project; logo from konfidence.cloud. "
        "Editorial use only; verify with the Konfidence project before "
        "external publication. konfidence.cloud",
        "OpenControlPlane — open-source project at open-control-plane.io. "
        "Editorial use only; verify with the project before external "
        "publication. open-control-plane.io",
    ])

    # ---- Hidden 2/2 — remaining marks + internal projects + sourcing -------
    s = prs.slides.add_slide(layouts["Plain"])
    s.element.set("show", "0")
    set_text(s, 1, "TRADEMARK & LICENSE NOTICES (2/2)")
    set_text(s, 2, "Third-party trademarks named for technical reference.")
    set_blue_box_bullets(s, 10, [
        "Kyma — SAP-originated open-source project at kyma-project.io. "
        "Editorial use only. kyma-project.io",
        "Hyperspace, RBSC, CSI, Greenhouse, Steampunk — internal SAP "
        "delivery infrastructure named for context; not third-party marks.",
        "Trivy, Grype, Sigstore, Helm, OCI, Kubernetes, kro, Flux, Argo CD — "
        "third-party trademarks named for technical reference; ownership "
        "remains with their respective projects and organisations.",
    ])


# =============================================================================
if __name__ == "__main__":
    if not shutil.which("rsvg-convert"):
        sys.exit("rsvg-convert not found; install via `brew install librsvg`")
    build()

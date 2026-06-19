"""
slide_3_native.py — Slide 3 "Meet OCM. One identity, every boundary."
rendered as native PowerPoint shapes (autoshapes, textboxes, freeform
Bezier paths), 1:1 with diagrams/03-meet-ocm-hub-and-spoke.svg.

Why native: SVG rasterization throws away typography and edit-in-PowerPoint
affordances. The hub-and-spoke composition is mostly autoshapes + curved
spokes + small flag/cloud-lock glyphs — all reachable from python-pptx
primitives plus a small bit of custGeom XML for the cubic Beziers.

Composition (mirrors the SVG viewBox 1760x560):

  HUB (centred at SVG (880, 250))
    - Two soft brand-blue halo circles (r=104 outer / r=88 inner)
    - "E-ticket" rounded rect: white body + light-grey (#F3F4F6) version
      sub-region, hairline #E5E7EB divider, brand-blue lock-badge circle,
      brand-blue + cyan name strips, "v1.0.0" label on the version stub.

  CLUSTER 1 — LEFT (artifact types)
    - 6 white pills with brand-blue stroke labelled
      OCI / Helm / npm / Binary / Config / "… any artifact type"
      (last one dashed-stroke italic).
    - Eyebrow "EVERY ARTIFACT TYPE" + brand-blue underline.
    - Solid #374151 cubic Bezier spokes, arrows landing on hub r=104.

  CLUSTER 2 — RIGHT (deployment boundaries)
    - 3 tiles: EU flag + "EU", US flag + "US", cloud-lock + "Sovereign Cloud".
    - Eyebrow "EVERY DEPLOYMENT BOUNDARY" + brand-blue underline.
    - Solid #374151 spokes flow OUTWARD from hub to tiles.

  CLUSTER 3 — BOTTOM (compliance frameworks)
    - 3 pills: DORA / NIS2 / CRA. Centred eyebrow above the pills.
    - Dashed #374151 spokes flow UPWARD from each pill into the hub.

  FOOTER
    - Italic mid-blue caption "plus FedRAMP/FISMA, BSI C5, SecNumCloud — and
      the regimes specific to your sector."

The whole 1760x560 SVG is mapped uniformly into the (x, y, w, h) slot the
builder hands us, centred horizontally and vertically.

Usage:
    from slide_3_native import add_hub_and_spoke_native_diagram
    add_hub_and_spoke_native_diagram(slide,
                                       x=60, y=240, w=1800, h=600,
                                       icons_dir=ICONS_DIR,
                                       rasterize_recolored=rasterize_svg_recolored)
"""
from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


PX = 9525  # EMU per px @ 96 dpi


def px(n: float) -> Emu:
    return Emu(int(n * PX))


# --- Palette (mirror of the SVG) -------------------------------------------
BLUE        = RGBColor(0x0F, 0x6B, 0xFF)   # brand-blue
BLUE_MID    = RGBColor(0x1D, 0x65, 0xB4)   # mid blue (pill strokes, body text accents)
BLUE_ACCENT = RGBColor(0x25, 0x7D, 0xDC)   # halo + left underline
BLUE_NIGHT  = RGBColor(0x0A, 0x15, 0x30)   # boundary tile body text + flag stroke
BLUE_LABEL  = RGBColor(0x0a, 0x3a, 0x99)   # "v1.0.0" stub label
CYAN        = RGBColor(0x5C, 0xD6, 0xFF)   # second name-strip on hub
SPOKE       = RGBColor(0x37, 0x41, 0x51)   # light-black spokes & arrowheads
GREY_TEXT   = RGBColor(0x6B, 0x72, 0x80)   # eyebrow grey
GREY_SOFT   = RGBColor(0xF3, 0xF4, 0xF6)   # version-stub fill
GREY_LINE   = RGBColor(0xE5, 0xE7, 0xEB)   # hairline divider
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)

# Flag colours
EU_BLUE   = RGBColor(0x00, 0x33, 0x99)
EU_YELLOW = RGBColor(0xFF, 0xCC, 0x00)
US_RED    = RGBColor(0xDC, 0x26, 0x26)
US_BLUE   = RGBColor(0x1D, 0x65, 0xB4)


# Source SVG canvas
SVG_W = 1760.0
SVG_H = 560.0
HUB_CX = 880.0  # hub centre in SVG coords (translate(802,220) + local (78,30))
HUB_CY = 250.0


# --- Type sizes -------------------------------------------------------------
# Editable-by-hand sizes — decoupled from the SVG-to-slide scale factor t.s
# so the slide stays legible and an editor can resize without dragging
# everything off-grid. The SVG was authored at 1760×560 where 9-10pt SVG
# glyphs read fine; rendered into a 1800×780 PowerPoint slot at scale ~1.0
# the same numbers come out 9-10pt in PowerPoint, which is unreadable for
# a presentation. Use these constants instead.
SZ_EYEBROW         = 11   # "EVERY ARTIFACT TYPE" / "EVERY DEPLOYMENT BOUNDARY" / "EVERY COMPLIANCE FRAMEWORK"
SZ_PILL_LABEL      = 12   # artifact-pills, compliance-pills (bold)
SZ_PILL_LABEL_SOFT = 11   # the "any artifact type" italic soft pill
SZ_TILE_LABEL      = 14   # boundary tiles ("EU", "US", "Sovereign Cloud")
SZ_FOOTER_NOTE     = 11   # "plus FedRAMP/FISMA, BSI C5..." italic line
SZ_VERSION_STUB    =  9   # "v1.0.0" badge — intentionally small (it's part of the icon)


# --- nsmap helpers ---------------------------------------------------------
NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NSMAP_A = {"a": NS_A}


def _qn(tag: str) -> str:
    """Qualified DrawingML tag name."""
    return f"{{{NS_A}}}{tag}"


# ---------------------------------------------------------------------------
# Generic primitives
# ---------------------------------------------------------------------------
def _styled_run(p, text, *, size_pt, bold=False, color=BLACK,
                font="Inter", italic=False, all_caps=False,
                letter_spacing=None):
    """Add a run with the OCM type system applied."""
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size_pt)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if all_caps or letter_spacing is not None:
        rPr = r._r.get_or_add_rPr()
        if all_caps:
            rPr.set("cap", "all")
        if letter_spacing is not None:
            rPr.set("spc", str(letter_spacing))
    return r


def _add_textbox(slide, x, y, w, h, *, anchor="t", align=None, margins=0,
                 fit_text=False):
    """Borderless textbox at the given coords, with consistent inset margins.

    fit_text=True enables PowerPoint's "Resize shape to fit text" so the box
    auto-shrinks (or grows) to its content. Use this for editable-by-hand
    labels where the SVG-derived box dims are wildly oversized for the
    actual text — without it, click-to-edit grabs an enormous selection
    rectangle around a short word.
    """
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = px(margins)
    tf.margin_top = tf.margin_bottom = px(margins)
    tf.word_wrap = True
    if fit_text:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    if anchor == "ctr":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "b":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    if align is not None:
        tf.paragraphs[0].alignment = align
    return tf


def _add_rounded_rect(slide, x, y, w, h, *,
                      fill=WHITE, stroke=BLUE_MID, stroke_pt=1.5,
                      corner_pct=0.04, dashed=False):
    """Rounded rectangle with explicit fill, stroke, and corner radius."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h),
    )
    shape.adjustments[0] = corner_pct
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = stroke
        shape.line.width = Pt(stroke_pt)
        if dashed:
            ln = shape.line._get_or_add_ln()
            # remove any existing prstDash
            for d in ln.findall(_qn("prstDash")):
                ln.remove(d)
            dash = etree.SubElement(ln, _qn("prstDash"))
            dash.set("val", "dash")
    return shape


def _add_rect(slide, x, y, w, h, *, fill=WHITE, stroke=None, stroke_pt=0.75):
    """Plain rectangle with explicit fill and optional stroke."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = stroke
        shape.line.width = Pt(stroke_pt)
    return shape


def _add_oval(slide, cx, cy, r, *, fill, stroke=None, stroke_pt=0.75,
              fill_alpha=None):
    """Circle (oval) centred on (cx, cy) with radius r in slide-px coords."""
    x, y = cx - r, cy - r
    w = h = r * 2
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, px(x), px(y), px(w), px(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        if fill_alpha is not None:
            # Inject alpha on the solidFill
            sf = shape.fill._xPr.find(_qn("solidFill"))
            if sf is not None:
                clr = sf.find(_qn("srgbClr"))
                if clr is not None:
                    # alpha as percentage * 1000 (0..100000)
                    a = etree.SubElement(clr, _qn("alpha"))
                    a.set("val", str(int(fill_alpha * 100000)))
    if stroke is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = stroke
        shape.line.width = Pt(stroke_pt)
    return shape


def _add_line(slide, x1, y1, x2, y2, *, color=BLACK, width_pt=1.0):
    """Straight line connector."""
    from pptx.enum.shapes import MSO_CONNECTOR
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, px(x1), px(y1), px(x2), px(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    return line


# ---------------------------------------------------------------------------
# Cubic Bezier path → custGeom shape (with optional arrow end + dashed stroke)
# ---------------------------------------------------------------------------
def _add_cubic_bezier_path(slide, segments, *,
                           color=SPOKE, width_pt=1.5,
                           dashed=False, arrow_end=True):
    """Add a cubic-bezier path as a custGeom shape.

    Parameters:
        segments: list of (x0, y0, c1x, c1y, c2x, c2y, x1, y1) tuples — each
            segment is a single cubic Bezier from (x0,y0) to (x1,y1) with
            control points (c1x,c1y) and (c2x,c2y). Coordinates are in
            slide-px (NOT viewBox-px). For multi-segment paths, segments are
            chained by python-pptx implicitly via repeated cubicBezTo
            elements; if your SVG has a single C command, pass one tuple.
        color, width_pt: stroke styling.
        dashed: render as a dashed stroke (used for compliance-cluster spokes).
        arrow_end: append an arrow head at the path terminus.

    Implementation: we emit a custGeom autoshape whose pathLst contains a
    moveTo + cubicBezTo per segment, with bounding box derived from the
    extreme x/y across all anchor and control points. python-pptx doesn't
    have a high-level cubic-bezier builder (FreeformBuilder is line-segment
    only), so this is the lowest-friction path.
    """
    if not segments:
        raise ValueError("segments must be non-empty")

    # Compute bbox covering anchors + controls so the custGeom rectangle
    # contains the entire curve.
    xs, ys = [], []
    for (x0, y0, c1x, c1y, c2x, c2y, x1, y1) in segments:
        xs.extend([x0, c1x, c2x, x1])
        ys.extend([y0, c1y, c2y, y1])
    min_x, max_x = min(xs), max(xs)
    min_y, max_y = min(ys), max(ys)
    bw = max(max_x - min_x, 1)
    bh = max(max_y - min_y, 1)

    # Use a freeform with a placeholder line then rewrite the path. Simpler:
    # build the sp XML by hand and append it directly to the slide spTree.
    spTree = slide.shapes._spTree

    nsdecl = (
        ' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="{NS_A}"'
    )

    # Path coord-space units. Use EMU directly — we already have full bw/bh
    # in slide-px and convert to EMU via px().
    path_w = int(bw * PX)
    path_h = int(bh * PX)
    off_x = int(min_x * PX)
    off_y = int(min_y * PX)
    ext_cx = int(bw * PX)
    ext_cy = int(bh * PX)

    def _emu_pt(x, y):
        return f'<a:pt x="{int((x - min_x) * PX)}" y="{int((y - min_y) * PX)}"/>'

    path_xml_parts = [f'<a:path w="{path_w}" h="{path_h}">']
    # moveTo first segment start
    x0, y0 = segments[0][0], segments[0][1]
    path_xml_parts.append(f"<a:moveTo>{_emu_pt(x0, y0)}</a:moveTo>")
    for (x0, y0, c1x, c1y, c2x, c2y, x1, y1) in segments:
        path_xml_parts.append(
            "<a:cubicBezTo>"
            f"{_emu_pt(c1x, c1y)}"
            f"{_emu_pt(c2x, c2y)}"
            f"{_emu_pt(x1, y1)}"
            "</a:cubicBezTo>"
        )
    path_xml_parts.append("</a:path>")
    path_xml = "".join(path_xml_parts)

    rgb_hex = "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2])

    dash_xml = '<a:prstDash val="dash"/>' if dashed else ""
    arrow_xml = (
        '<a:tailEnd type="triangle" w="med" len="med"/>' if arrow_end else ""
    )

    sp_xml = (
        f'<p:sp{nsdecl}>'
        '<p:nvSpPr>'
        f'<p:cNvPr id="Spoke_placeholder" name="Spoke"/>'
        '<p:cNvSpPr/>'
        '<p:nvPr/>'
        '</p:nvSpPr>'
        '<p:spPr>'
        f'<a:xfrm><a:off x="{off_x}" y="{off_y}"/><a:ext cx="{ext_cx}" cy="{ext_cy}"/></a:xfrm>'
        '<a:custGeom>'
        '<a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/>'
        '<a:rect l="l" t="t" r="r" b="b"/>'
        f'<a:pathLst>{path_xml}</a:pathLst>'
        '</a:custGeom>'
        '<a:noFill/>'
        f'<a:ln w="{int(width_pt * 12700)}" cap="rnd">'
        f'<a:solidFill><a:srgbClr val="{rgb_hex}"/></a:solidFill>'
        f'{dash_xml}'
        '<a:round/>'
        f'{arrow_xml}'
        '</a:ln>'
        '</p:spPr>'
        '<p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr/></a:p></p:txBody>'
        '</p:sp>'
    )

    # Assign a unique positive id — PowerPoint rejects id=0 and duplicates.
    new_shape_id = max(
        (int(el.get("id", 0)) for el in spTree.iter() if el.tag.endswith("}cNvPr")),
        default=0,
    ) + 1
    sp_xml = sp_xml.replace('id="Spoke_placeholder"', f'id="{new_shape_id}"')

    sp_el = etree.fromstring(sp_xml)
    spTree.append(sp_el)
    return sp_el


# ---------------------------------------------------------------------------
# Hub composition
# ---------------------------------------------------------------------------
def _draw_hub(slide, t):
    """Draw the OCM-component hub (halo + e-ticket icon + lock badge + name
    strips + version label). `t` is a coordinate transformer: t(svg_x, svg_y)
    → slide-px (x, y).  t.s scales an SVG length to slide-px."""
    cx, cy = t(HUB_CX, HUB_CY)

    # --- Halo: two soft brand-accent circles ------------------------------
    _add_oval(slide, cx, cy, 104 * t.s,
              fill=BLUE_ACCENT, stroke=None, fill_alpha=0.06)
    _add_oval(slide, cx, cy, 88 * t.s,
              fill=BLUE_ACCENT, stroke=None, fill_alpha=0.10)

    # --- E-ticket icon -----------------------------------------------------
    # Icon is anchored at hub-local origin (translate(802,220)) and spans
    # x=0..156, y=0..60 in that local frame. Hub-local centre is (78,30).
    icon_origin_svg = (802, 220)
    ix, iy = t(*icon_origin_svg)
    icon_w = 156 * t.s
    icon_h = 60 * t.s

    # Stub fill (light grey, stub area x=100..156)
    stub_x = ix + 100 * t.s
    _add_rect(slide, stub_x, iy, 56 * t.s, icon_h,
              fill=GREY_SOFT, stroke=None)
    # Body white panel (x=0..100), no stroke; the outer rounded-rect carries
    # the stroke for the whole 156-wide envelope.
    _add_rect(slide, ix, iy, 100 * t.s, icon_h,
              fill=WHITE, stroke=None)
    # Outer enclosing rounded rect with the unified stroke; transparent fill
    # so the body+stub colours show through.
    _add_rounded_rect(slide, ix, iy, icon_w, icon_h,
                      fill=None, stroke=BLUE_MID, stroke_pt=2.2,
                      corner_pct=10.0 / 60.0)
    # Hairline divider between body and stub at x=100
    _add_line(slide, ix + 100 * t.s, iy + 6 * t.s,
              ix + 100 * t.s, iy + 54 * t.s,
              color=GREY_LINE, width_pt=1.2 * t.s)

    # Lock badge (filled brand-blue circle at hub-local (24,30), r=14).
    lx, ly = t(802 + 24, 220 + 30)
    _add_oval(slide, lx, ly, 14 * t.s, fill=BLUE_MID, stroke=None)
    # Lock glyph: white outline lock (body rect + shackle arc + dot).
    # Approximated with a few simple primitives — the SVG draws three small
    # paths, we replicate the body+shackle silhouette with a rounded rect
    # for the body and a thin arc-ish shape for the shackle.
    body_w = 11 * t.s
    body_h = 7 * t.s
    body_x = lx - body_w / 2
    body_y = ly - body_h / 2 + 2 * t.s
    _add_rounded_rect(slide, body_x, body_y, body_w, body_h,
                      fill=WHITE, stroke=None, corner_pct=0.15)
    # Shackle: a thin horseshoe — render as two stroked vertical lines and a
    # connecting arc would require more freeform. Use a small rounded rect
    # outline above the body which reads as a shackle at this scale.
    shackle_w = 7 * t.s
    shackle_h = 6 * t.s
    shackle_x = lx - shackle_w / 2
    shackle_y = ly - body_h / 2 - shackle_h + 3.5 * t.s
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 px(shackle_x), px(shackle_y),
                                 px(shackle_w), px(shackle_h))
    sh.adjustments[0] = 0.45
    sh.fill.background()
    sh.line.color.rgb = WHITE
    sh.line.width = Pt(1.6 * t.s)
    # Keyhole dot (small brand-blue circle on the body)
    _add_oval(slide, lx, ly + 1 * t.s, 1.2 * t.s, fill=BLUE_MID, stroke=None)

    # Name strips on body (abstract bars).
    s1x, s1y = t(802 + 46, 220 + 22)
    _add_rounded_rect(slide, s1x, s1y, 44 * t.s, 5 * t.s,
                      fill=BLUE_MID, stroke=None, corner_pct=0.4)
    s2x, s2y = t(802 + 46, 220 + 34)
    _add_rounded_rect(slide, s2x, s2y, 32 * t.s, 5 * t.s,
                      fill=CYAN, stroke=None, corner_pct=0.4)

    # Version label "v1.0.0" — centred at hub-local (128, 36) on the stub.
    # Use a textbox roughly spanning the stub width and centred.
    vw = 56 * t.s
    vh = 20 * t.s
    vx = ix + (100 + 28 - 56 / 2) * t.s
    vy = iy + (36 - 14) * t.s
    tf = _add_textbox(slide, vx, vy, vw, vh, anchor="ctr",
                       align=PP_ALIGN.CENTER)
    p = tf.paragraphs[0]
    _styled_run(p, "v1.0.0", size_pt=SZ_VERSION_STUB,
                bold=True, color=BLUE_LABEL, letter_spacing="100")


# ---------------------------------------------------------------------------
# Cluster pieces
# ---------------------------------------------------------------------------
def _draw_artifact_pills(slide, t):
    """Cluster 1 — left pills + spokes."""
    # Eyebrow + underline. Tight box around the text so click-to-edit grabs
    # only what an editor expects (the eyebrow phrase), not 200 px of empty
    # space. fit_text auto-sizes the height to the actual line height of
    # SZ_EYEBROW.
    ex, ey = t(240, 38)
    tf = _add_textbox(slide, ex, ey - 14, 200, 22, fit_text=True)
    _styled_run(tf.paragraphs[0], "EVERY ARTIFACT TYPE",
                size_pt=SZ_EYEBROW, bold=True, color=GREY_TEXT,
                all_caps=True, letter_spacing="140")
    # Underline (colour matches SVG: #257DDC for the artifact eyebrow)
    ux1, uy = t(240, 50)
    ux2, _ = t(420, 50)
    _add_line(slide, ux1, uy, ux2, uy, color=BLUE_ACCENT, width_pt=2 * t.s)

    # Pill geometry: bumped to 44 px high (was 36) so 12pt text breathes
    # inside the rounded rectangle without crowding the corners. Spacing
    # between pills bumped from 52 to 54 SVG-units to keep the dashed
    # bottom pill from touching its neighbour at the new height.
    pill_w_svg = 160
    pill_h_svg = 44
    pills = [
        (62,  "OCI",                False),
        (116, "Helm",               False),
        (170, "npm",                False),
        (224, "Binary",             False),
        (278, "Config",             False),
        (332, "… any artifact type", True),
    ]
    for (sy, label, soft) in pills:
        x_px, y_px = t(240, sy)
        w_px = pill_w_svg * t.s
        h_px = pill_h_svg * t.s
        _add_rounded_rect(slide, x_px, y_px, w_px, h_px,
                          fill=WHITE, stroke=BLUE_MID, stroke_pt=1.5,
                          corner_pct=0.5, dashed=soft)
        tf = _add_textbox(slide, x_px, y_px, w_px, h_px,
                          anchor="ctr", align=PP_ALIGN.CENTER)
        p = tf.paragraphs[0]
        _styled_run(p, label,
                    size_pt=SZ_PILL_LABEL_SOFT if soft else SZ_PILL_LABEL,
                    bold=not soft,
                    italic=soft,
                    color=BLUE_MID)

    # Spokes — 6 cubic Bezier curves, anchors mirror the SVG path data.
    art_paths = [
        # (x0,y0,c1x,c1y,c2x,c2y,x1,y1)
        (400, 80,  560, 80,  640, 200, 779.3, 224),
        (400, 132, 560, 132, 640, 220, 776.5, 240),
        (400, 184, 560, 184, 660, 240, 776.0, 252),
        (400, 236, 560, 236, 660, 256, 777.0, 264),
        (400, 288, 560, 288, 660, 274, 779.3, 276),
    ]
    for seg in art_paths:
        seg_px = tuple(t.scale_pair(seg[i], seg[i + 1])
                       for i in range(0, len(seg), 2))
        flat = []
        for (x, y) in seg_px:
            flat.extend([x, y])
        _add_cubic_bezier_path(slide, [tuple(flat)],
                                color=SPOKE, width_pt=1.5 * t.s,
                                dashed=False, arrow_end=True)

    # Soft (dashed) "any artifact type" spoke
    soft_seg = (400, 340, 560, 340, 640, 305, 784, 290)
    seg_px = tuple(t.scale_pair(soft_seg[i], soft_seg[i + 1])
                   for i in range(0, len(soft_seg), 2))
    flat = []
    for (x, y) in seg_px:
        flat.extend([x, y])
    _add_cubic_bezier_path(slide, [tuple(flat)],
                            color=SPOKE, width_pt=1.3 * t.s,
                            dashed=True, arrow_end=True)


def _draw_eu_flag(slide, x, y, t):
    """Draw the EU flag glyph (36×24 SVG units) at slide-px (x, y)."""
    s = t.s
    _add_rect(slide, x, y, 36 * s, 24 * s,
              fill=EU_BLUE, stroke=BLUE_NIGHT, stroke_pt=0.75 * s)
    # 12 yellow dots, ring centred (18,12), r=8
    star_centres = [
        (18, 4), (22, 5.07), (24.93, 8), (26, 12), (24.93, 16),
        (22, 18.93), (18, 20), (14, 18.93), (11.07, 16), (10, 12),
        (11.07, 8), (14, 5.07),
    ]
    for (sx, sy) in star_centres:
        _add_oval(slide, x + sx * s, y + sy * s, 1.4 * s,
                  fill=EU_YELLOW, stroke=None)


def _draw_us_flag(slide, x, y, t):
    """Draw the US flag glyph (36×24 SVG units) at slide-px (x, y)."""
    s = t.s
    # White base rect (with night-blue 0.75 stroke)
    _add_rect(slide, x, y, 36 * s, 24 * s,
              fill=WHITE, stroke=BLUE_NIGHT, stroke_pt=0.75 * s)
    # 7 red stripes — every other 1.85 unit band
    stripe_h = 1.85
    for i in range(7):
        sy = i * 3.7
        _add_rect(slide, x, y + sy * s, 36 * s, stripe_h * s,
                  fill=US_RED, stroke=None)
    # Canton (blue field over upper-left)
    _add_rect(slide, x, y, 14.4 * s, 12.92 * s,
              fill=US_BLUE, stroke=None)
    # 4×3 white "stars" (dots)
    star_xs = [2.5, 6, 9.5, 12.5]
    star_ys = [2.4, 6.5, 10.5]
    for sy in star_ys:
        for sx in star_xs:
            _add_oval(slide, x + sx * s, y + sy * s, 0.7 * s,
                      fill=WHITE, stroke=None)


def _draw_cloud_lock(slide, x, y, t):
    """Cloud-lock glyph (~60×40 envelope) at slide-px (x, y).

    The SVG path uses a single cubic-Bezier cloud silhouette with control
    points outside the envelope — we approximate the silhouette here with a
    stack of overlapping ovals (4 blob circles + a base rectangle), which is
    the cleanest python-pptx-only rendering and reads identically at slide
    scale. The lock atop is drawn from primitives matching the SVG.
    """
    s = t.s
    # Cloud silhouette: base rectangle + 3 ovals make a serviceable cloud.
    # Envelope reference: SVG cloud spans roughly x=0..60, y=-5..32.
    # Base flat: y=18..32, x=0..60.
    _add_rounded_rect(slide, x, y + 18 * s, 60 * s, 14 * s,
                      fill=BLUE, stroke=None, corner_pct=0.45)
    # Left puff
    _add_oval(slide, x + 12 * s, y + 18 * s, 14 * s,
              fill=BLUE, stroke=None)
    # Centre puff (taller)
    _add_oval(slide, x + 30 * s, y + 12 * s, 18 * s,
              fill=BLUE, stroke=None)
    # Right puff
    _add_oval(slide, x + 47 * s, y + 18 * s, 13 * s,
              fill=BLUE, stroke=None)

    # Lock body (white rounded rect) — SVG: x=22 y=19 w=16 h=11
    _add_rounded_rect(slide, x + 22 * s, y + 19 * s, 16 * s, 11 * s,
                      fill=WHITE, stroke=None, corner_pct=0.10)
    # Lock shackle: thin white rounded-rect outline above the body
    sh_w = 10 * s
    sh_h = 10 * s
    sh_x = x + (30 - 5) * s
    sh_y = y + (19 - 8) * s
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 px(sh_x), px(sh_y), px(sh_w), px(sh_h))
    sh.adjustments[0] = 0.45
    sh.fill.background()
    sh.line.color.rgb = WHITE
    sh.line.width = Pt(2 * s)
    # Keyhole — small blue dot + thin tail
    _add_oval(slide, x + 30 * s, y + 24 * s, 1.2 * s,
              fill=BLUE, stroke=None)
    _add_rect(slide, x + 29.4 * s, y + 24 * s, 1.2 * s, 3.6 * s,
              fill=BLUE, stroke=None)


def _draw_boundary_tiles(slide, t):
    """Cluster 2 — right tiles + spokes."""
    # Eyebrow + underline (centre x = 1470 in SVG). fit_text shrinks the
    # textbox to the actual glyph extents so the editor doesn't grab a
    # 180-px-wide invisible box.
    eyx, eyy = t(1380, 38)
    tf = _add_textbox(slide, eyx, eyy - 14, 220, 22,
                       align=PP_ALIGN.CENTER, fit_text=True)
    _styled_run(tf.paragraphs[0], "EVERY DEPLOYMENT BOUNDARY",
                size_pt=SZ_EYEBROW, bold=True, color=GREY_TEXT,
                all_caps=True, letter_spacing="140")
    ux1, uy = t(1380, 50)
    ux2, _ = t(1560, 50)
    _add_line(slide, ux1, uy, ux2, uy, color=BLUE_MID, width_pt=2 * t.s)

    tile_w_svg = 180
    tile_h_svg = 80   # bumped from 70 so 14pt label + flag glyph breathe
    tiles_y = [80, 200, 320]
    for ty_svg in tiles_y:
        x_px, y_px = t(1380, ty_svg)
        _add_rounded_rect(slide,
                           x_px, y_px,
                           tile_w_svg * t.s, tile_h_svg * t.s,
                           fill=WHITE, stroke=BLUE_MID, stroke_pt=1.5,
                           corner_pct=10.0 / tile_h_svg)

    # EU tile — flag at SVG (1380+53, 80+23), label centred at (1380+103, 80+42)
    fx, fy = t(1380 + 53, 80 + 28)
    _draw_eu_flag(slide, fx, fy, t)
    lx, ly = t(1380 + 88, 80 + 30)
    tf = _add_textbox(slide, lx, ly, 80 * t.s, 28,
                       anchor="ctr", align=PP_ALIGN.LEFT, fit_text=True)
    _styled_run(tf.paragraphs[0], "EU",
                size_pt=SZ_TILE_LABEL, bold=True, color=BLUE_NIGHT)

    # US tile
    fx, fy = t(1380 + 53, 200 + 28)
    _draw_us_flag(slide, fx, fy, t)
    lx, ly = t(1380 + 88, 200 + 30)
    tf = _add_textbox(slide, lx, ly, 80 * t.s, 28,
                       anchor="ctr", align=PP_ALIGN.LEFT, fit_text=True)
    _styled_run(tf.paragraphs[0], "US",
                size_pt=SZ_TILE_LABEL, bold=True, color=BLUE_NIGHT)

    # Sovereign Cloud tile — two-line label so the box must accommodate
    # both lines at SZ_TILE_LABEL. fit_text would clip the second line, so
    # leave the box explicitly tall enough.
    cx_, cy_ = t(1380 + 14, 320 + 24)
    _draw_cloud_lock(slide, cx_, cy_, t)
    lx, ly = t(1380 + 80, 320 + 18)
    tf = _add_textbox(slide, lx, ly, 100 * t.s, 56, anchor="ctr",
                       align=PP_ALIGN.LEFT)
    _styled_run(tf.paragraphs[0], "Sovereign",
                size_pt=SZ_TILE_LABEL, bold=True, color=BLUE_NIGHT)
    p2 = tf.add_paragraph()
    _styled_run(p2, "Cloud",
                size_pt=SZ_TILE_LABEL, bold=True, color=BLUE_NIGHT)

    # Spokes — hub → tiles, arrows pointing OUTWARD onto the tile left edge.
    bnd_paths = [
        (970.07, 198, 1080, 180, 1220, 130, 1378, 115),
        (984,    250, 1080, 240, 1220, 230, 1378, 240),
        (970.07, 302, 1080, 320, 1220, 340, 1378, 360),
    ]
    for seg in bnd_paths:
        seg_px = tuple(t.scale_pair(seg[i], seg[i + 1])
                       for i in range(0, len(seg), 2))
        flat = []
        for (x, y) in seg_px:
            flat.extend([x, y])
        _add_cubic_bezier_path(slide, [tuple(flat)],
                                color=SPOKE, width_pt=1.75 * t.s,
                                dashed=False, arrow_end=True)


def _draw_compliance_pills(slide, t):
    """Cluster 3 — bottom pills + spokes."""
    # Eyebrow centred at (880, 520) + underline 750..1010. fit_text on the
    # eyebrow textbox so the click target is the headline phrase, not 260
    # px of empty centred space.
    ex, ey = t(880, 520)
    tf = _add_textbox(slide, ex - 130 * t.s, ey - 14,
                       260 * t.s, 22,
                       align=PP_ALIGN.CENTER, fit_text=True)
    _styled_run(tf.paragraphs[0], "EVERY COMPLIANCE FRAMEWORK",
                size_pt=SZ_EYEBROW, bold=True, color=GREY_TEXT,
                all_caps=True, letter_spacing="140")
    ux1, uy = t(750, 528)
    ux2, _ = t(1010, 528)
    _add_line(slide, ux1, uy, ux2, uy, color=BLUE_MID, width_pt=2 * t.s)

    # Pill height bumped 40 → 48 SVG-units so the 12pt bold label has
    # room above and below the baseline. Pill width unchanged — 180 svg
    # × t.s ≈ 184 slide-px, plenty for "DORA"/"NIS2"/"CRA".
    pill_w_svg = 180
    pill_h_svg = 48
    items = [
        (400,  "DORA"),  # centre x ≈ 490
        (790,  "NIS2"),  # centre x ≈ 880
        (1180, "CRA"),   # centre x ≈ 1270
    ]
    for (sx, label) in items:
        x_px, y_px = t(sx, 451)
        w_px = pill_w_svg * t.s
        h_px = pill_h_svg * t.s
        _add_rounded_rect(slide, x_px, y_px, w_px, h_px,
                          fill=WHITE, stroke=BLUE_MID, stroke_pt=1.5,
                          corner_pct=0.5)
        tf = _add_textbox(slide, x_px, y_px, w_px, h_px,
                           anchor="ctr", align=PP_ALIGN.CENTER)
        _styled_run(tf.paragraphs[0], label,
                    size_pt=SZ_PILL_LABEL, bold=True, color=BLUE_MID)

    # Spokes — pills → hub, arrows pointing UP into the hub. Dashed.
    cmp_paths = [
        (490,  451, 600, 410, 740, 360, 806.46, 323.54),
        (880,  451, 880, 420, 880, 380, 880,    354),
        (1270, 451, 1160, 410, 1020, 360, 953.54, 323.54),
    ]
    for seg in cmp_paths:
        seg_px = tuple(t.scale_pair(seg[i], seg[i + 1])
                       for i in range(0, len(seg), 2))
        flat = []
        for (x, y) in seg_px:
            flat.extend([x, y])
        _add_cubic_bezier_path(slide, [tuple(flat)],
                                color=SPOKE, width_pt=1.5 * t.s,
                                dashed=True, arrow_end=True)


def _draw_footer(slide, t):
    """Italic mid-blue caption at SVG (880, 552), centre-aligned."""
    fx, fy = t(880, 552)
    tf = _add_textbox(slide, fx - 600 * t.s, fy - 12,
                       1200 * t.s, 26, align=PP_ALIGN.CENTER,
                       fit_text=True)
    _styled_run(tf.paragraphs[0],
                "plus FedRAMP/FISMA, BSI C5, SecNumCloud — and the regimes specific to your sector.",
                size_pt=SZ_FOOTER_NOTE, italic=True, color=BLUE_MID,
                letter_spacing="20")


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------
class _Transformer:
    """Maps SVG (1760×560) coordinates uniformly into the slide-px slot."""

    def __init__(self, x, y, w, h):
        s_x = w / SVG_W
        s_y = h / SVG_H
        self.s = min(s_x, s_y)  # uniform scale, fit inside slot
        # Centre the scaled diagram in the slot.
        used_w = SVG_W * self.s
        used_h = SVG_H * self.s
        self.dx = x + (w - used_w) / 2
        self.dy = y + (h - used_h) / 2

    def __call__(self, svg_x, svg_y):
        return (self.dx + svg_x * self.s, self.dy + svg_y * self.s)

    def scale_pair(self, svg_x, svg_y):
        return self(svg_x, svg_y)


def add_hub_and_spoke_native_diagram(slide, *,
                                      x: float, y: float,
                                      w: float, h: float,
                                      icons_dir: Path,
                                      rasterize_recolored):
    """Render the slide-3 hub-and-spoke diagram natively into (x, y, w, h).

    Parameters:
        slide: python-pptx Slide.
        x, y, w, h: diagram slot in slide-px (96 dpi).
        icons_dir: directory of Tabler-style SVG icons. Currently unused —
            the hub icon and the cloud-lock are constructed from primitives
            (the SVG composes them by hand too) — but the parameter is kept
            so the signature matches the rest of the native-diagram modules.
        rasterize_recolored: callable (svg_path, target_w_px, color_hex) →
            Path. Same shape as slide_4b_native; unused here, kept for
            signature parity.

    Layout: the SVG viewBox 1760×560 is uniformly scaled into the slot and
    centred. All primitives below are expressed in SVG coords and pushed
    through the transformer.
    """
    t = _Transformer(x, y, w, h)

    # Order matters: spokes first (so pills/tiles overlay arrowheads at
    # their endpoints — hub circle is bigger than any anchor anyway, and
    # the artifact-pill right edge sits at x=400 which is the spoke start).
    _draw_artifact_pills(slide, t)
    _draw_boundary_tiles(slide, t)
    _draw_compliance_pills(slide, t)
    _draw_hub(slide, t)
    _draw_footer(slide, t)

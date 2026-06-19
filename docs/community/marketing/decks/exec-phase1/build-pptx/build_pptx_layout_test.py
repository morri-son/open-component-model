#!/usr/bin/env python3
"""
build_pptx_layout_test.py — A/B test deck for diagram-slide geometry.

Builds OCM-Layout-Test.pptx with 10 slides: 5 variants × 2 reference diagrams.
Lets you flip through the variants in PowerPoint and pick the one that
gives the best diagram-readability without wrecking the title rhythm.

Variants (Eyebrow 28pt and Title 64pt unchanged in size — only positions move):

  A — Status quo:    Eyebrow y=255, Title y=308, Diagram y=460 h=600
  B — Moderate:      Eyebrow y=135, Title y=188, Diagram y=340 h=680  (+13%)
  C — Recommended:   Eyebrow y= 75, Title y=128, Diagram y=280 h=740  (+23%)
  D — Maximum:       Eyebrow y= 40, Title y= 93, Diagram y=240 h=780  (+30%)
  E — Compact title: Eyebrow y= 40, Title y= 78 (48pt 1-line), Diagram y=200 h=820  (+37%)

Each test slide carries a footer caption "Variant X — diagram NxM, ratio X:Y"
so you can identify which variant you're looking at while flipping through.

Reference diagrams:
  - 03-meet-ocm-hub-and-spoke.svg  (1760×560, ratio 3.14:1, very wide)
  - 05-pack-sign-transport-deploy-v2.svg (1920×540, ratio 3.56:1, very wide)

Usage:
    .venv/bin/python build_pptx_layout_test.py
"""
from __future__ import annotations

import shutil
import subprocess
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt


SCRIPT_DIR = Path(__file__).resolve().parent
DECK_DIR = SCRIPT_DIR.parent
DIAGRAMS_DIR = DECK_DIR / "diagrams"
RASTER_DIR = SCRIPT_DIR / "_raster"
POTX_PATH = DECK_DIR / "OCM-Master.potx"
OUTPUT_PPTX = DECK_DIR / "OCM-Layout-Test.pptx"

RASTER_DIR.mkdir(exist_ok=True)


SLIDE_W_PX = 1920
SLIDE_H_PX = 1080
PX = 9525


def px(n: float) -> Emu:
    return Emu(int(n * PX))


class C:
    BLUE       = RGBColor(0x0F, 0x6B, 0xFF)
    BLUE_MID   = RGBColor(0x0A, 0x3A, 0x99)
    GREY_MID   = RGBColor(0x6B, 0x72, 0x80)
    BLACK      = RGBColor(0x00, 0x00, 0x00)


def rasterize_svg(svg_path: Path, target_w_px: int) -> Path:
    out = RASTER_DIR / (svg_path.stem + f"_{target_w_px}.png")
    if out.exists() and out.stat().st_mtime >= svg_path.stat().st_mtime:
        return out
    subprocess.run(
        ["rsvg-convert", "--width", str(target_w_px), "--keep-aspect-ratio",
         str(svg_path), "-o", str(out)],
        check=True, capture_output=True,
    )
    return out


def open_template_as_pptx() -> Presentation:
    tmp_pptx = RASTER_DIR / "_potx_loaded_test.pptx"
    with zipfile.ZipFile(POTX_PATH, "r") as src:
        data = {n: src.read(n) for n in src.namelist()}
    ct = data["[Content_Types].xml"].decode("utf-8")
    ct = ct.replace(
        "presentationml.template.main+xml",
        "presentationml.presentation.main+xml",
    )
    data["[Content_Types].xml"] = ct.encode("utf-8")
    with zipfile.ZipFile(tmp_pptx, "w", zipfile.ZIP_DEFLATED) as out:
        for name, blob in data.items():
            out.writestr(name, blob)
    return Presentation(str(tmp_pptx))


def find_placeholder(slide, idx: int):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(f"placeholder idx={idx} not found on layout {slide.slide_layout.name!r}")


def delete_placeholder(slide, idx: int):
    try:
        ph = find_placeholder(slide, idx)
        ph._element.getparent().remove(ph._element)
    except KeyError:
        pass


def add_eyebrow(slide, x_px: int, y_px: int, w_px: int, h_px: int, text: str):
    """Brand-blue ALL-CAPS 28pt eyebrow, letter-spaced. Mirrors the master
    layout's eyebrow style so test slides match production slides visually."""
    tb = slide.shapes.add_textbox(px(x_px), px(y_px), px(w_px), px(h_px))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = "Aptos"
    f.size = Pt(28)
    f.bold = True
    f.color.rgb = C.BLUE
    rPr = r._r.get_or_add_rPr()
    rPr.set("cap", "all")
    rPr.set("spc", "140")


def add_title(slide, x_px: int, y_px: int, w_px: int, h_px: int,
              text: str, size_pt: int = 64):
    tb = slide.shapes.add_textbox(px(x_px), px(y_px), px(w_px), px(h_px))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.line_spacing = 0.9
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = "Aptos Display"
    f.size = Pt(size_pt)
    f.bold = True
    f.color.rgb = C.BLACK


def add_variant_caption(slide, text: str):
    """Small grey caption at the very top-left so we can ID the variant
    while flipping through the deck. Lives outside the eyebrow/title zone
    so it never collides with the layout we're testing."""
    tb = slide.shapes.add_textbox(px(20), px(20), px(800), px(28))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(12)
    r.font.italic = True
    r.font.color.rgb = C.GREY_MID


def add_diagram_centered(slide, svg_path: Path,
                          slot_x: int, slot_y: int,
                          slot_w: int, slot_h: int):
    """Place the SVG inside (slot_x, slot_y, slot_w, slot_h), preserving
    the SVG's intrinsic aspect ratio and centring within the slot. Returns
    the rendered (w, h) so the caller can log the actual diagram footprint."""
    png = rasterize_svg(svg_path, target_w_px=slot_w)
    # Probe natural size by adding off-slide and reading height.
    pic = slide.shapes.add_picture(str(png), px(0), px(0), width=px(slot_w))
    if pic.height > px(slot_h):
        ratio = px(slot_h) / pic.height
        pic.height = px(slot_h)
        pic.width = int(pic.width * ratio)
    final_w_px = pic.width / PX
    final_h_px = pic.height / PX
    pic.left = px(slot_x + (slot_w - final_w_px) / 2)
    pic.top = px(slot_y + (slot_h - final_h_px) / 2)
    return (final_w_px, final_h_px)


# Five variants — same Eyebrow font/size and same Title font/size unless noted.
# (eyebrow_y, title_y, title_h, title_size_pt, diagram_y, diagram_h, label, hint)
VARIANTS = [
    (255, 308, 120, 64, 460, 600, "A — Status quo",
     "current production layout"),
    (135, 188, 120, 64, 340, 680, "B — Moderate",
     "everything 120px higher"),
    ( 75, 128, 120, 64, 280, 740, "C — Recommended",
     "everything 180px higher"),
    ( 40,  93, 120, 64, 240, 780, "D — Maximum",
     "everything 215px higher"),
    ( 40,  78,  80, 48, 200, 820, "E — Compact title",
     "title shrunk to 48pt 1-line, diagram fills the gain"),
]


# Slot horizontal: x=60, w=1800 (60px gutter on both sides). Same for all.
SLOT_X = 60
SLOT_W = 1800

DIAGRAMS = [
    # Two real diagrams from the production decks, picked as ratio extremes.
    # Anything in between (e.g. 04-sbom-inside-sbod.svg at 2.22:1) is covered
    # if both extremes render cleanly in the same layout.
    (DIAGRAMS_DIR / "03-meet-ocm-hub-and-spoke.svg",
     "MEET OCM",
     "One identity, every boundary."),  # 1760×560 ≈ 3.14:1 (very wide)
    (DIAGRAMS_DIR / "06-sovereign-airgap.svg",
     "SOVEREIGN-READY — AIR-GAP",
     "Trust travels with the component."),  # 1600×760 ≈ 2.10:1 (closer to square)
]


def build():
    prs = open_template_as_pptx()
    layouts = {l.name: l for l in prs.slide_masters[0].slide_layouts}
    plain = layouts["Plain"]

    for diag_path, eyebrow_text, title_text in DIAGRAMS:
        for (eyebrow_y, title_y, title_h, title_size,
             diagram_y, diagram_h, label, hint) in VARIANTS:
            s = prs.slides.add_slide(plain)
            # Strip every placeholder the Plain layout brought along — we draw
            # eyebrow, title, and diagram from scratch with test geometry.
            for idx in (1, 2, 10):
                delete_placeholder(s, idx)

            add_variant_caption(s, f"{label}  ·  {hint}")

            add_eyebrow(s, 120, eyebrow_y, 1680, 48, eyebrow_text)
            add_title(s, 120, title_y, 1680, title_h, title_text,
                      size_pt=title_size)

            w, h = add_diagram_centered(s, diag_path,
                                          SLOT_X, diagram_y,
                                          SLOT_W, diagram_h)
            ratio = w / h if h else 0
            stat = (f"slot {SLOT_W}×{diagram_h}  ·  diagram rendered "
                    f"{int(w)}×{int(h)}  (ratio {ratio:.2f}:1)")
            # Bottom caption with the actual rendered numbers.
            tb = s.shapes.add_textbox(px(20), px(SLIDE_H_PX - 50),
                                       px(SLIDE_W_PX - 40), px(28))
            tf = tb.text_frame
            tf.margin_left = tf.margin_right = 0
            tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            from pptx.enum.text import PP_ALIGN
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = stat
            r.font.name = "Aptos"
            r.font.size = Pt(10)
            r.font.italic = True
            r.font.color.rgb = C.GREY_MID

    prs.save(str(OUTPUT_PPTX))
    print(f"Wrote {OUTPUT_PPTX}")


if __name__ == "__main__":
    build()
